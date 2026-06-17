#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Extracteur de Fichiers Incorporés - Version Finale CORRIGÉE

CORRECTIONS APPLIQUÉES :
1. ✅ Activation Excel fonctionnelle (scripts VBS optimisés)
2. ✅ Fichiers ZIP remplacés dans documents modifiés
3. ✅ Nomenclature FJ corrigée (DER-348818_3_FJ_2 → FJ_3, FJ_4...)

Fonctionnalités:
- Détection et activation UNIQUEMENT des objets Excel embarqués
- Extraction avec ordre réel correct (PPTX, DOCX, XLSX)
- Remplacement direct (même nom que l'original)
- Indexation personnalisée via fichier de référence
- Rapport Excel détaillé
- Support: DOCX, DOC, XLSX, XLSM, XLS, PPTX, PPT, PDF, MSG
"""
import threading
import os
import sys
import shutil
import zipfile
import tempfile
from pathlib import Path
from datetime import datetime
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext
import subprocess
import xml.etree.ElementTree as ET
import re
import struct
import traceback
import time
try:
    from lxml import etree
except ImportError:
    install_package_safran("lxml")
    from lxml import etree
# Configuration repository Safran
SAFRAN_REPO = "--index-url https://artifacts.cloud.safran/repository/pypi-group/simple --trusted-host artifacts.cloud.safran"

def install_package_safran(package_name):
    """Installe un package depuis le repository Safran"""
    print(f"Installation de {package_name}...")
    cmd = f"pip install {package_name} {SAFRAN_REPO}"
    result = subprocess.run(cmd, shell=True, capture_output=True, text=True)
    return result.returncode == 0

# Import des bibliothèques
try:
    from docx import Document
    from docx.shared import RGBColor, Pt
except ImportError:
    install_package_safran("python-docx")
    from docx import Document
    from docx.shared import RGBColor, Pt

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
except ImportError:
    install_package_safran("openpyxl")
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

try:
    from pptx import Presentation
    from pptx.util import Pt as PPTPt, Inches
    from pptx.dml.color import RGBColor as PPTRGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    install_package_safran("python-pptx")
    from pptx import Presentation
    from pptx.util import Pt as PPTPt, Inches
    from pptx.dml.color import RGBColor as PPTRGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    import fitz
except ImportError:
    install_package_safran("PyMuPDF")
    import fitz

try:
    import olefile
except ImportError:
    install_package_safran("olefile")
    import olefile
class ExcelActivator:
    """
    Activation des objets Excel embarqués via win32com Python directement.
    Plus de scripts VBS — appel COM natif, plus fiable.
    """

    def __init__(self):
        self.log_messages = []

    def log(self, message):
        self.log_messages.append(message)
        print(message)

    def detect_excel_in_file(self, file_path):
        """Détecte si un fichier contient des objets Excel embarqués"""
        ext = Path(file_path).suffix.lower()
        try:
            if ext == '.pptx':
                return self._detect_excel_in_pptx(file_path)
            elif ext == '.docx':
                return self._detect_excel_in_docx(file_path)
            elif ext in ['.xlsx', '.xlsm']:
                return self._detect_excel_in_xlsx(file_path)
        except Exception:
            pass
        return False

    def _detect_excel_in_pptx(self, pptx_path):
        try:
            with zipfile.ZipFile(pptx_path, 'r') as zf:
                for name in zf.namelist():
                    if name.startswith('ppt/embeddings/') and 'oleObject' in name:
                        data = zf.read(name)
                        if self._is_excel_data(data):
                            return True
            return False
        except:
            return False

    def _detect_excel_in_docx(self, docx_path):
        try:
            with zipfile.ZipFile(docx_path, 'r') as zf:
                for name in zf.namelist():
                    if name.startswith('word/embeddings/') and 'oleObject' in name:
                        data = zf.read(name)
                        if self._is_excel_data(data):
                            return True
            return False
        except:
            return False

    def _detect_excel_in_xlsx(self, xlsx_path):
        try:
            with zipfile.ZipFile(xlsx_path, 'r') as zf:
                for name in zf.namelist():
                    if name.startswith('xl/embeddings/') and 'oleObject' in name:
                        data = zf.read(name)
                        if self._is_excel_data(data):
                            return True
            return False
        except:
            return False

    def _is_excel_data(self, data):
        try:
            if data[:8] == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
                import io
                ole = olefile.OleFileIO(io.BytesIO(data))
                for stream in ole.listdir():
                    name = '/'.join(stream)
                    if 'Workbook' in name or 'Book' in name:
                        ole.close()
                        return True
                ole.close()
            if data[:4] == b'PK\x03\x04':
                if b'xl/' in data[:4096]:
                    return True
            return False
        except:
            return False

    # =========================================================================
    # ACTIVATION DIRECTE via win32com — sans VBS
    # =========================================================================

    def activate_excel_in_file(self, file_path, vbs_scripts=None):
        """
        Active les objets Excel embarqués via win32com Python directement.
        Le paramètre vbs_scripts est ignoré (compatibilité avec l'ancien code).
        Retourne le nombre d'objets activés.
        """
        ext = Path(file_path).suffix.lower()

        try:
            import win32com.client
            import pythoncom
            pythoncom.CoInitialize()

            try:
                if ext == '.pptx':
                    return self._activate_in_pptx(file_path)
                elif ext == '.docx':
                    return self._activate_in_docx(file_path)
                elif ext in ['.xlsx', '.xlsm']:
                    return self._activate_in_xlsx(file_path)
                else:
                    return 0
            finally:
                pythoncom.CoUninitialize()

        except ImportError:
            self.log(f"  ⚠ pywin32 requis pour l'activation Excel")
            return 0
        except Exception as e:
            self.log(f"  ❌ Erreur activation {Path(file_path).name}: {e}")
            return 0
    def _kill_excel(self):
        """Force la fermeture d'Excel et de ses popups"""
        try:
            subprocess.run("taskkill /F /IM EXCEL.EXE", shell=True,
                        capture_output=True)
            time.sleep(0.5)
        except Exception:
            pass
        # Fermer aussi les éventuels popups Office restants
        try:
            subprocess.run("taskkill /F /IM SPLASHWIN.EXE", shell=True,
                        capture_output=True)
        except Exception:
            pass


    def _activate_in_pptx(self, pptx_path):
        """Active objets Excel dans un PPTX via PowerPoint COM"""
        import win32com.client
        count = 0
        ppt = None

        try:
            self.log(f"  🔄 Ouverture PowerPoint : {Path(pptx_path).name}")
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            ppt.Visible = True

            pres = ppt.Presentations.Open(
                str(Path(pptx_path).absolute()),
                ReadOnly=False,
                Untitled=False,
                WithWindow=True
            )
            time.sleep(3)

            for i in range(1, pres.Slides.Count + 1):
                slide = pres.Slides.Item(i)

                # Aller sur la slide
                try:
                    if pres.Windows.Count > 0:
                        pres.Windows.Item(1).View.GotoSlide(i)
                    time.sleep(0.5)
                except Exception:
                    pass

                for j in range(1, slide.Shapes.Count + 1):
                    shape = slide.Shapes.Item(j)

                    try:
                        if shape.Type != 7:  # 7 = msoEmbeddedOLEObject
                            continue

                        prog_id = shape.OLEFormat.ProgID
                        if 'Excel' not in str(prog_id):
                            continue

                        self.log(f"    → Excel détecté slide {i}, shape {j} ({prog_id})")

                        # Activer l'objet
                        shape.OLEFormat.DoVerb(1)
                        time.sleep(5)

                        # Fermer Excel proprement
                        try:
                            xl = win32com.client.GetObject(None, "Excel.Application")
                            if xl:
                                for wb in list(xl.Workbooks):
                                    try:
                                        wb.Save()
                                        wb.Close(True)
                                    except Exception:
                                        pass
                                xl.Quit()
                                time.sleep(2)
                        except Exception:
                            pass

                        self._kill_excel()
                        time.sleep(1)

                        count += 1
                        self.log(f"    ✅ Excel activé (slide {i})")

                    except Exception as shape_err:
                        self.log(f"    ⚠ Shape {j} slide {i}: {shape_err}")
                        self._kill_excel()

            # Sauvegarder et fermer
            try:
                pres.Save()
                time.sleep(2)
                pres.Close()
                ppt.Quit()
            except Exception:
                pass

            self.log(f"  ✅ PPTX : {count} objet(s) Excel activé(s)")
            return count

        except Exception as e:
            self.log(f"  ❌ Erreur activation PPTX : {e}")
            try:
                if ppt:
                    ppt.Quit()
            except Exception:
                pass
            self._kill_excel()
            return count

    def _activate_in_docx(self, docx_path):
        """Active objets Excel dans un DOCX via Word COM"""
        import win32com.client
        count = 0
        word = None

        try:
            self.log(f"  🔄 Ouverture Word : {Path(docx_path).name}")
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = True
            word.DisplayAlerts = 0

            doc = word.Documents.Open(
                str(Path(docx_path).absolute()),
                ReadOnly=False
            )
            time.sleep(3)

            # InlineShapes
            for shape in doc.InlineShapes:
                try:
                    if shape.Type != 4:  # 4 = wdInlineShapeEmbeddedOLEObject
                        continue
                    prog_id = shape.OLEFormat.ProgID
                    if 'Excel' not in str(prog_id):
                        continue

                    self.log(f"    → Excel InlineShape détecté ({prog_id})")
                    shape.OLEFormat.DoVerb(1)
                    time.sleep(5)

                    try:
                        xl = win32com.client.GetObject(None, "Excel.Application")
                        if xl:
                            for wb in list(xl.Workbooks):
                                try:
                                    wb.Save()
                                    wb.Close(True)
                                except Exception:
                                    pass
                            xl.Quit()
                            time.sleep(2)
                    except Exception:
                        pass

                    self._kill_excel()
                    time.sleep(1)
                    count += 1
                    self.log(f"    ✅ Excel activé (InlineShape)")

                except Exception as s_err:
                    self.log(f"    ⚠ InlineShape: {s_err}")
                    self._kill_excel()

            # Shapes flottantes
            for shape in doc.Shapes:
                try:
                    if shape.Type != 8:  # 8 = msoEmbeddedOLEObject
                        continue
                    prog_id = shape.OLEFormat.ProgID
                    if 'Excel' not in str(prog_id):
                        continue

                    self.log(f"    → Excel Shape flottant ({prog_id})")
                    shape.OLEFormat.DoVerb(1)
                    time.sleep(5)

                    try:
                        xl = win32com.client.GetObject(None, "Excel.Application")
                        if xl:
                            for wb in list(xl.Workbooks):
                                try:
                                    wb.Save()
                                    wb.Close(True)
                                except Exception:
                                    pass
                            xl.Quit()
                            time.sleep(2)
                    except Exception:
                        pass

                    self._kill_excel()
                    time.sleep(1)
                    count += 1
                    self.log(f"    ✅ Excel activé (Shape flottant)")

                except Exception as s_err:
                    self.log(f"    ⚠ Shape flottant: {s_err}")
                    self._kill_excel()

            try:
                doc.Save()
                time.sleep(1)
                doc.Close(False)
                word.Quit()
            except Exception:
                pass

            self.log(f"  ✅ DOCX : {count} objet(s) Excel activé(s)")
            return count

        except Exception as e:
            self.log(f"  ❌ Erreur activation DOCX : {e}")
            try:
                if word:
                    word.Quit()
            except Exception:
                pass
            self._kill_excel()
            return count

    def _activate_in_xlsx(self, xlsx_path):
        """Active objets Excel dans un XLSX/XLSM via Excel COM"""
        import win32com.client
        count = 0
        excel = None

        try:
            self.log(f"  🔄 Ouverture Excel : {Path(xlsx_path).name}")
            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = True
            excel.DisplayAlerts = False

            wb = excel.Workbooks.Open(
                str(Path(xlsx_path).absolute()),
                ReadOnly=False
            )
            time.sleep(3)

            for ws in wb.Worksheets:
                for shape in ws.Shapes:
                    try:
                        if shape.Type != 8:  # 8 = msoEmbeddedOLEObject
                            continue
                        prog_id = shape.OLEFormat.ProgID
                        if 'Excel' not in str(prog_id):
                            continue

                        self.log(f"    → Excel embarqué feuille '{ws.Name}' ({prog_id})")
                        shape.OLEFormat.DoVerb(1)
                        time.sleep(4)

                        self._kill_excel()
                        time.sleep(1)
                        count += 1
                        self.log(f"    ✅ Excel activé (feuille {ws.Name})")

                    except Exception as s_err:
                        self.log(f"    ⚠ Shape feuille {ws.Name}: {s_err}")

            try:
                wb.Save()
                wb.Close(True)
                excel.Quit()
            except Exception:
                pass

            self.log(f"  ✅ XLSX : {count} objet(s) Excel activé(s)")
            return count

        except Exception as e:
            self.log(f"  ❌ Erreur activation XLSX : {e}")
            try:
                if excel:
                    excel.Quit()
            except Exception:
                pass
            return count

    def create_vbs_scripts(self, temp_dir):
        """
        Conservé pour compatibilité — ne fait rien.
        L'activation se fait maintenant directement via Python win32com.
        """
        return {}


class SheetColumnSelectorDialog:
    """Dialogue pour sélectionner la feuille et les colonnes du fichier Excel"""
    
    def __init__(self, parent, excel_file_path):
        self.result = None
        self.excel_file_path = excel_file_path
        
        self.dialog = tk.Toplevel(parent)
        self.dialog.title("Sélection de la feuille et des colonnes")
        self.dialog.geometry("600x500")
        self.dialog.transient(parent)
        self.dialog.grab_set()
        
        self.dialog.update_idletasks()
        x = (self.dialog.winfo_screenwidth() // 2) - (300)
        y = (self.dialog.winfo_screenheight() // 2) - (250)
        self.dialog.geometry(f"600x500+{x}+{y}")
        
        self.setup_ui()
        self.load_excel_data()
    
    def setup_ui(self):
        """Configure l'interface du dialogue"""
        main_frame = ttk.Frame(self.dialog, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        title_label = tk.Label(main_frame, 
                              text="📊 Configuration du fichier de référence",
                              font=('Arial', 14, 'bold'), fg='#2c3e50')
        title_label.pack(pady=(0, 20))
        
        sheet_frame = ttk.LabelFrame(main_frame, text="1. Sélectionner la feuille Excel", padding="10")
        sheet_frame.pack(fill=tk.X, pady=(0, 15))
        
        self.sheet_var = tk.StringVar()
        self.sheet_combo = ttk.Combobox(sheet_frame, textvariable=self.sheet_var, 
                                        state='readonly', width=50, font=('Arial', 10))
        self.sheet_combo.pack(fill=tk.X, pady=5)
        self.sheet_combo.bind('<<ComboboxSelected>>', self.on_sheet_selected)
        
        columns_frame = ttk.LabelFrame(main_frame, text="2. Sélectionner les colonnes", padding="10")
        columns_frame.pack(fill=tk.BOTH, expand=True, pady=(0, 15))
        
        filename_subframe = ttk.Frame(columns_frame)
        filename_subframe.pack(fill=tk.X, pady=5)
        
        tk.Label(filename_subframe, text="📄 Colonne 'Nom' :", 
                font=('Arial', 10, 'bold'), fg='#2c3e50').pack(side=tk.LEFT, padx=(0, 10))
        
        self.filename_col_var = tk.StringVar()
        self.filename_combo = ttk.Combobox(filename_subframe, textvariable=self.filename_col_var,
                                           state='readonly', width=30, font=('Arial', 10))
        self.filename_combo.pack(side=tk.LEFT, fill=tk.X, expand=True)
        
        children_subframe = ttk.Frame(columns_frame)
        children_subframe.pack(fill=tk.X, pady=5)
        
        tk.Label(children_subframe, text="🔢 Colonne 'Nombre d'enfants' :", 
                font=('Arial', 10, 'bold'), fg='#2c3e50').pack(side=tk.LEFT, padx=(0, 10))
        
        self.children_col_var = tk.StringVar()
        self.children_combo = ttk.Combobox(children_subframe, textvariable=self.children_col_var,
                                          state='readonly', width=30, font=('Arial', 10))
        self.children_combo.pack(side=tk.LEFT, fill=tk.X, expand=True)
        
        preview_frame = ttk.LabelFrame(columns_frame, text="Aperçu des données (5 premières lignes)", 
                                       padding="10")
        preview_frame.pack(fill=tk.BOTH, expand=True, pady=(10, 0))
        
        self.preview_text = scrolledtext.ScrolledText(preview_frame, height=8, width=60,
                                                     bg='#ecf0f1', font=('Consolas', 9), wrap=tk.NONE)
        self.preview_text.pack(fill=tk.BOTH, expand=True)
        
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill=tk.X, pady=(10, 0))
        
        tk.Button(button_frame, text="✓ Valider", command=self.validate,
                 bg='#27ae60', fg='white', font=('Arial', 11, 'bold'),
                 padx=20, pady=8).pack(side=tk.LEFT, padx=5)
        
        tk.Button(button_frame, text="✗ Annuler", command=self.cancel,
                 bg='#e74c3c', fg='white', font=('Arial', 11, 'bold'),
                 padx=20, pady=8).pack(side=tk.LEFT, padx=5)
    
    def load_excel_data(self):
        """Charge les feuilles du fichier Excel"""
        try:
            self.wb = openpyxl.load_workbook(self.excel_file_path, read_only=True)
            sheet_names = self.wb.sheetnames
            
            self.sheet_combo['values'] = sheet_names
            if sheet_names:
                self.sheet_combo.current(0)
                self.on_sheet_selected(None)
        
        except Exception as e:
            messagebox.showerror("Erreur", f"Impossible de lire le fichier Excel:\n{e}")
            self.dialog.destroy()
    
    def on_sheet_selected(self, event):
        """Appelé quand une feuille est sélectionnée"""
        sheet_name = self.sheet_var.get()
        if not sheet_name:
            return
        
        try:
            ws = self.wb[sheet_name]
            
            first_row = next(ws.iter_rows(min_row=1, max_row=1, values_only=True))
            
            columns = []
            for idx, cell_value in enumerate(first_row, 1):
                col_letter = openpyxl.utils.get_column_letter(idx)
                col_name = str(cell_value) if cell_value else f"Colonne {col_letter}"
                columns.append(f"{col_letter} - {col_name}")
            
            self.filename_combo['values'] = columns
            self.children_combo['values'] = columns
            
            if len(columns) > 0:
                self.filename_combo.current(0)
            if len(columns) > 1:
                self.children_combo.current(1)
            
            self.update_preview()
        
        except Exception as e:
            messagebox.showerror("Erreur", f"Erreur lors de la lecture de la feuille:\n{e}")
    
    def update_preview(self):
        """Met à jour l'aperçu des données"""
        sheet_name = self.sheet_var.get()
        if not sheet_name:
            return
        
        try:
            ws = self.wb[sheet_name]
            
            self.preview_text.delete(1.0, tk.END)
            
            for idx, row in enumerate(ws.iter_rows(min_row=1, max_row=6, values_only=True)):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                
                if idx == 0:
                    self.preview_text.insert(tk.END, row_text + "\n", "header")
                    self.preview_text.insert(tk.END, "-" * 80 + "\n")
                else:
                    self.preview_text.insert(tk.END, row_text + "\n")
            
            self.preview_text.tag_config("header", font=('Consolas', 9, 'bold'))
        
        except Exception as e:
            self.preview_text.insert(tk.END, f"Erreur d'aperçu: {e}\n")
    
    def validate(self):
        """Valide la sélection"""
        sheet_name = self.sheet_var.get()
        filename_col = self.filename_col_var.get()
        children_col = self.children_col_var.get()
        
        if not sheet_name:
            messagebox.showwarning("Attention", "Veuillez sélectionner une feuille!")
            return
        
        if not filename_col:
            messagebox.showwarning("Attention", "Veuillez sélectionner la colonne 'Nom'!")
            return
        
        if not children_col:
            messagebox.showwarning("Attention", "Veuillez sélectionner la colonne 'Nombre d'enfants'!")
            return
        
        filename_letter = filename_col.split(' - ')[0]
        children_letter = children_col.split(' - ')[0]
        
        filename_idx = openpyxl.utils.column_index_from_string(filename_letter)
        children_idx = openpyxl.utils.column_index_from_string(children_letter)
        
        self.result = {
            'sheet_name': sheet_name,
            'filename_col_idx': filename_idx,
            'children_col_idx': children_idx,
            'filename_col_letter': filename_letter,
            'children_col_letter': children_letter
        }
        
        self.wb.close()
        self.dialog.destroy()
    
    def cancel(self):
        """Annule la sélection"""
        self.result = None
        self.wb.close()
        self.dialog.destroy()
    
    def get_result(self):
        """Retourne le résultat de la sélection"""
        self.dialog.wait_window()
        return self.result


class IndexationManager:
    """Gestionnaire d'indexation basé sur fichier de référence"""
    
    def __init__(self, reference_file=None, config=None):
        self.reference_data = {}
        if reference_file and config:
            self.load_reference_file(reference_file, config)
    
    def load_reference_file(self, file_path, config):
        """Charge le fichier de référence Excel avec configuration personnalisée"""
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            ws = wb[config['sheet_name']]
            
            filename_col = config['filename_col_idx']
            children_col = config['children_col_idx']
            
            for row in ws.iter_rows(min_row=2, values_only=True):
                if len(row) >= max(filename_col, children_col):
                    filename = row[filename_col - 1]
                    nb_enfants = row[children_col - 1]
                    
                    if filename and nb_enfants is not None:
                        try:
                            filename = str(filename).strip()
                            nb_enfants = int(nb_enfants)
                            self.reference_data[filename] = nb_enfants
                        except (ValueError, TypeError):
                            continue
            
            wb.close()
            print(f"✓ Référence chargée : {len(self.reference_data)} fichiers")
            print(f"  Feuille: {config['sheet_name']}")
            print(f"  Colonne noms: {config['filename_col_letter']}")
            print(f"  Colonne enfants: {config['children_col_letter']}")
            
        except Exception as e:
            print(f"⚠ Erreur chargement référence : {e}")
    
    def get_start_index(self, filename):
        """Retourne l'index de départ pour un fichier"""
        stem = Path(filename).stem

        # Extraire uniquement la partie jusqu'à _FJ (inclus), ignorer le reste
        fj_match = re.match(r'^(.*?_FJ)', stem)
        search_key = fj_match.group(1) if fj_match else stem

        # Chercher la clé tronquée directement dans les données de référence
        if search_key in self.reference_data:
            return self.reference_data[search_key]

        # Fallback : comparer avec les stems des clés de référence (tronqués aussi)
        for ref_name, nb_enfants in self.reference_data.items():
            ref_stem = Path(ref_name).stem
            ref_fj_match = re.match(r'^(.*?_FJ)', ref_stem)
            ref_key = ref_fj_match.group(1) if ref_fj_match else ref_stem
            if ref_key == search_key:
                return nb_enfants

        return 1
class EmbeddedFileExtractor:
    """Extracteur spécialisé pour les fichiers incorporés avec ordre corrigé"""
    
    def __init__(self, indexation_manager=None):
        self.extracted_count = 0
        self.log_messages = []
        self.indexation_manager = indexation_manager
        self.current_file_extractions = {}
        
        self.report_data = {
            'files_processed': [],
            'files_extracted': [],
            'errors': []
        }
        self.current_source_file = None
        self.recursion_depth = 0  # ← AJOUTER CETTE LIGNE
        self.total_input_files = 0
        self.manual_start_index = None
                # ── AJOUTER CES DEUX LIGNES ICI ──────────────────────────────
        self._msg_already_processed = set()
        self._fj_floor = 0
        self._extracted_content_hashes = set()
        self._processed_absolute_paths = set()
        self._content_hash_to_name = {}
        self._used_fj_numbers = set()
    def log(self, message):
        """Enregistre un message de log"""
        self.log_messages.append(message)
        print(message)
    def _kill_excel(self):
        """Force la fermeture d'Excel"""
        try:
            subprocess.run("taskkill /F /IM EXCEL.EXE", shell=True, capture_output=True)
            time.sleep(0.5)
        except Exception:
            pass
    def _extract_msg_attachments_only(self, msg_path, output_dir):
        """
        Extrait uniquement les pièces jointes d'un MSG (sans conversion PDF).
        Utilisé avant convert_msg_to_pdf pour avoir la liste extracted_names.
        Retourne la liste des noms de fichiers extraits.
        """
        extracted_names = []
        try:
            import win32com.client
            import pythoncom

            pythoncom.CoInitialize()
            try:
                SUPPORTED_EXTENSIONS = {
                    '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls',
                    '.pptx', '.ppt', '.msg', '.txt', '.zip', '.7z', '.pptm',
                    '.htm', '.html'
                }

                outlook = win32com.client.Dispatch("Outlook.Application")
                msg_com = outlook.Session.OpenSharedItem(str(Path(msg_path).absolute()))

                attachments = msg_com.Attachments
                nb_att = attachments.Count
                self.log(f"    📎 {nb_att} pièce(s) jointe(s) dans le MSG encapsulé")

                for i in range(1, nb_att + 1):
                    attachment = attachments.Item(i)
                    att_name = attachment.FileName

                    if not att_name:
                        continue

                    att_ext = Path(att_name).suffix.lower()
                    if att_ext not in SUPPORTED_EXTENSIONS:
                        _IMG_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff',
                                     '.webp', '.svg', '.ico', '.emf', '.wmf'}
                        if att_ext not in _IMG_EXTS:
                            self.log(f"    ⚠️ Extension non supportée ignorée : {att_name}")
                        else:
                            self.log(f"    ⏭️ Image ignorée (MSG): {att_name}")
                        continue

                    self.extracted_count += 1
                    output_name = self._generate_output_name(
                        Path(msg_path).stem, att_name, att_ext
                    )
                    output_name_clean = output_name.strip()
                    output_path = Path(output_dir) / output_name_clean

                    try:
                        attachment.SaveAsFile(str(output_path.absolute()))
                        self.log(f"    ✅ PJ extraite : {output_name_clean}")
                        extracted_names.append(output_name_clean)

                        self.add_to_report('extracted', {
                            'source_file': Path(msg_path).name,
                            'extracted_file': output_name_clean,
                            'position': f'Pièce jointe {i}',
                            'type': att_ext,
                            'status': 'Succès'
                        })
                    except Exception as save_err:
                        self.log(f"    ❌ Erreur sauvegarde PJ : {save_err}")

                msg_com.Close(0)

            finally:
                pythoncom.CoUninitialize()

        except ImportError:
            self.log(f"    ⚠️ pywin32 requis pour extraction PJ MSG")
        except Exception as e:
            self.log(f"    ⚠️ Erreur extraction PJ MSG : {e}")

        return extracted_names
    def process_extracted_files_recursively(self, output_dir, processed_signatures=None, depth=0, max_depth=3):
        """
        Pour chaque fichier extrait dans output_dir,
        relance process_file() exactement comme au premier passage.
        3 niveaux de récursivité maximum.
        """
        import hashlib

        if processed_signatures is None:
            processed_signatures = set()

        if depth >= max_depth:
            self.log(f"  ⚠️ Profondeur max ({max_depth}) atteinte — arrêt récursion")
            return 0

        SUPPORTED = {
            '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls',
            '.pptx', '.pptm', '.ppt', '.msg', '.zip',
            '.htm', '.html'
        }

        output_dir = Path(output_dir)

        # ── Collecter les fichiers nouveaux (pas encore traités) ──────────────
        files_to_process = []

        for file_path in sorted(output_dir.iterdir()):
            if not file_path.is_file():
                continue
            if file_path.suffix.lower() not in SUPPORTED:
                continue

            try:
                stat = file_path.stat()
                with open(file_path, 'rb') as f:
                    head = f.read(2048)
                sig = f"{file_path.name}|{stat.st_size}|{hashlib.md5(head).hexdigest()}"
            except Exception:
                sig = str(file_path.absolute())

            if sig not in processed_signatures:
                files_to_process.append((file_path, sig))

        if not files_to_process:
            self.log(f"  ✓ Niveau {depth + 1} : rien de nouveau — récursion terminée")
            return 0

        self.log(f"\n  🔄 RÉCURSION niveau {depth + 1} : {len(files_to_process)} fichier(s) à traiter")

        total_new = 0

        for file_path, sig in files_to_process:
            processed_signatures.add(sig)

            # ── Skip MSG déjà traités par _extract_msg_attachments_only ──
       
       




            # Skip fichiers deja traites par chemin absolu
            abs_path = str(file_path.resolve())
            if abs_path in self._processed_absolute_paths:
                self.log(f"\n    [SKIP] Niv.{depth + 1} - Deja traite : {file_path.name}")
                continue
            self._processed_absolute_paths.add(abs_path)

            # ── Skip MSG déjà traités par _extract_msg_attachments_only ──
            if abs_path in self._msg_already_processed:





                
                self.log(f"\n    ⏭️ [Niv.{depth + 1}] MSG déjà traité (PJ extraites) : {file_path.name}")
                continue

            self.log(f"\n    📄 [Niv.{depth + 1}] {file_path.name}")

            try:
                extracted = self.process_file(file_path, output_dir)
                count = len(extracted) if extracted else 0
                total_new += count

                if count > 0:
                    self.log(f"      → {count} fichier(s) extrait(s)")
                else:
                    self.log(f"      → Aucun fichier incorporé")

                # Supprimer le fichier MSG du dossier de sortie après traitement
                if file_path.suffix.lower() == '.msg' and file_path.exists():
                    try:
                        file_path.unlink()
                        self.log(f"      🗑️ MSG supprimé du dossier de sortie: {file_path.name}")
                    except Exception as del_err:
                        self.log(f"      ⚠️ Impossible de supprimer le MSG: {del_err}")

            except Exception as e:
                self.log(f"      ❌ Erreur : {e}")
                # Supprimer quand même le MSG même si erreur de traitement
                if file_path.suffix.lower() == '.msg' and file_path.exists():
                    try:
                        file_path.unlink()
                        self.log(f"      🗑️ MSG supprimé (après erreur): {file_path.name}")
                    except Exception as del_err:
                        self.log(f"      ⚠️ Impossible de supprimer le MSG: {del_err}")

        # ── Si des fichiers ont été extraits, descendre d'un niveau ──────────
        if total_new > 0:
            self.log(f"\n  🔁 {total_new} nouveau(x) extrait(s) → passage niveau {depth + 2}...")
            deeper = self.process_extracted_files_recursively(
                output_dir, processed_signatures, depth + 1, max_depth
            )
            total_new += deeper

        return total_new
# ============================================================================
    # CORRECTION 1 : Mapper correctement les positions OLE
    # ============================================================================

    def get_ole_positions_mapping(self, file_path):  # ← AJOUT self
        """
        Retourne un dictionnaire {file_num: position_document}
        
        Exemple pour PPTX :
        - Slide 1, objet 1 → position 1
        - Slide 1, objet 2 → position 2
        - Slide 2, objet 1 → position 3
        """
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pptx':
            return self.get_pptx_ole_positions(file_path)  # ← AJOUT self.
        elif ext == '.docx':
            return self.get_docx_ole_positions(file_path)  # ← AJOUT self.
        elif ext in ['.xlsx', '.xlsm']:
            return self.get_xlsx_ole_positions(file_path)  # ← AJOUT self.
        
        return {}


    def get_pptx_ole_positions(self, pptx_path):
        """
        Pour PPTX : parcourt les slides dans l'ordre et numérote les objets OLE
        Retourne : {file_num: position_document}
        """
        position_mapping = {}
        
        try:
            with zipfile.ZipFile(pptx_path, 'r') as zf:
                # Lister les slides dans l'ordre
                slide_files = [f for f in zf.namelist() 
                            if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                slide_files.sort(key=lambda x: int(re.search(r'slide(\d+)', x).group(1)))
                
                current_position = 0
                
                for slide_file in slide_files:
                    slide_num = int(re.search(r'slide(\d+)', slide_file).group(1))
                    
                    # Lire le XML de la slide
                    xml_content = zf.read(slide_file).decode('utf-8')
                    
                    # Trouver tous les objets OLE dans l'ordre d'apparition
                    ole_matches = re.findall(r'<p:oleObj[^>]*r:id="([^"]+)"', xml_content)
                    
                    # Pour chaque objet OLE, trouver le fichier correspondant
                    rels_file = slide_file.replace('slides/slide', 'slides/_rels/slide').replace('.xml', '.xml.rels')
                    
                    if rels_file in zf.namelist():
                        rels_xml = zf.read(rels_file)
                        rels_root = ET.fromstring(rels_xml)
                        
                        ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                        
                        for r_id in ole_matches:
                            current_position += 1
                            
                            # Trouver le fichier correspondant
                            for rel in rels_root.findall('r:Relationship', ns):
                                if rel.get('Id') == r_id:
                                    target = rel.get('Target')
                                    
                                    # Extraire le numéro du fichier
                                    match = re.search(r'oleObject(\d+)', target)
                                    if match:
                                        file_num = int(match.group(1))
                                        position_mapping[file_num] = current_position
                                        self.log(f"        🔍 PPTX Mapping: oleObject{file_num} → Position {current_position} (Slide {slide_num})")
                                    break
            
            return position_mapping
            
        except Exception as e:
            self.log(f"    ⚠️ Erreur mapping PPTX: {e}")
            return {}
    def get_docx_ole_positions_exact(self, docx_path, temp_path=None):
        """
        Détecte positions EXACTES avec numéro de page estimé
        Retourne : {nom_fichier: {'position': X, 'page': Y, 'para_element': element, ...}}
        """
        position_mapping = {}
        
        try:
            from lxml import etree
            
            _temp_dir_created = None
            if temp_path is None:
                import tempfile
                _temp_dir_created = tempfile.mkdtemp()
                try:
                    with zipfile.ZipFile(docx_path, 'r') as zf:
                        zf.extractall(_temp_dir_created)
                    temp_path = Path(_temp_dir_created)
                except Exception:
                    shutil.rmtree(_temp_dir_created, ignore_errors=True)
                    _temp_dir_created = None
                    raise
            
            # Lire XMLs
            doc_xml_path = temp_path / 'word' / 'document.xml'
            rels_xml_path = temp_path / 'word' / '_rels' / 'document.xml.rels'
            
            doc_tree = etree.parse(str(doc_xml_path))
            doc_root = doc_tree.getroot()
            
            rels_tree = etree.parse(str(rels_xml_path))
            rels_root = rels_tree.getroot()
            
            # Namespaces
            ns = {
                'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'o': 'urn:schemas-microsoft-com:office:office',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'rel': 'http://schemas.openxmlformats.org/package/2006/relationships',
                'v': 'urn:schemas-microsoft-com:vml'
            }
            
            # Mapping r:id → fichier
            rid_to_file = {}
            for rel in rels_root.xpath('rel:Relationship', namespaces={'rel': ns['rel']}):
                r_id = rel.get('Id')
                target = rel.get('Target')
                if target and 'embeddings/' in target:
                    filename = target.split('/')[-1]
                    rid_to_file[r_id] = filename
            
            current_position = 0
            current_page = 1  # Estimation
            paragraph_count = 0
            
            # Fonction récursive
            def analyze_element(element, context=""):
                nonlocal current_position, current_page, paragraph_count
                
                # Détecter sauts de page
                if element.tag == f"{{{ns['w']}}}br":
                    br_type = element.get(f"{{{ns['w']}}}type")
                    if br_type == "page":
                        current_page += 1
                
                # Compter paragraphes (estimation page)
                if element.tag == f"{{{ns['w']}}}p":
                    paragraph_count += 1
                    # Environ 30 paragraphes par page
                    current_page = (paragraph_count // 30) + 1
                
                # Chercher objets OLE
                for obj in element.xpath('.//w:object', namespaces=ns):
                    ole_obj = obj.find('.//o:OLEObject', namespaces=ns)
                    if ole_obj is not None:
                        r_id = ole_obj.get(f"{{{ns['r']}}}id")
                        
                        if r_id and r_id in rid_to_file:
                            current_position += 1
                            filename = rid_to_file[r_id]
                            
                            position_mapping[filename] = {
                                'position': current_position,
                                'page': current_page,
                                'context': context,
                                'element': obj,
                                'paragraph': element if element.tag == f"{{{ns['w']}}}p" else None
                            }
                            
                            self.log(f"        🔍 {filename} → Position {current_position} (Page ~{current_page}, {context})")
                
                # Chercher VML shapes
                for pict in element.xpath('.//w:pict', namespaces=ns):
                    for shape in pict.xpath('.//v:shape', namespaces=ns):
                        for imagedata in shape.xpath('.//v:imagedata', namespaces=ns):
                            r_id = imagedata.get(f"{{{ns['r']}}}id")
                            
                            if r_id and r_id in rid_to_file:
                                current_position += 1
                                filename = rid_to_file[r_id]
                                
                                position_mapping[filename] = {
                                    'position': current_position,
                                    'page': current_page,
                                    'context': f"{context} (VML)",
                                    'element': pict,
                                    'paragraph': element if element.tag == f"{{{ns['w']}}}p" else None
                                }
                                
                                self.log(f"        🔍 {filename} → Position {current_position} (Page ~{current_page}, {context})")
            
            # Parcourir body
            body = doc_root.find('.//w:body', namespaces=ns)
            if body is not None:
                for idx, element in enumerate(body):
                    tag = element.tag.replace(f"{{{ns['w']}}}", "")
                    
                    if tag == 'p':
                        analyze_element(element, f"Paragraphe")
                    elif tag == 'tbl':
                        for row_idx, row in enumerate(element.xpath('.//w:tr', namespaces=ns), 1):
                            for cell_idx, cell in enumerate(row.xpath('.//w:tc', namespaces=ns), 1):
                                analyze_element(cell, f"Tableau L{row_idx}C{cell_idx}")
                    else:
                        analyze_element(element, f"Element {tag}")
            
            self.log(f"        ✅ Mapping DOCX : {len(position_mapping)} fichier(s)")

            return position_mapping

        except Exception as e:
            self.log(f"    ⚠️ Erreur mapping: {e}")
            return {}
        finally:
            if _temp_dir_created:
                shutil.rmtree(_temp_dir_created, ignore_errors=True)

    def get_xlsx_ole_positions_exact(self, xlsx_path, temp_path):
        """
        Détecte positions EXACTES (feuille + cellule approximative)
        """
        position_mapping = {}
        
        try:
            from lxml import etree
            
            ns = {
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'
            }
            
            # Lire workbook pour noms de feuilles
            workbook_path = temp_path / 'xl' / 'workbook.xml'
            wb_tree = etree.parse(str(workbook_path))
            wb_root = wb_tree.getroot()
            
            sheet_names = {}
            for sheet in wb_root.xpath('//ns:sheet', namespaces={'ns': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}):
                sheet_id = sheet.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                sheet_name = sheet.get('name')
                if sheet_id and sheet_name:
                    sheet_names[sheet_id] = sheet_name
            
            # Parcourir feuilles
            worksheets_dir = temp_path / 'xl' / 'worksheets'
            sheet_files = sorted(worksheets_dir.glob('sheet*.xml'), key=lambda x: int(re.search(r'sheet(\d+)', x.name).group(1)))
            
            current_position = 0
            
            for sheet_file in sheet_files:
                sheet_num = int(re.search(r'sheet(\d+)', sheet_file.name).group(1))
                
                # Trouver nom de la feuille
                rels_file = worksheets_dir / '_rels' / f'sheet{sheet_num}.xml.rels'
                
                if not rels_file.exists():
                    continue
                
                rels_tree = etree.parse(str(rels_file))
                rels_root = rels_tree.getroot()
                
                for rel in rels_root.xpath('rel:Relationship', namespaces={'rel': ns['rel']}):
                    target = rel.get('Target')
                    
                    if target and 'embeddings/' in target:
                        current_position += 1
                        filename = target.split('/')[-1]
                        
                        # Trouver nom de la feuille
                        sheet_name = f'Feuille {sheet_num}'
                        for sid, sname in sheet_names.items():
                            if f'sheet{sheet_num}' in sid.lower():
                                sheet_name = sname
                                break
                        
                        position_mapping[filename] = {
                            'position': current_position,
                            'sheet_number': sheet_num,
                            'sheet_name': sheet_name
                        }
                        
                        self.log(f"        🔍 {filename} → Position {current_position} (Feuille: {sheet_name})")
            
            self.log(f"        ✅ Mapping XLSX : {len(position_mapping)} fichier(s)")
            
            return position_mapping
            
        except Exception as e:
            self.log(f"    ⚠️ Erreur mapping: {e}")
            return {}


    def add_to_report(self, report_type, data):
        """Ajoute une entrée au rapport (sans doublons)"""
        if report_type == 'processed':
            # Éviter les doublons pour les fichiers traités aussi
            file_name = data.get('file_name', '')
            is_duplicate = any(
                existing.get('file_name') == file_name 
                for existing in self.report_data['files_processed']
            )
            
            if not is_duplicate:
                self.report_data['files_processed'].append(data)
                
        elif report_type == 'extracted':
            # Vérifier unicité par source_file + position + extracted_file
            is_duplicate = False
            
            if 'source_file' in data and 'position' in data:
                for existing in self.report_data['files_extracted']:
                    if (existing.get('source_file') == data['source_file'] and 
                        existing.get('position') == data['position'] and
                        existing.get('extracted_file') == data.get('extracted_file')):
                        is_duplicate = True
                        break
            
            if not is_duplicate:
                self.report_data['files_extracted'].append(data)
                
        elif report_type == 'error':
            # Éviter les doublons d'erreurs
            error_msg = data.get('error', '')
            file_name = data.get('file_name', '')
            
            is_duplicate = any(
                existing.get('error') == error_msg and 
                existing.get('file_name') == file_name
                for existing in self.report_data['errors']
            )
            
            if not is_duplicate:
                self.report_data['errors'].append(data)
    def reset_extraction_tracking(self, source_filename, output_dir=None):
        """
        Réinitialise COMPLÈTEMENT le suivi des extractions pour chaque nouveau fichier source.
        """
        self.current_file_extractions = {}
        self.current_source_file = source_filename
        self.extracted_count = 0

        # ── RÉINITIALISATION COMPLÈTE pour chaque nouveau fichier ────────
        self._fj_floor = 0
        self._used_fj_numbers = set()
        self._extracted_content_hashes = set()
        self._content_hash_to_name = {}

        # ── Charger TOUS les numéros FJ existants depuis le dossier ──────
        if output_dir:
            directory = Path(output_dir)
            if directory.exists():
                for fp in directory.iterdir():
                    if fp.is_file():
                        m = re.search(r'_FJ_(\d+)', fp.stem)
                        if m:
                            n = int(m.group(1))
                            self._used_fj_numbers.add(n)
                            if n > self._fj_floor:
                                self._fj_floor = n
            if self._fj_floor > 0:
                self.log(f"  📊 Plancher FJ : {self._fj_floor} (prochain libre = FJ_{self._fj_floor + 1})")
            if self._used_fj_numbers:
                self.log(f"  📊 Index FJ réservés : {sorted(self._used_fj_numbers)}")

        # ── Log référence / index manuel ──────────────────────────────────
        if self.indexation_manager:
            start_index = self.indexation_manager.get_start_index(source_filename)
            if start_index > 1:
                self.log(f"  📊 RÉFÉRENCE TROUVÉE : Départ à FJ_{start_index}")
            else:
                self.log(f"  ℹ️ Pas de référence : Numérotation d'après le nom du fichier")
        else:
            self.log(f"  ℹ️ Pas de gestionnaire de référence")
    def is_file_not_image(self, filename, data=None):
        """
        Retourne False (= ignorer) pour :
        - images classiques (png, jpg, emf, wmf...)
        - bitmaps OLE (Paint.Picture, MSPaintBrush...)
        - formules mathématiques (Equation.3, MathType...)
        """
        # ── Extensions image classiques ───────────────────────────────────
        image_extensions = {
            '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff',
            '.svg', '.ico', '.webp', '.emf', '.wmf'
        }
        ext = Path(filename).suffix.lower()
        if ext in image_extensions:
            return False

        # ── Noms de fichiers typiques des bitmaps/équations OLE ───────────
        # Word/PowerPoint nomme ces objets avec des patterns reconnaissables
        ignorable_name_patterns = [
            'equation', 'mathtype', 'msequation',
            'paint', 'bitmap', 'msbitmap',
            'microsoft_equation', 'microsoft equation',
            'corel', 'wmf', 'emf',
        ]
        name_lower = filename.lower()
        for pattern in ignorable_name_patterns:
            if pattern in name_lower:
                return False

        # ── Détection par magic bytes ──────────────────────────────────────
        if data and len(data) > 4:
            image_signatures = [
                b'\xff\xd8\xff',           # JPEG
                b'\x89PNG',                # PNG
                b'GIF87a',                 # GIF
                b'GIF89a',                 # GIF
                b'BM',                     # BMP
                b'II\x2a\x00',             # TIFF little-endian
                b'MM\x00\x2a',             # TIFF big-endian
            ]
            for sig in image_signatures:
                if data[:len(sig)] == sig:
                    return False

            # ── OLE bitmap / équation : détecter via streams OLE ──────────
            if data[:4] == b'\xd0\xcf\x11\xe0':
                try:
                    import io
                    test_ole = olefile.OleFileIO(io.BytesIO(data[:min(len(data), 512000)]))
                    streams = ['/'.join(e) for e in test_ole.listdir()]
                    test_ole.close()

                    # Streams caractéristiques des équations Word
                    equation_streams = [
                        'Ole10Native', '\x01Ole10Native',
                        'Equation Native', 'MathType',
                        '\x01CompObj',
                    ]
                    # Vérifier le CompObj pour lire le ProgID
                    for comp_name in ['\x01CompObj', 'CompObj']:
                        if any(comp_name in s for s in streams):
                            try:
                                test_ole2 = olefile.OleFileIO(io.BytesIO(data[:min(len(data), 512000)]))
                                comp_data = test_ole2.openstream(comp_name).read()
                                test_ole2.close()
                                comp_str = comp_data.decode('latin-1', errors='ignore').lower()

                                ignorable_progids = [
                                    'equation', 'mathtype', 'paint.picture',
                                    'msbitmap', 'paintbrush', 'ms_clipart',
                                    'photoeditor', 'msdraw',
                                ]
                                for pid in ignorable_progids:
                                    if pid in comp_str:
                                        return False
                            except Exception:
                                pass

                    # Stream "Equation Native" = formule Word
                    if any('equation native' in s.lower() for s in streams):
                        return False

                except Exception:
                    pass

        return True

    def clean_bin_files_from_embeddings(self, file_path):
        """
        NOUVEAU : Supprime les fichiers .bin de 0 Ko dans les embeddings avant traitement
        Évite les problèmes lors de l'extraction
        """
        try:
            file_path = Path(file_path)
            ext = file_path.suffix.lower()
            
            # Vérifier si c'est un fichier ZIP Office
            if ext not in ['.docx', '.pptx', '.xlsx', '.xlsm']:
                return 0
            
            bin_deleted = 0
            
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                
                # Extraire le fichier
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_path)
                
                # Trouver le dossier embeddings selon le type
                if ext == '.pptx':
                    embeddings_dir = temp_path / "ppt" / "embeddings"
                elif ext == '.docx':
                    embeddings_dir = temp_path / "word" / "embeddings"
                elif ext in ['.xlsx', '.xlsm']:
                    embeddings_dir = temp_path / "xl" / "embeddings"
                else:
                    return 0
                
                # Si le dossier embeddings existe
                if embeddings_dir.exists():
                    for embed_file in embeddings_dir.iterdir():
                        if embed_file.is_file():
                            # Vérifier si c'est un .bin de 0 Ko
                            if embed_file.suffix.lower() == '.bin' and embed_file.stat().st_size == 0:
                                self.log(f"    🗑️ Suppression fichier .bin vide: {embed_file.name}")
                                embed_file.unlink()
                                bin_deleted += 1
                    
                    # Si des .bin ont été supprimés, recréer le fichier
                    if bin_deleted > 0:
                        self.log(f"    ♻️ Reconstruction du fichier sans les .bin vides...")
                        
                        # Créer fichier temporaire de sortie
                        temp_output = temp_path / f"{file_path.stem}_cleaned{file_path.suffix}"
                        
                        # Recréer le ZIP
                        with zipfile.ZipFile(temp_output, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                            for root, dirs, files in os.walk(temp_path):
                                for file in files:
                                    file_path_full = Path(root) / file
                                    arcname = file_path_full.relative_to(temp_path)
                                    zip_out.write(file_path_full, arcname)
                        
                        # Remplacer le fichier original
                        shutil.copy2(temp_output, file_path)
                        self.log(f"    ✅ Fichier nettoyé: {bin_deleted} .bin supprimé(s)")
            
            return bin_deleted
            
        except Exception as e:
            self.log(f"    ⚠️ Erreur nettoyage .bin: {e}")
            return 0    
    # ============================================================================
    # MÉTHODES POUR DÉTERMINER L'ORDRE RÉEL DES OBJETS OLE
    # ============================================================================
    def convert_legacy_to_modern_format(self, file_path):
        """
        Convertit les anciens formats Office vers les nouveaux formats
        .doc → .docx / .ppt → .pptx / .xls → .xlsx (.xlsm si VBA)

        CORRECTION POPUP ReportINI.xls :
        - Renomme temporairement le dossier XLSTART avant de lancer Excel
        - Désactive add-ins COM via registre temporairement
        - Lance Excel avec /automation /e /safemode
        - DisplayAlerts=False AVANT toute ouverture de fichier
        - Restaure XLSTART après la conversion
        """
        try:
            file_path = Path(file_path)
            ext = file_path.suffix.lower()

            if ext not in ['.doc', '.ppt', '.xls']:
                return file_path, False

            self.log(f"    🔄 Conversion format ancien → moderne: {file_path.name}")

            import win32com.client
            import pythoncom

            pythoncom.CoInitialize()

            try:
                # ── .doc → .docx ──────────────────────────────────────────
                if ext == '.doc':
                    new_path = file_path.with_suffix('.docx')
                    _conv_result = [None]
                    _conv_error  = [None]

                    def _do_doc_convert():
                        import win32com.client as _wcc
                        import pythoncom as _pycom
                        _pycom.CoInitialize()
                        _word = None
                        try:
                            _word = _wcc.DispatchEx("Word.Application")
                            _word.Visible = False
                            _word.DisplayAlerts = 0
                            _word.AutomationSecurity = 3  # désactive macros sans dialogue
                            _doc = _word.Documents.Open(
                                str(file_path.absolute()),
                                ConfirmConversions=False,
                                ReadOnly=True,
                                AddToRecentFiles=False,
                                PasswordDocument="",
                                PasswordTemplate="",
                                Revert=False,
                                NoEncodingDialog=True,
                            )
                            _doc.SaveAs(str(new_path.absolute()), FileFormat=16)
                            _doc.Close(False)
                            _conv_result[0] = new_path
                        except Exception as _e:
                            _conv_error[0] = str(_e)
                        finally:
                            try:
                                if _word is not None:
                                    _word.Quit()
                            except Exception:
                                pass
                            try:
                                _pycom.CoUninitialize()
                            except Exception:
                                pass

                    _t = threading.Thread(target=_do_doc_convert, daemon=True)
                    _t.start()
                    _t.join(timeout=90)

                    if _t.is_alive():
                        try:
                            subprocess.run(['taskkill', '/F', '/IM', 'WINWORD.EXE'],
                                           capture_output=True, timeout=10)
                        except Exception:
                            pass
                        self.log(f"    ⚠️ Timeout 90s — Word ne répond pas, conversion annulée")
                        return file_path, False

                    if _conv_result[0] is not None:
                        self.log(f"    ✅ Converti: {file_path.name} → {new_path.name}")
                        return _conv_result[0], True
                    else:
                        self.log(f"    ❌ Erreur conversion .doc: {_conv_error[0]}")
                        return file_path, False

                # ── .ppt → .pptx ──────────────────────────────────────────
                elif ext == '.ppt':
                    new_path = file_path.with_suffix('.pptx')
                    _conv_result = [None]
                    _conv_error  = [None]

                    def _do_ppt_convert():
                        import win32com.client as _wcc
                        import pythoncom as _pycom
                        _pycom.CoInitialize()
                        _ppt = None
                        try:
                            _ppt = _wcc.DispatchEx("PowerPoint.Application")
                            try:
                                _ppt.DisplayAlerts = 0
                            except Exception:
                                pass
                            try:
                                _ppt.AutomationSecurity = 3
                            except Exception:
                                pass
                            _pres = _ppt.Presentations.Open(
                                str(file_path.absolute()),
                                ReadOnly=True,
                                Untitled=False,
                                WithWindow=False
                            )
                            _pres.SaveAs(str(new_path.absolute()), FileFormat=24)
                            _pres.Close()
                            _conv_result[0] = new_path
                        except Exception as _e:
                            _conv_error[0] = str(_e)
                        finally:
                            try:
                                if _ppt is not None:
                                    _ppt.Quit()
                            except Exception:
                                pass
                            try:
                                _pycom.CoUninitialize()
                            except Exception:
                                pass

                    _t = threading.Thread(target=_do_ppt_convert, daemon=True)
                    _t.start()
                    _t.join(timeout=90)

                    if _t.is_alive():
                        try:
                            subprocess.run(['taskkill', '/F', '/IM', 'POWERPNT.EXE'],
                                           capture_output=True, timeout=10)
                        except Exception:
                            pass
                        self.log(f"    ⚠️ Timeout 90s — PowerPoint ne répond pas, conversion annulée")
                        return file_path, False

                    if _conv_result[0] is not None:
                        self.log(f"    ✅ Converti: {file_path.name} → {new_path.name}")
                        return _conv_result[0], True
                    else:
                        self.log(f"    ❌ Erreur conversion .ppt: {_conv_error[0]}")
                        return file_path, False

                # ── .xls → .xlsx / .xlsm ──────────────────────────────────
                elif ext == '.xls':
                    new_path = file_path.with_suffix('.xlsx')
                    excel = None
                    wb = None
                    xlstart_renamed = []

                    try:
                        # ══════════════════════════════════════════════════
                        # ÉTAPE 1 : Tuer Excel + Neutraliser XLSTART
                        # ══════════════════════════════════════════════════
                        subprocess.run("taskkill /F /IM EXCEL.EXE",
                                       shell=True, capture_output=True)
                        time.sleep(1)

                        # Renommer temporairement TOUS les dossiers XLSTART
                        # pour empêcher le chargement de ReportINI.xls
                        xlstart_renamed = self._disable_xlstart_folders()

                        # ══════════════════════════════════════════════════
                        # ÉTAPE 2 : Lancer Excel en mode sécurisé
                        # ══════════════════════════════════════════════════
                        excel_exe = self._find_excel_exe()

                        if excel_exe:
                            # /automation → désactive add-ins et macros au démarrage
                            # /e          → désactive le classeur personnel (PERSONAL.XLSB)
                            self.log(f"    🔄 Lancement Excel /automation /e")
                            subprocess.Popen(
                                [excel_exe, '/automation', '/e'],
                                creationflags=subprocess.CREATE_NO_WINDOW
                            )
                            time.sleep(4)

                            try:
                                excel = win32com.client.GetActiveObject("Excel.Application")
                            except Exception:
                                self.log(f"    ⚠️ GetActiveObject échoué, fallback DispatchEx")
                                excel = win32com.client.DispatchEx("Excel.Application")
                        else:
                            self.log(f"    ⚠️ Excel introuvable, tentative DispatchEx")
                            excel = win32com.client.DispatchEx("Excel.Application")

                        # ══════════════════════════════════════════════════
                        # ÉTAPE 3 : Supprimer TOUS les popups AVANT Open
                        # ══════════════════════════════════════════════════
                        excel.Visible = False
                        self._suppress_all_excel_alerts(excel)

                        # Désactiver les add-ins COM chargés
                        self._disable_excel_addins(excel)

                        # ══════════════════════════════════════════════════
                        # ÉTAPE 4 : Ouvrir le fichier .xls
                        # ══════════════════════════════════════════════════
                        self.log(f"    📂 Ouverture: {file_path.name}")
                        wb = excel.Workbooks.Open(
                            str(file_path.absolute()),
                            UpdateLinks=0,
                            ReadOnly=False,
                            IgnoreReadOnlyRecommended=True,
                            CorruptLoad=1  # xlNormalLoad
                        )

                        # Remettre les alertes à False après Open
                        self._suppress_all_excel_alerts(excel)

                        # ══════════════════════════════════════════════════
                        # ÉTAPE 5 : Sauvegarder
                        # ══════════════════════════════════════════════════
                        try:
                            wb.SaveAs(str(new_path.absolute()), FileFormat=51)
                            self.log(f"    ✅ Converti: {file_path.name} → {new_path.name}")
                        except Exception as save_err:
                            self.log(f"    ⚠️ xlsx impossible ({save_err}), tentative xlsm...")
                            new_path = file_path.with_suffix('.xlsm')
                            try:
                                self._suppress_all_excel_alerts(excel)
                                wb.SaveAs(str(new_path.absolute()), FileFormat=52)
                                self.log(f"    ✅ Converti (xlsm): {file_path.name} → {new_path.name}")
                            except Exception as save_err2:
                                self.log(f"    ❌ Échec sauvegarde xlsm: {save_err2}")
                                return file_path, False

                        return new_path, True

                    except Exception as e:
                        self.log(f"    ❌ Erreur conversion .xls: {e}")
                        return file_path, False

                    finally:
                        # ══════════════════════════════════════════════════
                        # NETTOYAGE : Fermer Excel + Restaurer XLSTART
                        # ══════════════════════════════════════════════════
                        try:
                            if wb is not None:
                                wb.Close(False)
                        except Exception:
                            pass
                        try:
                            if excel is not None:
                                excel.DisplayAlerts = False
                                excel.Quit()
                        except Exception:
                            pass
                        time.sleep(1)
                        subprocess.run("taskkill /F /IM EXCEL.EXE",
                                       shell=True, capture_output=True)

                        # Restaurer les dossiers XLSTART renommés
                        self._restore_xlstart_folders(xlstart_renamed)

                return file_path, False

            finally:
                pythoncom.CoUninitialize()

        except ImportError:
            self.log(f"    ⚠️ Module pywin32 requis pour la conversion")
            return file_path, False
        except Exception as e:
            self.log(f"    ⚠️ Erreur conversion: {e}")
            return file_path, False

    # ──────────────────────────────────────────────────────────────────────
    # HELPERS ANTI-POPUP
    # ──────────────────────────────────────────────────────────────────────

    def _suppress_all_excel_alerts(self, excel):
        """Désactive TOUS les popups et alertes Excel"""
        for attr, val in [
            ('DisplayAlerts', False),
            ('AskToUpdateLinks', False),
            ('EnableEvents', False),
            ('Interactive', False),           # empêche les popups interactifs
            ('ScreenUpdating', False),
        ]:
            try:
                setattr(excel, attr, val)
            except Exception:
                pass
        try:
            excel.AutomationSecurity = 3  # msoAutomationSecurityForceDisable
        except Exception:
            pass

    def _disable_excel_addins(self, excel):
        """Désactive les add-ins COM qui pourraient afficher des popups"""
        try:
            for i in range(excel.AddIns.Count, 0, -1):
                try:
                    addin = excel.AddIns.Item(i)
                    if addin.Installed:
                        addin.Installed = False
                        self.log(f"      🔕 Add-in désactivé: {addin.Name}")
                except Exception:
                    pass
        except Exception:
            pass

        # Désactiver aussi les COM Add-ins
        try:
            for i in range(excel.COMAddIns.Count, 0, -1):
                try:
                    comaddin = excel.COMAddIns.Item(i)
                    if comaddin.Connect:
                        comaddin.Connect = False
                        self.log(f"      🔕 COM Add-in déconnecté: {comaddin.Description}")
                except Exception:
                    pass
        except Exception:
            pass

    def _disable_xlstart_folders(self):
        """
        Renomme temporairement les dossiers XLSTART pour empêcher
        le chargement automatique de fichiers comme ReportINI.xls.
        Retourne la liste des (original_path, renamed_path) pour restauration.
        """
        renamed = []

        # Chercher tous les emplacements possibles de XLSTART
        xlstart_candidates = []

        # 1. Via variable d'environnement APPDATA
        appdata = os.environ.get('APPDATA', '')
        if appdata:
            xlstart_candidates.append(
                Path(appdata) / 'Microsoft' / 'Excel' / 'XLSTART'
            )

        # 2. Dans les dossiers Office classiques
        program_files = [
            os.environ.get('ProgramFiles', r'C:\Program Files'),
            os.environ.get('ProgramFiles(x86)', r'C:\Program Files (x86)')
        ]
        for pf in program_files:
            if pf:
                for version in ['Office16', 'Office15', 'Office14']:
                    xlstart_candidates.append(
                        Path(pf) / 'Microsoft Office' / 'root' / version / 'XLSTART'
                    )
                    xlstart_candidates.append(
                        Path(pf) / 'Microsoft Office' / version / 'XLSTART'
                    )

        # 3. Chercher via le registre (AltStartup)
        try:
            import winreg
            for hive in [winreg.HKEY_CURRENT_USER, winreg.HKEY_LOCAL_MACHINE]:
                try:
                    key = winreg.OpenKey(
                        hive,
                        r'Software\Microsoft\Office\16.0\Excel\Options'
                    )
                    alt_startup, _ = winreg.QueryValueEx(key, 'AltStartup')
                    winreg.CloseKey(key)
                    if alt_startup and Path(alt_startup).exists():
                        xlstart_candidates.append(Path(alt_startup))
                except Exception:
                    pass
        except ImportError:
            pass

        # Renommer chaque dossier XLSTART trouvé
        for xlstart_path in xlstart_candidates:
            if xlstart_path.exists() and xlstart_path.is_dir():
                disabled_path = xlstart_path.with_name('XLSTART_DISABLED_TEMP')
                try:
                    # Ne pas renommer si déjà renommé
                    if disabled_path.exists():
                        continue
                    xlstart_path.rename(disabled_path)
                    renamed.append((disabled_path, xlstart_path))
                    self.log(f"      🔇 XLSTART neutralisé: {xlstart_path}")
                except PermissionError:
                    self.log(f"      ⚠️ Permission refusée pour {xlstart_path}")
                except Exception as e:
                    self.log(f"      ⚠️ Impossible de renommer {xlstart_path}: {e}")

        return renamed

    def _restore_xlstart_folders(self, renamed_list):
        """Restaure les dossiers XLSTART renommés"""
        for disabled_path, original_path in renamed_list:
            try:
                if disabled_path.exists():
                    disabled_path.rename(original_path)
                    self.log(f"      🔊 XLSTART restauré: {original_path}")
            except Exception as e:
                self.log(f"      ⚠️ Impossible de restaurer {original_path}: {e}")
    def _find_excel_exe(self):
        """Trouve le chemin de l'exécutable Excel installé"""
        import winreg
        try:
            # Chercher dans le registre
            key = winreg.OpenKey(
                winreg.HKEY_LOCAL_MACHINE,
                r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\excel.exe"
            )
            path, _ = winreg.QueryValueEx(key, "")
            winreg.CloseKey(key)
            if path and Path(path).exists():
                return path
        except Exception:
            pass

        # Chemins courants Office
        candidates = [
            r"C:\Program Files\Microsoft Office\root\Office16\EXCEL.EXE",
            r"C:\Program Files\Microsoft Office\Office16\EXCEL.EXE",
            r"C:\Program Files (x86)\Microsoft Office\root\Office16\EXCEL.EXE",
            r"C:\Program Files (x86)\Microsoft Office\Office16\EXCEL.EXE",
            r"C:\Program Files\Microsoft Office\root\Office15\EXCEL.EXE",
            r"C:\Program Files\Microsoft Office\Office15\EXCEL.EXE",
        ]
        for c in candidates:
            if Path(c).exists():
                return c

        return None
    def get_ole_order_from_pptx(self, pptx_path):
        """Détermine l'ordre réel des objets OLE dans un PPTX"""
        ole_order = []
        
        try:
            with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                slide_files = [f for f in zip_ref.namelist() 
                              if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                slide_files.sort(key=lambda x: int(re.search(r'slide(\d+)', x).group(1)))
                
                for slide_file in slide_files:
                    slide_num = int(re.search(r'slide(\d+)', slide_file).group(1))
                    
                    xml_content = zip_ref.read(slide_file)
                    
                    ole_matches = re.findall(r'<p:oleObj[^>]*r:id="([^"]+)"', 
                                            xml_content.decode('utf-8'))
                    
                    for idx, r_id in enumerate(ole_matches):
                        ole_order.append({
                            'slide': slide_num,
                            'r_id': r_id,
                            'position_in_slide': idx + 1
                        })
                
                for slide_file in slide_files:
                    slide_num = int(re.search(r'slide(\d+)', slide_file).group(1))
                    rels_file = slide_file.replace('slides/slide', 'slides/_rels/slide').replace('.xml', '.xml.rels')
                    
                    if rels_file in zip_ref.namelist():
                        rels_xml = zip_ref.read(rels_file)
                        root = ET.fromstring(rels_xml)
                        
                        ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                        
                        for relationship in root.findall('r:Relationship', ns):
                            r_id = relationship.get('Id')
                            target = relationship.get('Target')
                            
                            if 'embeddings/' in target:
                                match = re.search(r'oleObject(\d+)', target)
                                if match:
                                    file_num = int(match.group(1))
                                    
                                    for ole_info in ole_order:
                                        if ole_info['slide'] == slide_num and ole_info['r_id'] == r_id:
                                            ole_info['file_num'] = file_num
            
            ole_order_sorted = sorted(
                [o for o in ole_order if 'file_num' in o],
                key=lambda x: (x['slide'], x['position_in_slide'])
            )
            
            file_num_to_doc_position = {}
            for doc_position, ole_info in enumerate(ole_order_sorted, 1):
                file_num_to_doc_position[ole_info['file_num']] = doc_position
            
            if file_num_to_doc_position:
                self.log(f"    📋 Ordre PowerPoint: {file_num_to_doc_position}")
            
            return file_num_to_doc_position
        
        except Exception as e:
            self.log(f"    ⚠ Erreur ordre PPTX: {e}")
            return {}
    
    def get_ole_order_from_docx(self, docx_path):
        """Détermine l'ordre réel des objets OLE dans un DOCX"""
        ole_order = []
        
        try:
            with zipfile.ZipFile(docx_path, 'r') as zip_ref:
                doc_xml = zip_ref.read('word/document.xml')
                root = ET.fromstring(doc_xml)
                
                position = 0
                for obj in root.iter('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}object'):
                    position += 1
                    ole_obj = obj.find('.//{urn:schemas-microsoft-com:office:office}OLEObject')
                    if ole_obj is not None:
                        r_id = ole_obj.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                        if r_id:
                            ole_order.append({'position': position, 'r_id': r_id})
                
                if 'word/_rels/document.xml.rels' in zip_ref.namelist():
                    rels_xml = zip_ref.read('word/_rels/document.xml.rels')
                    rels_root = ET.fromstring(rels_xml)
                    
                    ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                    
                    for relationship in rels_root.findall('r:Relationship', ns):
                        r_id = relationship.get('Id')
                        target = relationship.get('Target')
                        
                        if 'embeddings/' in target:
                            match = re.search(r'oleObject(\d+)', target)
                            if match:
                                file_num = int(match.group(1))
                                
                                for ole_info in ole_order:
                                    if ole_info['r_id'] == r_id:
                                        ole_info['file_num'] = file_num
            
            file_num_to_doc_position = {}
            for ole_info in ole_order:
                if 'file_num' in ole_info:
                    file_num_to_doc_position[ole_info['file_num']] = ole_info['position']
            
            if file_num_to_doc_position:
                self.log(f"    📋 Ordre Word: {file_num_to_doc_position}")
            
            return file_num_to_doc_position
        
        except Exception as e:
            self.log(f"    ⚠ Erreur ordre DOCX: {e}")
            return {}
    def get_max_fj_number_in_directory(self, directory):
        """Trouve le plus grand numéro FJ dans un dossier"""
        max_num = 0
        
        try:
            directory = Path(directory)
            if not directory.exists():
                return 0
            
            for file_path in directory.iterdir():
                if file_path.is_file():
                    match = re.search(r'_FJ_(\d+)', file_path.stem)
                    if match:
                        num = int(match.group(1))
                        if num > max_num:
                            max_num = num
            
            if max_num > 0:
                self.log(f"    📊 Plus grand FJ trouvé: FJ_{max_num}")
            return max_num
        
        except Exception as e:
            self.log(f"    ⚠ Erreur détection FJ max: {e}")
            return 0
    
    def recursive_extract_from_directory(self, directory, max_depth=3):
        """Extraction récursive automatique sur tous les fichiers d'un dossier"""
        directory = Path(directory)
        
        if not hasattr(self, 'recursion_depth'):
            self.recursion_depth = 0
        
        if self.recursion_depth >= max_depth:
            self.log(f"  ⚠️ Profondeur maximale {max_depth} atteinte")
            return
        
        # CORRECTION : Suivre les fichiers déjà traités par signature unique
        if not hasattr(self, 'processed_files_signatures'):
            self.processed_files_signatures = set()
        
        self.recursion_depth += 1
        indent = "  " * self.recursion_depth
        
        self.log(f"{indent}🔄 EXTRACTION RÉCURSIVE niveau {self.recursion_depth}: {directory.name}")
        
        # Trouver le max FJ avant extraction
        max_fj = self.get_max_fj_number_in_directory(directory)
        
        SUPPORTED_EXTENSIONS = {
            '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls', 
            '.pptx', '.pptm''.ppt', '.msg', '.zip','.htm', '.html' 
        }
        
        # Extensions à exclure explicitement (ne jamais extraire)
        EXCLUDED_EXTENSIONS = {
        '.equation',      # ← NOUVEAU : formules mathématiques Word
        '.bitmap_ole',    # ← NOUVEAU : bitmaps OLE (Paint, MSBitmap...)
         '.ole', '.bin', '.dat', '.tmp','txt'
        # CATIA
        '.catpart', '.catproduct', '.catdrawing',
        '.catanalysis', '.catprocess', '.catmaterial',
        '.cgr', '.model', '.exp', '.dlv', '.dlv3', '.catalog',
        # Autres formats CAO non archivables
        '.pmat', '.rpx', '.3dxml',
        }
        
        # ========================================
        # ÉTAPE 1 : CONVERTIR LES FICHIERS LEGACY
        # ========================================
        legacy_files = []
        for file_path in directory.iterdir():
            if file_path.is_file() and file_path.suffix.lower() in ['.doc', '.xls', '.ppt']:
                if not file_path.stem.endswith('_original'):
                    legacy_files.append(file_path)
        
        if legacy_files:
            self.log(f"{indent}  🔄 {len(legacy_files)} fichier(s) legacy à convertir")
            
            for legacy_file in legacy_files:
                self.log(f"{indent}    📄 Conversion: {legacy_file.name}")
                
                try:
                    # Appeler la conversion
                    conversion_result = self.convert_legacy_to_modern_format(legacy_file)
                    
                    # CORRECTION CRITIQUE : Gérer le tuple retourné
                    if conversion_result:
                        if isinstance(conversion_result, tuple):
                            converted_file = conversion_result[0]
                        else:
                            converted_file = conversion_result

                        if converted_file and converted_file.exists():
                            self.log(f"{indent}      ✅ Converti: {legacy_file.name} → {converted_file.name}")
                            # Supprimer le fichier legacy original du dossier
                            try:
                                legacy_file.unlink()
                                self.log(f"{indent}      🗑️ Supprimé: {legacy_file.name}")
                            except Exception as del_err:
                                self.log(f"{indent}      ⚠️ Impossible de supprimer {legacy_file.name}: {del_err}")
                        else:
                            self.log(f"{indent}      ⚠️ Conversion échouée: {legacy_file.name}")
                    else:
                        self.log(f"{indent}      ⚠️ Conversion impossible: {legacy_file.name}")
                        
                except Exception as conv_error:
                    self.log(f"{indent}      ❌ Erreur conversion {legacy_file.name}: {conv_error}")
        
        # ========================================
        # ÉTAPE 2 : COLLECTER LES FICHIERS À TRAITER
        # ========================================
        files_to_process = []
        for file_path in directory.iterdir():
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                if not file_path.stem.endswith('_original'):
                    # CORRECTION : Créer une signature unique pour chaque fichier
                    try:
                        file_stat = file_path.stat()
                        
                        # Lire les premiers 2048 octets pour créer une empreinte
                        with open(file_path, 'rb') as f:
                            file_start = f.read(2048)
                        
                        # Créer une signature unique basée sur : nom + taille + date + hash du contenu
                        import hashlib
                        file_hash = hashlib.md5(file_start).hexdigest()
                        file_signature = f"{file_path.name}|{file_stat.st_size}|{file_stat.st_mtime:.0f}|{file_hash}"
                        
                        # Vérifier si ce fichier a déjà été traité
                        if file_signature not in self.processed_files_signatures:
                            files_to_process.append((file_path, file_signature))
                            self.log(f"{indent}  📝 Fichier nouveau: {file_path.name}")
                        else:
                            self.log(f"{indent}  ⏭️ Fichier déjà traité (signature identique): {file_path.name}")
                    
                    except Exception as sig_error:
                        self.log(f"{indent}  ⚠️ Erreur lecture signature pour {file_path.name}: {sig_error}")
                        # En cas d'erreur, utiliser juste le chemin absolu comme signature
                        file_signature = str(file_path.absolute())
                        if file_signature not in self.processed_files_signatures:
                            files_to_process.append((file_path, file_signature))
        
        if not files_to_process:
            self.log(f"{indent}  ℹ️ Aucun nouveau fichier à traiter")
            self.recursion_depth -= 1
            return
        
        self.log(f"{indent}  📋 {len(files_to_process)} nouveau(x) fichier(s) à analyser")
        
        # Variable pour détecter si de nouveaux fichiers ont été extraits
        new_files_extracted = False
        total_extracted_files = []
        
        # ========================================
        # ÉTAPE 3 : TRAITER CHAQUE FICHIER
        # ========================================
        for file_path, file_signature in files_to_process:
            # CORRECTION : Marquer le fichier comme traité AVANT de le traiter
            self.processed_files_signatures.add(file_signature)
            
            self.log(f"{indent}  🔍 Analyse: {file_path.name}")
            
            # ========================================
            # NOUVEAU : NETTOYAGE .BIN AVANT EXTRACTION RÉCURSIVE
            # ========================================
            if file_path.suffix.lower() in ['.docx', '.pptx', '.xlsx', '.xlsm']:
                try:
                    bin_cleaned = self.clean_bin_files_from_embeddings(file_path)
                    if bin_cleaned > 0:
                        self.log(f"{indent}    ✅ Pré-nettoyage: {bin_cleaned} fichier(s) .bin vide(s) supprimé(s)")
                except Exception as clean_error:
                    self.log(f"{indent}    ⚠️ Erreur nettoyage .bin: {clean_error}")
            
            # IMPORTANT: Ajuster extracted_count pour continuer la numérotation
            self.extracted_count = max_fj + 1
            
            # Extraire
            extracted_files = self.process_file(file_path, directory)
            
            # Afficher les logs avec indentation
            original_log_messages = self.log_messages.copy()
            self.log_messages = []
            
            for log_msg in original_log_messages:
                self.log(f"{indent}  {log_msg}")
            
            if extracted_files:
                # CORRECTION : Compter seulement les fichiers vraiment extraits (pas vides, pas non supportés)
                extracted_count = 0
                for item in extracted_files:
                    if isinstance(item, tuple) and len(item) >= 2:
                        filename = item[1] if len(item) >= 2 else item
                        
                        # Ne compter QUE les fichiers réellement extraits
                        # Exclure : fichiers vides, fichiers non supportés
                        if isinstance(filename, str):
                            if "vide" not in filename.lower() and "supporté" not in filename.lower():
                                extracted_count += 1
                                total_extracted_files.append(item)
                
                if extracted_count > 0:
                    self.log(f"{indent}  ✅ {extracted_count} fichier(s) extrait(s) de {file_path.name}")
                    new_files_extracted = True
                    
                    # Mettre à jour le max FJ
                    new_max = self.get_max_fj_number_in_directory(directory)
                    if new_max > max_fj:
                        max_fj = new_max
                else:
                    self.log(f"{indent}  ℹ️ Aucun fichier avec contenu extrait de {file_path.name}")
        
        # ========================================
        # ÉTAPE 4 : DÉCIDER SI ON CONTINUE LA RÉCURSION
        # ========================================
        # Compter le nombre total de fichiers réellement extraits
        files_with_content = len(total_extracted_files)
        
        if files_with_content > 0 and new_files_extracted and self.recursion_depth < max_depth:
            self.log(f"{indent}  🔄 {files_with_content} fichier(s) avec contenu extrait(s), poursuite de la récursion...")
            self.recursive_extract_from_directory(directory, max_depth)
        else:
            if not new_files_extracted:
                self.log(f"{indent}  ✓ Aucun nouveau fichier extrait, arrêt de la récursion")
            elif files_with_content == 0:
                self.log(f"{indent}  ✓ Fichiers sans contenu incorporé, arrêt de la récursion")
            else:
                self.log(f"{indent}  ✓ Profondeur maximale atteinte ou plus de fichiers à traiter")
        
        self.recursion_depth -= 1
    def _generate_output_name(self, base_name, original_filename, file_ext):
        """
        Gère la nomenclature avec PRIORITÉ AU FICHIER DE RÉFÉRENCE / INDEX MANUEL

        RÈGLE DE PRIORITÉ :
        1. Index manuel saisi dans l'interface  → priorité absolue
        2. Fichier de référence Excel           → par fichier source
        3. Nom se terminant par _FJ_X           → utiliser X + position
        4. Sinon                                → extracted_count normal

        CAS 3 (sans _FJ_X) : on remplace juste le compteur par l'index,
        SANS ajouter "_FJ_" (évite les doublons comme DER-259935_1_FJ_FJ_6)
        """

        # ====================================================================
        # ÉTAPE 0 : RÉCUPÉRER LA BONNE EXTENSION depuis original_filename
        # ====================================================================
        if original_filename and original_filename.strip() and original_filename != "embedded_file":
            original_ext = Path(original_filename).suffix.lower()
            if original_ext and original_ext not in ['', '.bin', '.dat', '.tmp']:
                file_ext = original_ext

        # ====================================================================
        # ÉTAPE 1 : DÉTERMINER reference_start_index
        #   Priorité : index manuel > fichier de référence
        # ====================================================================
        reference_start_index = None

        if self.manual_start_index is not None:
            reference_start_index = self.manual_start_index
            self.log(f"      🔢 Index manuel utilisé : départ FJ_{reference_start_index}")

        elif self.indexation_manager and self.current_source_file:
            reference_start_index = self.indexation_manager.get_start_index(
                self.current_source_file
            )

        # ====================================================================
        # ÉTAPE 2 : GÉNÉRER LE NOM DE SORTIE
        # ====================================================================

        # Vérifier si le nom se termine déjà par _FJ_X  (ex: DER-362949_2_FJ_4)
        fj_pattern = r'^(.*_FJ)_(\d+)$'
        match = re.match(fj_pattern, base_name)

        if match:
            # ================================================================
            # CAS 1 & 2 : Nom avec _FJ_X  → remplacer X par le bon numéro
            # ================================================================
            prefix = match.group(1)                    # ex: "DER-362949_2_FJ"
            file_number_in_name = int(match.group(2))  # ex: 4

            if reference_start_index is not None:
                next_num = reference_start_index + (self.extracted_count - 1)
                self.log(f"      ✅ INDEX APPLIQUÉ : FJ_{next_num} "
                        f"(base={reference_start_index}, position={self.extracted_count})")
            else:
                # Prochain numéro séquentiel après le dernier assigné
                next_num = self._fj_floor + 1
                self.log(f"      ℹ️ FICHIER : FJ_{next_num} "
                        f"(après plancher FJ_{self._fj_floor})")

            # ── Éviter collision avec un FJ déjà utilisé ─────────────────
            while next_num in self._used_fj_numbers:
                next_num += 1
                self.log(f"      🔒 COLLISION ÉVITÉE → FJ_{next_num}")

            self._used_fj_numbers.add(next_num)
            self._fj_floor = max(self._fj_floor, next_num)

            # Résultat : DER-362949_2_FJ_5.pdf  (pas de double _FJ_)
            output_name = f"{prefix}_{next_num}{file_ext}"

        else:
            # ================================================================
            # CAS 3 : Nom sans _FJ_X  (ex: DER-259935_1_FJ  ou  Document_1)
            # On remplace simplement le compteur par l'index —
            # SANS ajouter "_FJ_" pour éviter DER-259935_1_FJ_FJ_6
            # ================================================================

            if reference_start_index is not None:
                next_num = reference_start_index + (self.extracted_count - 1)
                self.log(f"      ✅ INDEX APPLIQUÉ : _{next_num} "
                        f"(base={reference_start_index}, position={self.extracted_count})")
            else:
                next_num = self.extracted_count

            # ── Plancher global ───────────────────────────────────────────
            if next_num <= self._fj_floor:
                next_num = self._fj_floor + 1
                self.log(f"      🔒 PLANCHER APPLIQUÉ → FJ_{next_num} (évite collision)")

            # ── Éviter collision avec un FJ déjà utilisé ─────────────────
            while next_num in self._used_fj_numbers:
                next_num += 1
                self.log(f"      🔒 COLLISION ÉVITÉE → FJ_{next_num}")

            self._used_fj_numbers.add(next_num)
            self._fj_floor = next_num

            output_name = f"{base_name}_{next_num}{file_ext}"

        return output_name
    def extract_ole10native(self, ole):
        """Extrait le contenu d'un objet Ole10Native avec extraction robuste du nom de fichier"""
        try:
            stream = None
            if ole.exists('\x01Ole10Native'):
                stream = ole.openstream('\x01Ole10Native')
            elif ole.exists('Ole10Native'):
                stream = ole.openstream('Ole10Native')
            else:
                return None
            
            data = stream.read()
            
            if len(data) < 20:
                return None
            
            offset = 0
            filename = "embedded_file"
            original_filename = None
            
            # ══════════════════════════════════════════════════════════════════
            # HELPER : vérifie si des données sont un 7z et corrige le nom
            # ══════════════════════════════════════════════════════════════════
            def fix_7z_filename(file_data, orig_name):
                """Si file_data est un 7z, force l'extension .7z sur orig_name"""
                if file_data and len(file_data) >= 6 and file_data[:6] == b'\x37\x7a\xbc\xaf\x27\x1c':
                    if orig_name:
                        stem = orig_name
                        if '.' in orig_name:
                            stem = orig_name.rsplit('.', 1)[0]
                        corrected = stem + '.7z'
                        self.log(f"      🔄 OVERRIDE 7z : '{orig_name}' → '{corrected}'")
                        return corrected
                    return "embedded_file.7z"
                return orig_name
            
            # ──────────────────────────────────────────────────────────────
            # MÉTHODE 1 : Format Ole10Native réel
            # [4] uint32 total_size
            # [2] uint16 type (02 00 = embedded)
            # [N] ANSI+\x00 → nom fichier original
            # [N] ANSI+\x00 → chemin fichier original
            # [4] octets inconnus
            # [4] uint32 longueur chemin temp
            # [N] chemin temp
            # [4] uint32 taille données réelles
            # [N] DONNÉES RÉELLES ← ce qu'on veut
            # ──────────────────────────────────────────────────────────────
            try:
                if len(data) > 10:
                    offset = 0

                    # Sauter uint32 total_size
                    total_sz = struct.unpack_from('<I', data, offset)[0]
                    if 10 < total_sz <= len(data) + 8:
                        offset = 4

                    # Lire type uint16 (optionnel, certains streams l'omettent)
                    type_val = struct.unpack_from('<H', data, offset)[0]
                    if type_val in (1, 2):
                        offset += 2

                    # Nom fichier original (ANSI, null-terminated)
                    end = data.find(b'\x00', offset)
                    if 0 < end - offset < 260:
                        candidate = data[offset:end].decode('latin-1', errors='ignore').strip()
                        if candidate and '.' in candidate:
                            original_filename = candidate
                            self.log(f"      📝 Nom extrait (méthode 1): {original_filename}")
                        offset = end + 1
                    else:
                        raise ValueError("Nom fichier introuvable")

                    # Chemin original (ANSI, null-terminated)
                    end = data.find(b'\x00', offset)
                    if end == -1 or end - offset > 1000:
                        raise ValueError("Chemin original introuvable")
                    offset = end + 1

                    # 4 octets inconnus
                    offset += 4

                    # Longueur chemin temp (uint32)
                    if offset + 4 > len(data):
                        raise ValueError("Chemin temp manquant")
                    temp_len = struct.unpack_from('<I', data, offset)[0]
                    offset += 4
                    offset += temp_len  # sauter le chemin temp

                    # Taille données réelles (uint32)
                    if offset + 4 > len(data):
                        raise ValueError("Taille données manquante")
                    data_size = struct.unpack_from('<I', data, offset)[0]
                    offset += 4

                    if data_size == 0:
                        raise ValueError("Données vides")

                    # Extraire les données réelles
                    if offset + data_size <= len(data):
                        file_data = data[offset:offset + data_size]
                    else:
                        file_data = data[offset:]  # prendre ce qui reste

                    self.log(f"      ✓ Ole10Native parsé : {len(file_data):,} bytes réels / {len(data):,} bruts")

                    # Override 7z si nécessaire
                    corrected = fix_7z_filename(file_data, original_filename or filename)
                    if corrected != (original_filename or filename):
                        original_filename = corrected

                    if original_filename:
                        return file_data, original_filename
                    return file_data, filename

            except Exception as parse_error:
                self.log(f"      ⚠️ Méthode 1 échouée: {parse_error}")

            # ──────────────────────────────────────────────────────────────
            # MÉTHODE 2 : Recherche du nom par pattern
            # ──────────────────────────────────────────────────────────────
            if not original_filename:
                try:
                    import re
                    pattern = rb'([A-Za-z0-9_\-\.]+\.(7z|pmat|pro|tab|txt|dat|csv|ini|cfg|xml|rpx|json|zip|pdf|msg|docx|xlsx|pptx))'
                    matches = re.findall(pattern, data[:2048], re.IGNORECASE)
                    
                    if matches:
                        for match in matches:
                            candidate = match[0].decode('ascii', errors='ignore')
                            if len(candidate) > 3 and '.' in candidate:
                                original_filename = candidate
                                self.log(f"      📝 Nom détecté (pattern): {original_filename}")
                                break
                except Exception as pattern_error:
                    self.log(f"      ⚠️ Recherche pattern échouée: {pattern_error}")
            
            # ──────────────────────────────────────────────────────────────
            # MÉTHODE 3 : Recherche de signatures connues
            # ──────────────────────────────────────────────────────────────
            self.log(f"      🔍 Recherche signatures...")
            signatures = [
                (b'\x37\x7a\xbc\xaf\x27\x1c', '7z'),
                (b'\xd0\xcf\x11\xe0', 'OLE/MSG'),
                (b'%PDF', 'PDF'),
                (b'PK\x03\x04', 'ZIP/Office'),
                (b'{\rtf', 'RTF'),
            ]
            
            for start in range(0, min(len(data), 2000)):
                for sig, sig_name in signatures:
                    if data[start:start+len(sig)] == sig:
                        file_data = data[start:]
                        self.log(f"      ✓ Signature {sig_name} à {start}")
                        
                        if sig_name == '7z':
                            corrected = fix_7z_filename(file_data, original_filename or filename)
                            if corrected != (original_filename or filename):
                                original_filename = corrected
                        
                        if original_filename:
                            return file_data, original_filename
                        return file_data, filename
            
            # ──────────────────────────────────────────────────────────────
            # MÉTHODE 4 : Vérifier si c'est du texte
            # ──────────────────────────────────────────────────────────────
            if len(data) > 100:
                
                # Override 7z avant test texte (latin-1 accepte tout)
                if data[:6] == b'\x37\x7a\xbc\xaf\x27\x1c':
                    self.log(f"      ✓ 7z détecté (magic bytes avant test texte)")
                    corrected = fix_7z_filename(data, original_filename or filename)
                    return data, corrected
                
                try:
                    test_text = data[:500].decode('utf-8', errors='strict')
                    self.log(f"      ✓ Fichier texte détecté (UTF-8)")
                    if original_filename:
                        return data, original_filename
                    return data, filename
                except:
                    pass
                
                for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                    try:
                        test_text = data[:500].decode(encoding, errors='strict')
                        self.log(f"      ✓ Fichier texte détecté ({encoding})")
                        if original_filename:
                            return data, original_filename
                        return data, filename
                    except:
                        continue
            
            return None
            
        except Exception as e:
            self.log(f"      ⚠ Erreur Ole10Native: {e}")
            return None
    def extract_file_from_package(self, package_data):
        """Extrait fichier du Package"""
        try:
            signatures = [
                (b'%PDF', '.pdf'),
                (b'PK\x03\x04', '.zip'),
                (b'\xd0\xcf\x11\xe0', '.ole'),
                (b'{\rtf', '.rtf'),
            ]
            
            for signature, ext in signatures:
                pos = package_data.find(signature, 0, 2000)
                if pos >= 0:
                    return package_data[pos:]
            
            if len(package_data) > 260:
                for start_pos in range(0, min(512, len(package_data)), 4):
                    if package_data[start_pos:start_pos+4] == b'\x00\x00\x00\x00':
                        content_start = start_pos
                        while content_start < len(package_data) and package_data[content_start] == 0:
                            content_start += 1
                        
                        if content_start < len(package_data) - 100:
                            test_data = package_data[content_start:]
                            for signature, ext in signatures:
                                if test_data.startswith(signature):
                                    return test_data
            
            file_type = self.detect_file_type(package_data)
            if file_type not in ['.bin', '.dat']:
                return package_data
            
            for offset in [20, 24, 28, 32, 36, 40, 48, 64, 78, 128, 256, 512]:
                if offset < len(package_data):
                    test_data = package_data[offset:]
                    for signature, ext in signatures:
                        if test_data.startswith(signature):
                            return test_data
            
            return None
            
        except Exception as e:
            return None
    
    def extract_actual_content(self, data):
        """Extrait contenu réel"""
        signatures = {
            b'PK\x03\x04': 0,
            b'%PDF': 0,
            b'\xd0\xcf\x11\xe0': 0,
            b'{\rtf': 0,
        }
        
        for signature, offset in signatures.items():
            pos = data.find(signature)
            if pos >= 0:
                return data[pos:]
        
        return data
    
    def detect_file_type(self, data):
        """Détecte type fichier"""
        if not data or len(data) < 4:
            return '.bin'
        
        signatures = {
            b'PK\x03\x04': None,
            b'\xd0\xcf\x11\xe0': None,
            b'%PDF': '.pdf',
            b'{\\rtf': '.rtf',
            b'\xff\xd8\xff': '.jpg',
            b'\x89PNG': '.png',
            b'GIF87a': '.gif',
            b'GIF89a': '.gif',
            b'\x37\x7a\xbc\xaf\x27\x1c': '.7z',
        }
        
        for signature, ext in signatures.items():
            if data[:len(signature)] == signature:
                if ext:
                    return ext
                
                if signature.startswith(b'PK'):
                    return self.detect_office_xml_type(data)
                
                if signature == b'\xd0\xcf\x11\xe0':
                    return self.detect_ole_type(data)
        
        try:
            test_sample = data[:min(2000, len(data))]
            decoded = test_sample.decode('utf-8', errors='strict')
            
            printable_ratio = sum(c.isprintable() or c.isspace() for c in decoded) / len(decoded)
            
            if printable_ratio > 0.8:
                if decoded.lstrip().startswith('<?xml'):
                    return '.xml'
                elif decoded.lstrip().startswith('{') or decoded.lstrip().startswith('['):
                    return '.json'
                elif decoded.lstrip().startswith('<html') or decoded.lstrip().startswith('<!DOCTYPE'):
                    return '.html'
                elif '\t' in decoded[:500] or ',' in decoded[:500]:
                    lines = decoded.split('\n')[:5]
                    if any(',' in line or '\t' in line for line in lines):
                        return '.csv'
                else:
                    return '.txt'  # ← CORRECTION : retourner .txt normal
        except UnicodeDecodeError:
            pass
        
        for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
            try:
                test_sample = data[:min(2000, len(data))]
                decoded = test_sample.decode(encoding, errors='strict')
                
                printable_ratio = sum(c.isprintable() or c.isspace() for c in decoded) / len(decoded)
                
                if printable_ratio > 0.8:
                    return '.txt'  # ← CORRECTION : retourner .txt normal
            except:
                continue

        return '.bin'  # ← CORRECTION : retourner .bin normal
    
    def detect_office_xml_type(self, data):
        """Détecte type Office XML — distingue pptx/pptm via vbaProject.bin"""
        try:
            import io, zipfile

            # Tentative d'analyse ZIP pour détection précise pptm vs pptx
            try:
                with zipfile.ZipFile(io.BytesIO(data)) as zf:
                    names = zf.namelist()
                    if any('word/' in n for n in names):
                        return '.docx'
                    elif any('xl/' in n for n in names):
                        return '.xlsx'
                    elif any('ppt/' in n for n in names):
                        # pptm contient un module VBA → vbaProject.bin
                        if 'ppt/vbaProject.bin' in names:
                            return '.pptm'
                        return '.pptx'
                    return '.zip'
            except Exception:
                # Fallback : analyse par magic bytes bruts (ZIP partiellement lisible)
                content_check = data[:4096]
                if b'word/' in content_check:
                    return '.docx'
                elif b'xl/' in content_check:
                    return '.xlsx'
                elif b'ppt/' in content_check:
                    if b'vbaProject.bin' in content_check:
                        return '.pptm'
                    return '.pptx'
                return '.zip'

        except Exception:
            return '.zip'
        
    def detect_ole_type(self, data):
        """Détecte type OLE avec lecture du vrai nom depuis Ole10Native/CompObj.
        Ignore silencieusement bitmaps OLE et formules mathématiques."""

        # ── Mapping ProgID → extension ────────────────────────────────────
        PROGID_TO_EXT = {
            # CATIA
            "catia.part":        ".catpart",
            "catia.product":     ".catproduct",
            "catia.drawing":     ".catdrawing",
            "catia.analysis":    ".catanalysis",
            "catia.process":     ".catprocess",
            "catia.material":    ".catmaterial",
            "catia.catalog":     ".catalog",
            # SolidWorks
            "solidworks.part":   ".sldprt",
            "solidworks.assembly": ".sldasm",
            "solidworks.drawing":  ".slddrw",
            # AutoCAD
            "autocad.drawing":   ".dwg",
            "autocad.r18":       ".dwg",
            # NX / Unigraphics
            "nx.part":           ".prt",
            # Pro/E Creo
            "proe.part":         ".prt",
            # RPX
            "rpx":               ".rpx",
            # 3DXML
            "3dxml":             ".3dxml",
            # Images / formules (à ignorer)
            "equation.3":        ".equation",
            "equation.2":        ".equation",
            "equation.1":        ".equation",
            "mathtype":          ".equation",
            "ms equation":       ".equation",
            "paint.picture":     ".bitmap_ole",
            "paintbrush":        ".bitmap_ole",
            "msbitmap":          ".bitmap_ole",
            "ms_clipart":        ".bitmap_ole",
            "msdraw":            ".bitmap_ole",
            "photoeditor":       ".bitmap_ole",
        }

        try:
            with tempfile.NamedTemporaryFile(suffix=".ole", delete=False) as temp:
                temp.write(data[:min(len(data), 1_000_000)])
                temp_path = temp.name

            try:
                ole = olefile.OleFileIO(temp_path)
                streams = ole.listdir()
                stream_list = ["/".join(entry) for entry in streams]

                # ════════════════════════════════════════════════════════
                # ÉTAPE 0 : CompObj → ProgID
                # ════════════════════════════════════════════════════════
                progid_detected = None
                for comp_name in ["\\x01CompObj", "CompObj"]:
                    if ole.exists(comp_name):
                        try:
                            comp_data = ole.openstream(comp_name).read()
                            comp_str  = comp_data.decode("latin-1", errors="ignore").lower()

                            # Vérifier formules / bitmaps en premier
                            for kw in ["equation.3", "equation.2", "equation.1",
                                       "equation native", "mathtype", "ms equation",
                                       "microsoft equation"]:
                                if kw in comp_str:
                                    ole.close(); os.unlink(temp_path)
                                    return ".equation"

                            for kw in ["paint.picture", "paintbrush", "msbitmap",
                                       "ms_clipart", "msdraw", "photoeditor",
                                       "coreldraw", "ms paint", "image document"]:
                                if kw in comp_str:
                                    ole.close(); os.unlink(temp_path)
                                    return ".bitmap_ole"

                            # Mapper ProgID → extension
                            for prog_key, ext in PROGID_TO_EXT.items():
                                if prog_key in comp_str:
                                    progid_detected = ext
                                    self.log(f"      🔍 ProgID détecté : {prog_key} → {ext}")
                                    break

                        except Exception:
                            pass
                        break

                # Stream Equation Native
                if any("equation native" in s.lower() for s in stream_list):
                    ole.close(); os.unlink(temp_path)
                    return ".equation"

                # ════════════════════════════════════════════════════════
                # ÉTAPE 1 : MSG
                # ════════════════════════════════════════════════════════
                msg_score = 0
                msg_indicators = {
                    "__properties_version1.0": 10,
                    "__nameid_version1.0":     10,
                    "__substg1.0_":             3,
                    "__recip_version1.0":       5,
                    "__attach_version1.0":      5,
                }
                for sn in stream_list:
                    for ind, score in msg_indicators.items():
                        if ind in sn:
                            msg_score += score

                if msg_score >= 10:
                    ole.close(); os.unlink(temp_path)
                    return ".msg"

                # ════════════════════════════════════════════════════════
                # ÉTAPE 2 : Types Office directs
                # ════════════════════════════════════════════════════════
                if ole.exists("Workbook") or ole.exists("Book"):
                    ole.close(); os.unlink(temp_path); return ".xls"
                if ole.exists("WordDocument"):
                    ole.close(); os.unlink(temp_path); return ".doc"
                if ole.exists("PowerPoint Document"):
                    ole.close(); os.unlink(temp_path); return ".ppt"

                # ════════════════════════════════════════════════════════
                # ÉTAPE 3 : Ole10Native → NOM DE FICHIER RÉEL
                #
                # C'est ici que la correction principale intervient.
                # Format Ole10Native (début) :
                #   [0x02 0x00]                  ← header 2 octets
                #   [len_filename : uint16 LE]   ← longueur nom (avec \0)
                #   [filename : bytes]            ← nom ASCII/Latin-1
                #   [len_filepath : uint16 LE]
                #   [filepath : bytes]
                #   ...
                # Certaines implémentations commencent directement par
                # le nom en ANSI sans header — on essaie les deux.
                # ════════════════════════════════════════════════════════
                original_name_from_native = None

                for native_name in ["\\x01Ole10Native", "Ole10Native"]:
                    if ole.exists(native_name):
                        try:
                            raw = ole.openstream(native_name).read()
                            self.log(f"      🔍 Ole10Native présent ({len(raw)} bytes)")

                            # ── Méthode A : header standard 4 octets (taille totale) ──
                            # Format : uint32 total_size, puis :
                            #   uint16 type (02 00 = linked/embedded)
                            #   pascal-string filename
                            if len(raw) > 6:
                                offset = 0
                                # Sauter le uint32 de taille totale si présent
                                if len(raw) > 4:
                                    # Tester si raw[0:4] ressemble à une taille
                                    import struct
                                    total_sz = struct.unpack_from("<I", raw, 0)[0]
                                    if 10 < total_sz <= len(raw) + 8:
                                        offset = 4  # on saute le uint32

                                # Lire le type (2 octets 02 00 = embedded)
                                if raw[offset:offset+2] in (b"\\x02\\x00", b"\\x01\\x00"):
                                    offset += 2

                                # Lire le nom (chaîne ANSI terminée par \\x00)
                                end = raw.find(b"\\x00", offset)
                                if 0 < end - offset < 260:
                                    candidate = raw[offset:end].decode(
                                        "latin-1", errors="ignore"
                                    ).strip()
                                    if candidate and "." in candidate:
                                        original_name_from_native = candidate
                                        self.log(
                                            f"      📝 Nom Ole10Native (méthode A): "
                                            f"{original_name_from_native}"
                                        )

                            # ── Méthode B : recherche pattern extension ──────────
                            if not original_name_from_native:
                                import re as _re
                                # Chercher dans les 2048 premiers octets
                                pattern = rb"([A-Za-z0-9_\\-\\. ]+\\.([A-Za-z0-9]{2,10}))\\x00"
                                for m in _re.finditer(pattern, raw[:2048]):
                                    candidate = m.group(1).decode("latin-1", errors="ignore").strip()
                                    ext_cand  = m.group(2).decode("latin-1", errors="ignore").lower()
                                    # Ignorer les extensions génériques Windows
                                    if ext_cand not in ("dll", "exe", "sys", "ini",
                                                        "tmp", "log", "lnk"):
                                        original_name_from_native = candidate
                                        self.log(
                                            f"      📝 Nom Ole10Native (méthode B): "
                                            f"{original_name_from_native}"
                                        )
                                        break

                        except Exception as native_err:
                            self.log(f"      ⚠️ Erreur lecture Ole10Native: {native_err}")
                        break  # un seul stream Ole10Native

                # ── Déduire l'extension depuis le nom trouvé ─────────────
                if original_name_from_native:
                    ext_from_name = Path(original_name_from_native).suffix.lower()
                    if ext_from_name and ext_from_name != ".ole":
                        self.log(
                            f"      ✅ Extension réelle depuis Ole10Native : "
                            f"{ext_from_name}"
                        )
                        ole.close(); os.unlink(temp_path)
                        return ext_from_name

                # ── Utiliser le ProgID si trouvé ─────────────────────────
                if progid_detected:
                    ole.close(); os.unlink(temp_path)
                    return progid_detected

                # ════════════════════════════════════════════════════════
                # ÉTAPE 4 : Package stream (fallback)
                # ════════════════════════════════════════════════════════
                if ole.exists("Package"):
                    try:
                        pkg_data = ole.openstream("Package").read()
                        search_limit = min(len(pkg_data), 4096)
                        for start in range(search_limit):
                            if pkg_data[start:start+4] == b"%PDF":
                                ole.close(); os.unlink(temp_path); return ".pdf"
                            if pkg_data[start:start+4] == b"PK\\x03\\x04":
                                detected = self.detect_office_xml_type(pkg_data[start:])
                                ole.close(); os.unlink(temp_path); return detected
                    except Exception as pkg_err:
                        self.log(f"      ⚠️ Erreur Package: {pkg_err}")

                ole.close(); os.unlink(temp_path)
                return ".ole"

            except Exception as ole_error:
                self.log(f"      ⚠️ Erreur ouverture OLE: {ole_error}")
                return ".ole"
            finally:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
                return ".ole"

        except Exception as e:
            self.log(f"      ⚠️ Erreur detect_ole_type: {e}")
            return ".ole"
    def extract_ole_content(self, ole_data, output_dir, base_name, position=None):
        """
        Extrait le contenu d'un fichier OLE.
        Ignore silencieusement : images BMP, bitmaps OLE, formules mathématiques.
        """
        original_filename = None

        # ========================================
        # VÉRIFICATION TAILLE 0 Ko
        # ========================================
        if len(ole_data) == 0:
            self.log(f"  ⏭️ Fichier encapsulé vide (0 Ko) ignoré à la position {position}")
            if position:
                self.current_file_extractions[position] = "Fichier vide (0 Ko) - ignoré"
            return None

        # ========================================
        # LISTES EXTENSIONS
        # ========================================
        SUPPORTED_EXTENSIONS = {
            '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls',
            '.pptx','.pptm', '.ppt', '.msg', '.txt', '.zip', '.7z','.htm', '.html' 
        }
        EXCLUDED_EXTENSIONS = {
            '.equation',   # formules mathématiques Word
            '.bitmap_ole', # bitmaps OLE (Paint, MSBitmap...)
            '.ole', '.bin', '.dat', '.tmp',
            # CATIA
            '.catpart', '.catproduct', '.catdrawing',
            '.catanalysis', '.catprocess', '.catmaterial',
            '.cgr', '.model', '.exp', '.dlv', '.dlv3', '.catalog',
            '.pmat', '.rpx', '.3dxml',
        }
        # Extensions image — on ignore silencieusement sans rapport d'erreur
        IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp',
                            '.tiff', '.emf', '.wmf', '.svg', '.ico'}

        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.ole') as temp_ole:
                temp_ole.write(ole_data)
                temp_ole_path = temp_ole.name

            try:
                ole = olefile.OleFileIO(temp_ole_path)
                streams = ole.listdir()

                self.log(f"    🔍 Analyse OLE...")
                stream_list = ['/'.join(entry) for entry in streams]

                # ========================================
                # ÉTAPE 0 : BITMAP OLE / FORMULE MATHÉMATIQUE
                # Lecture du stream CompObj → ProgID réel
                # ========================================
                for comp_name in ['\x01CompObj', 'CompObj']:
                    if ole.exists(comp_name):
                        try:
                            comp_data = ole.openstream(comp_name).read()
                            comp_str = comp_data.decode('latin-1', errors='ignore').lower()

                            equation_keywords = [
                                'equation.3', 'equation.2', 'equation.1',
                                'equation native', 'mathtype', 'ms equation',
                                'microsoft equation',
                            ]
                            for kw in equation_keywords:
                                if kw in comp_str:
                                    self.log(f"    ⏭️ Formule mathématique ({kw}) → ignorée")
                                    ole.close()
                                    os.unlink(temp_ole_path)
                                    if position:
                                        self.current_file_extractions[position] = "Formule mathématique - ignorée"
                                    return None

                            bitmap_keywords = [
                                'paint.picture', 'paintbrush', 'msbitmap',
                                'ms_clipart', 'msdraw', 'photoeditor',
                                'coreldraw', 'ms paint', 'image document',
                            ]
                            for kw in bitmap_keywords:
                                if kw in comp_str:
                                    self.log(f"    ⏭️ Bitmap OLE ({kw}) → ignoré")
                                    ole.close()
                                    os.unlink(temp_ole_path)
                                    if position:
                                        self.current_file_extractions[position] = "Bitmap OLE - ignoré"
                                    return None

                        except Exception:
                            pass
                        break

                # Stream Equation Native sans CompObj lisible
                if any('equation native' in s.lower() for s in stream_list):
                    self.log(f"    ⏭️ Formule mathématique (Equation Native) → ignorée")
                    ole.close()
                    os.unlink(temp_ole_path)
                    if position:
                        self.current_file_extractions[position] = "Formule mathématique - ignorée"
                    return None

                # ========================================
                # DÉTECTION MSG RENFORCÉE
                # ========================================
                is_msg = False
                msg_score = 0
                msg_indicators = {
                    '__properties_version1.0': 10,
                    '__nameid_version1.0': 10,
                    '__substg1.0_': 3,
                    '__recip_version1.0': 5,
                    '__attach_version1.0': 5,
                    'Root Entry': 2,
                }

                for stream_name in stream_list:
                    for indicator, score in msg_indicators.items():
                        if indicator in stream_name:
                            msg_score += score
                            self.log(f"      🔍 Indicateur MSG trouvé: {indicator} (+{score})")

                if msg_score >= 10:
                    is_msg = True
                    self.log(f"      ✅ MSG CONFIRMÉ (score: {msg_score})")

                # ========================================
                # EXTRACTION MSG
                # ========================================
                if is_msg:
                    self.log(f"    📧 Fichier MSG Outlook détecté - EXTRACTION")
                    self.extracted_count += 1

                    output_name = self._generate_output_name(base_name, None, '.msg')
                    output_path = Path(output_dir) / output_name

                    with open(output_path, 'wb') as f:
                        f.write(ole_data)

                    if position:
                        self.current_file_extractions[position] = output_name

                    self.log(f"  ✅ MSG extrait: {output_name}")

                    # ── Extraire les PJ du MSG AVANT la conversion PDF ──────────
                    self.log(f"  🔄 Extraction des pièces jointes du MSG encapsulé...")
                    msg_attachments = self._extract_msg_attachments_only(output_path, output_dir)
                    # ── Marquer ce MSG comme déjà traité (skip récursion) ──
                    self._msg_already_processed.add(str(output_path.resolve()))
                    if msg_attachments:
                        self.log(f"  ✅ {len(msg_attachments)} PJ extraite(s) : {msg_attachments}")
                    else:
                        self.log(f"  ℹ️ Aucune pièce jointe dans ce MSG")

                    # ── Conversion MSG → PDF avec références PJ ─────────────────
                    self.log(f"  🔄 Tentative conversion MSG → PDF...")
                    try:
                        pdf_path = self.convert_msg_to_pdf(output_path, output_dir, msg_attachments)
                        if pdf_path and pdf_path.exists():
                            self.log(f"  ✅ PDF créé: {pdf_path.name}")
                        else:
                            self.log(f"  ⚠️ PDF non créé — suppression MSG quand même")
                        # Toujours supprimer le MSG du dossier de sortie
                        if output_path.exists():
                            try:
                                output_path.unlink()
                                self.log(f"  🗑️ MSG supprimé: {output_path.name}")
                            except Exception as del_err:
                                self.log(f"  ⚠️ Impossible de supprimer le MSG: {del_err}")
                    except Exception as pdf_error:
                        self.log(f"  ⚠️ Erreur conversion PDF: {pdf_error}")
                        # Supprimer quand même le MSG même si conversion échoue
                        if output_path.exists():
                            try:
                                output_path.unlink()
                                self.log(f"  🗑️ MSG supprimé (après erreur PDF): {output_path.name}")
                            except Exception as del_err:
                                self.log(f"  ⚠️ Impossible de supprimer le MSG: {del_err}")

                    ole.close()
                    os.unlink(temp_ole_path)
                    return output_name

                # ========================================
                # OLE10NATIVE
                # ========================================
                if ole.exists('\x01Ole10Native') or ole.exists('Ole10Native'):
                    try:
                        self.log(f"    📦 Ole10Native détecté")
                        result = self.extract_ole10native(ole)

                        if result:
                            if isinstance(result, tuple):
                                file_data, original_filename = result
                            else:
                                file_data = result
                                original_filename = None

                            # ── ÉTAPE 0-bis : BMP dans Ole10Native ──────────────
                            # Doit être AVANT la détection d'extension et le fallback texte
                            if file_data[:2] == b'BM' and len(file_data) > 14:
                                try:
                                    import struct
                                    bmp_size = struct.unpack_from('<I', file_data, 2)[0]
                                    if bmp_size > 100:
                                        self.log(f"    ⏭️ Image BMP ignorée (magic bytes BM, taille={bmp_size})")
                                        ole.close()
                                        os.unlink(temp_ole_path)
                                        if position:
                                            self.current_file_extractions[position] = "Image BMP - ignorée"
                                        return None
                                except Exception:
                                    pass

                            # ── Vérifier ratio binaire AVANT fallback .txt ────────
                            # Évite de traiter un fichier binaire comme texte
                            def _is_binary(data, sample=512):
                                chunk = data[:sample]
                                if len(chunk) == 0:
                                    return False
                                non_print = sum(
                                    1 for b in chunk
                                    if b < 9 or (14 <= b <= 31)
                                )
                                return (non_print / len(chunk)) > 0.05

                            # ── Priorité au vrai nom de fichier OLE ───────────────
                            GENERIQUE = {'', '.bin', '.dat', '.tmp', '.ole'}

                            if original_filename and Path(original_filename).suffix.lower() not in GENERIQUE:
                                file_ext = Path(original_filename).suffix.lower()
                                self.log(f"    📄 Extension réelle (nom fichier OLE): {file_ext} ← {original_filename}")
                            else:
                                # Pas de nom fiable → détection par signature binaire
                                file_ext = self.detect_file_type(file_data)

                                # Si detect_file_type retourne .txt mais que le contenu
                                # est trop binaire → c'est probablement une image/bitmap
                                if file_ext == '.txt' and _is_binary(file_data):
                                    self.log(f"    ⏭️ Données binaires non reconnues (ratio > 5%) → ignorées")
                                    ole.close()
                                    os.unlink(temp_ole_path)
                                    if position:
                                        self.current_file_extractions[position] = "Données binaires non reconnues - ignorées"
                                    return None

                                self.log(f"    📄 Type détecté (signature binaire): {file_ext}")

                            # ── OVERRIDE magic bytes 7z ───────────────────────────
                            if file_data[:6] == b'\x37\x7a\xbc\xaf\x27\x1c':
                                self.log(f"    🔄 OVERRIDE: magic bytes 7z → .7z (était: {file_ext})")
                                file_ext = '.7z'

                            # ── DÉTECTION EXCEL ROBUSTE ───────────────────────────
                            if file_data[:8] == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
                                try:
                                    import io as _io
                                    test_ole = olefile.OleFileIO(_io.BytesIO(file_data))

                                    # Vérifier bitmap/équation dans l'OLE imbriqué
                                    for cn in ['\x01CompObj', 'CompObj']:
                                        if test_ole.exists(cn):
                                            try:
                                                cd = test_ole.openstream(cn).read()
                                                cs = cd.decode('latin-1', errors='ignore').lower()
                                                for kw in ['equation.3', 'equation.2', 'equation.1',
                                                           'mathtype', 'paint.picture', 'paintbrush',
                                                           'msbitmap', 'msdraw']:
                                                    if kw in cs:
                                                        self.log(f"    ⏭️ OLE imbriqué bitmap/équation ({kw}) → ignoré")
                                                        test_ole.close()
                                                        ole.close()
                                                        os.unlink(temp_ole_path)
                                                        if position:
                                                            self.current_file_extractions[position] = "Bitmap/Équation OLE - ignoré"
                                                        return None
                                            except Exception:
                                                pass
                                            break

                                    has_workbook = any(
                                        'workbook' in '/'.join(s).lower() or 'book' in '/'.join(s).lower()
                                        for s in test_ole.listdir()
                                    )
                                    test_ole.close()

                                    if has_workbook:
                                        self.log(f"      ✅ Excel binaire (.xls) détecté")
                                        file_ext = '.xls'

                                except Exception as test_error:
                                    self.log(f"      ⚠ Test Excel: {test_error}")

                            elif file_data[:4] == b'PK\x03\x04':
                                try:
                                    import io as _io, zipfile
                                    with zipfile.ZipFile(_io.BytesIO(file_data)) as zf:
                                        if any('xl/' in name for name in zf.namelist()):
                                            self.log(f"      ✅ Excel moderne (.xlsx) détecté")
                                            file_ext = '.xlsx'
                                except Exception as zip_error:
                                    self.log(f"      ⚠ Test ZIP: {zip_error}")

                            # ── Vérifier MSG ──────────────────────────────────────
                            if file_data[:4] == b'\xd0\xcf\x11\xe0':
                                try:
                                    with tempfile.NamedTemporaryFile(delete=False, suffix='.tmp') as tmp:
                                        tmp.write(file_data)
                                        tmp_path = tmp.name

                                    test_ole2 = olefile.OleFileIO(tmp_path)
                                    test_streams2 = ['/'.join(e) for e in test_ole2.listdir()]

                                    inner_msg_score = 0
                                    for sn in test_streams2:
                                        if '__properties_version1.0' in sn: inner_msg_score += 10
                                        if '__nameid_version1.0' in sn:    inner_msg_score += 10
                                        if '__substg1.0_' in sn:           inner_msg_score += 3
                                        if '__recip_version1.0' in sn:     inner_msg_score += 5

                                    test_ole2.close()
                                    os.unlink(tmp_path)

                                    if inner_msg_score >= 10:
                                        self.log(f"      ✅ MSG détecté dans Ole10Native (score: {inner_msg_score})")
                                        file_ext = '.msg'

                                except Exception as test_err:
                                    self.log(f"      ⚠️ Test MSG: {test_err}")

                            # ── Vérification extension supportée ─────────────────
                            if file_ext in IMAGE_EXTENSIONS:
                                self.log(f"    ⏭️ Image ignorée ({file_ext})")
                                ole.close()
                                os.unlink(temp_ole_path)
                                if position:
                                    self.current_file_extractions[position] = f"Image {file_ext} - ignorée"
                                return None

                            if file_ext in EXCLUDED_EXTENSIONS:
                                label = ("Formule mathématique" if file_ext == '.equation'
                                         else "Bitmap OLE" if file_ext == '.bitmap_ole'
                                         else f"Type non supporté: {file_ext}")
                                self.log(f"    ⏭️ {label} → ignoré")
                                ole.close()
                                os.unlink(temp_ole_path)
                                if position:
                                    self.current_file_extractions[position] = f"{label} - ignoré"
                                return None

                            if file_ext not in SUPPORTED_EXTENSIONS:
                                self.log(f"  ⚠️ EXTENSION NON SUPPORTÉE: {file_ext}")
                                self.log(f"  📌 Le fichier reste dans le document parent")

                                if position:
                                    self.current_file_extractions[position] = f"Type non supporté: {file_ext}"

                                display_name = original_filename if original_filename else f"Fichier {file_ext}"
                                self.add_to_report('extracted', {
                                    'source_file': self.current_source_file or 'Inconnu',
                                    'extracted_file': f'[Non supporté: {display_name}]',
                                    'position': position if position else 'N/A',
                                    'type': file_ext,
                                    'status': '⚠️ Type non supporté - Conservé dans document'
                                })

                                ole.close()
                                os.unlink(temp_ole_path)
                                return None

                            # ── Déduplication AVANT assignation de numéro ────────
                            import hashlib
                            _content_hash = hashlib.md5(file_data).hexdigest()
                            if _content_hash in self._extracted_content_hashes:
                                _existing = self._content_hash_to_name.get(_content_hash)
                                if _existing:
                                    self.log(f"    ⏭️ Contenu déjà extrait → réutilise {_existing}")
                                    ole.close()
                                    os.unlink(temp_ole_path)
                                    return _existing
                                self.log(f"    ⏭️ Contenu identique déjà extrait (MD5={_content_hash[:8]}...), ignoré")
                                ole.close()
                                os.unlink(temp_ole_path)
                                return None

                            # ── Extraction effective ──────────────────────────────
                            self._extracted_content_hashes.add(_content_hash)
                            self.extracted_count += 1

                            output_name = self._generate_output_name(base_name, original_filename, file_ext)
                            output_path = Path(output_dir) / output_name
                            self._content_hash_to_name[_content_hash] = output_name

                            with open(output_path, 'wb') as f:
                                f.write(file_data)

                            if position:
                                self.current_file_extractions[position] = output_name

                            self.log(f"  ✓ Fichier extrait: {output_name}")

                            if file_ext == '.zip':
                                self.log(f"  📌 ZIP sera remplacé dans document modifié")

                            if file_ext == '.msg':
                                # Marquer MSG pour eviter retraitement en recursion
                                self._msg_already_processed.add(str(output_path.resolve()))
                                self._processed_absolute_paths.add(str(output_path.resolve()))
                                # ── Extraire les PJ du MSG AVANT la conversion PDF ──────────
                                self.log(f"  🔄 Extraction des pièces jointes du MSG encapsulé...")
                                msg_attachments = self._extract_msg_attachments_only(output_path, output_dir)
                                if msg_attachments:
                                    self.log(f"  ✅ {len(msg_attachments)} PJ extraite(s) : {msg_attachments}")
                                else:
                                    self.log(f"  ℹ️ Aucune pièce jointe dans ce MSG")

                                # ── Conversion MSG → PDF avec références PJ ─────────────────
                                self.log(f"  🔄 Tentative conversion MSG → PDF...")
                                try:
                                    pdf_path = self.convert_msg_to_pdf(output_path, output_dir, msg_attachments)
                                    if pdf_path and pdf_path.exists():
                                        self.log(f"  ✅ PDF créé: {pdf_path.name}")
                                    else:
                                        self.log(f"  ⚠️ PDF non créé — suppression MSG quand même")
                                    # Toujours supprimer le MSG du dossier de sortie
                                    if output_path.exists():
                                        try:
                                            output_path.unlink()
                                            self.log(f"  🗑️ MSG supprimé: {output_path.name}")
                                        except Exception as del_err:
                                            self.log(f"  ⚠️ Impossible de supprimer le MSG: {del_err}")
                                except Exception as pdf_error:
                                    self.log(f"  ⚠️ Erreur conversion PDF: {pdf_error}")
                                    # Supprimer quand même le MSG même si conversion échoue
                                    if output_path.exists():
                                        try:
                                            output_path.unlink()
                                            self.log(f"  🗑️ MSG supprimé (après erreur PDF): {output_path.name}")
                                        except Exception as del_err:
                                            self.log(f"  ⚠️ Impossible de supprimer le MSG: {del_err}")

                            ole.close()
                            os.unlink(temp_ole_path)
                            return output_name

                    except Exception as ole10_error:
                        self.log(f"    ⚠ Erreur Ole10Native: {ole10_error}")

                # ========================================
                # PACKAGE
                # ========================================
                if ole.exists('Package'):
                    try:
                        self.log(f"    📦 Package détecté")
                        package_stream = ole.openstream('Package')
                        package_data = package_stream.read()
                        actual_file = self.extract_file_from_package(package_data)

                        if actual_file:
                            file_ext = self.detect_file_type(actual_file)
                            self.log(f"    📄 Type Package: {file_ext}")

                            if file_ext in IMAGE_EXTENSIONS:
                                self.log(f"    ⏭️ Image ignorée ({file_ext})")
                                ole.close()
                                os.unlink(temp_ole_path)
                                return None

                            if file_ext in EXCLUDED_EXTENSIONS or file_ext not in SUPPORTED_EXTENSIONS:
                                self.log(f"  ⚠️ EXTENSION NON SUPPORTÉE: {file_ext}")
                                self.log(f"  📌 Le fichier reste dans le document parent")

                                if position:
                                    self.current_file_extractions[position] = f"Type non supporté: {file_ext}"

                                display_name = original_filename if original_filename else f"Fichier {file_ext}"
                                self.add_to_report('extracted', {
                                    'source_file': self.current_source_file or 'Inconnu',
                                    'extracted_file': f'[Non supporté: {display_name}]',
                                    'position': position if position else 'N/A',
                                    'type': file_ext,
                                    'status': '⚠️ Type non supporté - Conservé dans document'
                                })

                                ole.close()
                                os.unlink(temp_ole_path)
                                return None

                            self.extracted_count += 1

                            output_name = self._generate_output_name(base_name, None, file_ext)
                            output_path = Path(output_dir) / output_name

                            with open(output_path, 'wb') as f:
                                f.write(actual_file)

                            if position:
                                self.current_file_extractions[position] = output_name

                            self.log(f"  ✓ Fichier extrait: {output_name}")

                            if file_ext == '.zip':
                                self.log(f"  📌 ZIP sera remplacé dans document modifié")

                            ole.close()
                            os.unlink(temp_ole_path)
                            return output_name

                    except Exception as pkg_error:
                        self.log(f"    ⚠ Erreur Package: {pkg_error}")

# ========================================
                # EXCEL DIRECT — reconstruction d'un .xls lisible
                # ========================================
                if ole.exists('Workbook') or ole.exists('Book'):
                    try:
                        self.log(f"    📊 Excel binaire - SUPPORTÉ")
                        self.extracted_count += 1

                        output_name = self._generate_output_name(base_name, None, '.xls')
                        output_path = Path(output_dir) / output_name

                        # ── Extraire uniquement le stream Workbook ──────────
                        wb_stream_name = 'Workbook' if ole.exists('Workbook') else 'Book'
                        wb_bytes = ole.openstream(wb_stream_name).read()
                        self.log(f"      📖 Stream '{wb_stream_name}' : {len(wb_bytes):,} bytes")

                        # ── Reconstruire un OLE minimal autonome ────────────
                        # Un .xls valide = OLE avec UNIQUEMENT le stream Workbook
                        # On utilise le module 'compoundfiles' ou on passe par
                        # une écriture directe du format BIFF8 via xlwt
                        rebuilt = False

                        # MÉTHODE 1 : xlrd + xlutils.copy (conserve mise en forme)
                        try:
                            import xlrd
                            import xlutils.copy as xlutils_copy

                            # Lire depuis l'OLE COMPLET (pas juste le stream)
                            # xlrd sait lire un OLE embarqué directement
                            wb_read = xlrd.open_workbook(
                                file_contents=ole_data,  # ← OLE COMPLET
                                formatting_info=True
                            )
                            wb_write = xlutils_copy.copy(wb_read)
                            wb_write.save(str(output_path))

                            self.log(f"      ✅ .xls reconstruit via xlrd+xlutils")
                            rebuilt = True

                        except ImportError as imp_err:
                            self.log(f"      ⚠️ xlrd/xlutils absent: {imp_err}")
                        except Exception as xlrw_err:
                            self.log(f"      ⚠️ xlrd+xlutils échoué: {xlrw_err}")
                        # MÉTHODE 2 : Excel COM sur OLE COMPLET
                        if not rebuilt:
                            try:
                                import win32com.client
                                import pythoncom

                                # Écrire l'OLE COMPLET dans tmp (pas le stream seul)
                                with tempfile.NamedTemporaryFile(
                                    delete=False, suffix='.xls'
                                ) as tmp_f:
                                    tmp_f.write(ole_data)  # ← OLE COMPLET
                                    tmp_path = tmp_f.name

                                # Débloquer le fichier tmp avant ouverture COM
                                subprocess.run(
                                    ['powershell', '-Command',
                                     f'Unblock-File -Path "{tmp_path}"'],
                                    capture_output=True, timeout=5
                                )

                                pythoncom.CoInitialize()
                                try:
                                    xl = win32com.client.DispatchEx("Excel.Application")
                                    xl.Visible = False
                                    xl.DisplayAlerts = False
                                    xl.AskToUpdateLinks = False
                                    xl.AutomationSecurity = 3  # ForceDisable macros

                                    wb_com = xl.Workbooks.Open(
                                        tmp_path,
                                        UpdateLinks=0,
                                        ReadOnly=True,
                                        IgnoreReadOnlyRecommended=True,
                                        CorruptLoad=1
                                    )
                                    xl.DisplayAlerts = False

                                    wb_com.SaveAs(
                                        str(output_path.absolute()),
                                        FileFormat=56  # xlExcel8 = .xls
                                    )
                                    wb_com.Close(False)
                                    xl.Quit()
                                    time.sleep(1)
                                    self._kill_excel()

                                    if os.path.exists(tmp_path):
                                        os.unlink(tmp_path)

                                    self.log(f"      ✅ .xls reconstruit via Excel COM")
                                    rebuilt = True

                                finally:
                                    pythoncom.CoUninitialize()

                            except ImportError:
                                self.log(f"      ⚠️ pywin32 absent")
                            except Exception as com_err:
                                self.log(f"      ⚠️ Excel COM échoué: {com_err}")
                                self._kill_excel()
                                if 'tmp_path' in locals() and os.path.exists(tmp_path):
                                    os.unlink(tmp_path)
                        # MÉTHODE 3 : fallback OLE brut + Unblock
                        if not rebuilt:
                            with open(output_path, 'wb') as f:
                                f.write(ole_data)
                            self.log(f"      ⚠️ OLE brut écrit (fallback)")

                        # Débloquer dans tous les cas
                        try:
                            subprocess.run(
                                ['powershell', '-Command',
                                 f'Unblock-File -Path "{str(output_path)}"'],
                                capture_output=True, timeout=5
                            )
                            self.log(f"      ✅ Unblock-File appliqué")
                        except Exception:
                            pass

                        if position:
                            self.current_file_extractions[position] = output_name

                        self.log(f"  ✓ Excel extrait: {output_name}")

                        ole.close()
                        os.unlink(temp_ole_path)
                        return output_name

                    except Exception as excel_err:
                        self.log(f"    ⚠️ Erreur extraction Excel: {excel_err}")

                # ========================================
                # WORD DIRECT
                # ========================================
                if ole.exists('WordDocument'):
                    try:
                        self.log(f"    📝 Word binaire - SUPPORTÉ")
                        self.extracted_count += 1

                        output_name = self._generate_output_name(base_name, None, '.doc')
                        output_path = Path(output_dir) / output_name

                        with open(output_path, 'wb') as f:
                            f.write(ole_data)

                        if position:
                            self.current_file_extractions[position] = output_name

                        self.log(f"  ✓ Word extrait: {output_name}")

                        ole.close()
                        os.unlink(temp_ole_path)
                        return output_name
                    except Exception:
                        pass

                # ========================================
                # POWERPOINT DIRECT
                # ========================================
                if ole.exists('PowerPoint Document') or ole.exists('Current User'):
                    try:
                        self.log(f"    📽️ PowerPoint binaire - SUPPORTÉ")
                        self.extracted_count += 1

                        output_name = self._generate_output_name(base_name, None, '.ppt')
                        output_path = Path(output_dir) / output_name

                        with open(output_path, 'wb') as f:
                            f.write(ole_data)

                        if position:
                            self.current_file_extractions[position] = output_name

                        self.log(f"  ✓ PowerPoint extrait: {output_name}")

                        ole.close()
                        os.unlink(temp_ole_path)
                        return output_name
                    except Exception:
                        pass

                # ========================================
                # CONTENTS / OBJDATA
                # ========================================
                for entry in streams:
                    stream_name = '/'.join(entry)
                    if 'contents' in stream_name.lower() or 'objdata' in stream_name.lower():
                        try:
                            stream_data = ole.openstream(entry).read()
                            if len(stream_data) > 100:
                                actual_file = self.extract_actual_content(stream_data)
                                file_ext = self.detect_file_type(actual_file)

                                if file_ext in IMAGE_EXTENSIONS:
                                    self.log(f"    ⏭️ Image ignorée ({file_ext})")
                                    continue

                                if file_ext in EXCLUDED_EXTENSIONS or file_ext not in SUPPORTED_EXTENSIONS:
                                    self.log(f"  ⚠️ EXTENSION NON SUPPORTÉE: {file_ext}")
                                    self.log(f"  📌 Le fichier reste dans le document parent")

                                    if position:
                                        self.current_file_extractions[position] = f"Type non supporté: {file_ext}"

                                    display_name = original_filename if original_filename else f"Fichier {file_ext}"
                                    self.add_to_report('extracted', {
                                        'source_file': self.current_source_file or 'Inconnu',
                                        'extracted_file': f'[Non supporté: {display_name}]',
                                        'position': position if position else 'N/A',
                                        'type': file_ext,
                                        'status': '⚠️ Type non supporté - Conservé dans document'
                                    })

                                    ole.close()
                                    os.unlink(temp_ole_path)
                                    return None

                                self.extracted_count += 1

                                output_name = self._generate_output_name(base_name, None, file_ext)
                                output_path = Path(output_dir) / output_name

                                with open(output_path, 'wb') as f:
                                    f.write(actual_file)

                                if position:
                                    self.current_file_extractions[position] = output_name

                                self.log(f"  ✓ Fichier extrait: {output_name}")

                                if file_ext == '.zip':
                                    self.log(f"  📌 ZIP sera remplacé dans document modifié")

                                ole.close()
                                os.unlink(temp_ole_path)
                                return output_name
                        except Exception:
                            continue

                # ========================================
                # DERNIER RECOURS
                # ========================================
                if len(ole_data) > 1000:
                    self.log(f"    ⚠ Type OLE non identifié")

                    file_ext = self.detect_file_type(ole_data)

                    if file_ext in IMAGE_EXTENSIONS:
                        self.log(f"    ⏭️ Image ignorée ({file_ext})")
                        ole.close()
                        os.unlink(temp_ole_path)
                        return None

                    if file_ext in EXCLUDED_EXTENSIONS or file_ext not in SUPPORTED_EXTENSIONS:
                        self.log(f"  ⚠️ EXTENSION NON SUPPORTÉE: {file_ext}")
                        self.log(f"  📌 Le fichier reste dans le document parent")

                        if position:
                            self.current_file_extractions[position] = f"Type non supporté: {file_ext}"

                        display_name = original_filename if original_filename else f"Fichier {file_ext}"
                        self.add_to_report('extracted', {
                            'source_file': self.current_source_file or 'Inconnu',
                            'extracted_file': f'[Non supporté: {display_name}]',
                            'position': position if position else 'N/A',
                            'type': file_ext,
                            'status': '⚠️ Type non supporté - Conservé dans document'
                        })

                        ole.close()
                        os.unlink(temp_ole_path)
                        return None

                    self.extracted_count += 1

                    output_name = self._generate_output_name(base_name, None, file_ext)
                    output_path = Path(output_dir) / output_name

                    with open(output_path, 'wb') as f:
                        f.write(ole_data)

                    if position:
                        self.current_file_extractions[position] = output_name

                    self.log(f"  ✓ Fichier sauvegardé: {output_name}")

                    ole.close()
                    os.unlink(temp_ole_path)
                    return output_name

                ole.close()
                os.unlink(temp_ole_path)

            except Exception as ole_error:
                self.log(f"    ⚠ Erreur OLE: {ole_error}")

                file_ext = self.detect_file_type(ole_data)

                if file_ext in IMAGE_EXTENSIONS:
                    self.log(f"    ⏭️ Image ignorée ({file_ext})")
                    if 'temp_ole_path' in locals():
                        try: os.unlink(temp_ole_path)
                        except Exception: pass
                    return None

                if file_ext in EXCLUDED_EXTENSIONS or file_ext not in SUPPORTED_EXTENSIONS:
                    self.log(f"  ⚠️ EXTENSION NON SUPPORTÉE: {file_ext}")
                    self.log(f"  📌 Le fichier reste dans le document parent")

                    if position:
                        self.current_file_extractions[position] = "Ce type fichier n'est pas supporté par l'extraction de l'outil"

                    if 'temp_ole_path' in locals():
                        try: os.unlink(temp_ole_path)
                        except Exception: pass

                    return None

                self.extracted_count += 1

                output_name = self._generate_output_name(base_name, None, file_ext)
                output_path = Path(output_dir) / output_name

                with open(output_path, 'wb') as f:
                    f.write(ole_data)

                if position:
                    self.current_file_extractions[position] = output_name

                self.log(f"  ✓ Fichier sauvegardé: {output_name}")

                if 'temp_ole_path' in locals():
                    try: os.unlink(temp_ole_path)
                    except Exception: pass

                return output_name

        except Exception as e:
            self.log(f"  ❌ Erreur extraction: {str(e)}")
            if 'temp_ole_path' in locals():
                try: os.unlink(temp_ole_path)
                except Exception: pass
    def extract_ole_from_docx(self, docx_path, output_dir):
        """
        Extrait fichiers Word avec POSITION EXACTE conservée
        - Remplace fichiers extraits par texte rouge À LA MÊME POSITION
        - Conserve fichiers non supportés avec message d'avertissement
        CORRECTION : Vérifications filename + mapping par nom de fichier
        """
        # Nettoyage préalable
        bin_cleaned = self.clean_bin_files_from_embeddings(docx_path)
        if bin_cleaned > 0:
            self.log(f"  ✅ Pré-traitement: {bin_cleaned} fichier(s) .bin vide(s) supprimé(s)")
        
        self.reset_extraction_tracking(Path(docx_path).name, output_dir)
        Path(output_dir).mkdir(parents=True, exist_ok=True)

        extracted_files = []

        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)

                # Décompresser le DOCX
                with zipfile.ZipFile(docx_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_path)
                
                embeddings_dir = temp_path / "word" / "embeddings"
                
                if embeddings_dir.exists():
                    # ========================================
                    # ÉTAPE 1 : MAPPER POSITIONS EXACTES
                    # ========================================
                    self.log(f"\n  🔍 MAPPING EXACT DES POSITIONS DANS LE DOCUMENT...")
                    
                    position_mapping = self.get_docx_ole_positions_exact(docx_path, temp_path)
                    
                    if not position_mapping:
                        self.log(f"    ⚠️ Aucun mapping trouvé")
                        # Copier l'original
                        original_copy = Path(output_dir) / Path(docx_path).name
                        shutil.copy2(docx_path, original_copy)
                        self.log(f"  ✅ Copie créée: {original_copy.name}")
                        return []
                    
                    # ========================================
                    # ÉTAPE 2 : TRAITER CHAQUE FICHIER
                    # ========================================
                    self.log(f"\n  📦 EXTRACTION ET CLASSIFICATION...")
                    
                    for filename in sorted(position_mapping.keys(), key=lambda x: position_mapping[x]['position']):
                        location = position_mapping[filename]
                        position = location['position']
                        page_estimate = location.get('page', '?')
                        context = location.get('context', '')
                        
                        # ✅ VÉRIFICATION CRITIQUE : filename doit être valide
                        if not filename or filename.strip() == '':
                            self.log(f"    ❌ ERREUR CRITIQUE : filename vide pour position {position}, IGNORÉ")
                            continue
                        
                        embed_path = embeddings_dir / filename
                        
                        if not embed_path.exists():
                            self.log(f"    ⚠️ Fichier {filename} introuvable dans embeddings/")
                            continue
                        
                        self.log(f"\n  🔍 Position {position} (Page ~{page_estimate}, {context})")
                        self.log(f"    📂 Fichier: {filename}")
                        
                        # Lire le fichier
                        with open(embed_path, 'rb') as f:
                            file_data = f.read()
                        
                        self.log(f"    📦 Taille: {len(file_data)} bytes")
                        
                        # ========================================
                        # CAS 1 : FICHIER VIDE
                        # ========================================
                        if len(file_data) == 0:
                            self.log(f"    ⚠️ Fichier VIDE détecté")
                            
                            # ✅ VÉRIFICATION : filename valide
                            if not filename or filename.strip() == '':
                                self.log(f"    ❌ ERREUR : filename vide, ignoré")
                                continue
                            
                            extracted_files.append({
                                'position': position,
                                'filename': filename,
                                'action': 'keep_empty',
                                'page': page_estimate,
                                'context': context,
                                'location': location
                            })
                            
                            self.add_to_report('extracted', {
                                'source_file': Path(docx_path).name,
                                'extracted_file': '[Fichier vide - 0 Ko]',
                                'position': f'Page {page_estimate}',
                                'type': 'Vide',
                                'status': 'Ignoré - Taille 0 Ko'
                            })
                            continue
                        
                        # ========================================
                        # CAS 2 : IMAGE (ignorer)
                        # ========================================
                        if not self.is_file_not_image(filename, file_data):
                            self.log(f"    ⏭️ Image ignorée")
                            continue
                        
                        # ========================================
                        # CAS 3 : FICHIER OLE
                        # ========================================
                        if file_data[:4] == b'\xd0\xcf\x11\xe0':
                            self.log(f"    📦 Fichier OLE détecté")
                            
                            # Détecter type
                            file_ext = self.detect_file_type(file_data)
                            self.log(f"    🔍 Type détecté: {file_ext}")
                            
                            # Vérifier si supporté
                            SUPPORTED_EXTENSIONS = {
                                '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls', 
                                '.pptx', '.ppt','.pptm', '.msg', '.txt', '.zip', '.7z','.htm', '.html' 
                            }
                            
                            if file_ext not in SUPPORTED_EXTENSIONS and file_ext != '.ole':
                                # NON SUPPORTÉ → GARDER
                                self.log(f"    ⚠️ Type NON SUPPORTÉ: {file_ext}")
                                
                                # ✅ VÉRIFICATION : filename valide
                                if not filename or filename.strip() == '':
                                    self.log(f"    ❌ ERREUR : filename vide, ignoré")
                                    continue
                                
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'keep_unsupported',
                                    'page': page_estimate,
                                    'context': context,
                                    'location': location,
                                    'file_type': file_ext
                                })
                                
                                self.add_to_report('extracted', {
                                    'source_file': Path(docx_path).name,
                                    'extracted_file': f'[Non supporté: {file_ext}]',
                                    'position': f'Page {page_estimate}',
                                    'type': file_ext,
                                    'status': 'Conservé dans le document'
                                })
                                continue
                            
                            # SUPPORTÉ → EXTRAIRE
                            self.log(f"    ✅ Type SUPPORTÉ, extraction...")
                            
                            extracted_name = self.extract_ole_content(
                                file_data, output_dir, Path(docx_path).stem, position
                            )
                            
                            if extracted_name:
                                self.log(f"    ✅ Extrait: {extracted_name}")
                                
                                # ✅ VÉRIFICATION : filename valide
                                if not filename or filename.strip() == '':
                                    self.log(f"    ❌ ERREUR : filename vide, ignoré")
                                    continue
                                
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'replace',
                                    'page': page_estimate,
                                    'context': context,
                                    'location': location,
                                    'extracted_name': extracted_name
                                })
                                
                                self.add_to_report('extracted', {
                                    'source_file': Path(docx_path).name,
                                    'extracted_file': extracted_name,
                                    'position': f'Page {page_estimate}',
                                    'type': Path(extracted_name).suffix,
                                    'status': 'Succès'
                                })
                            else:
                                # Échec extraction mais type supporté
                                self.log(f"    ⚠️ Extraction échouée")
                                
                                # ✅ VÉRIFICATION : filename valide
                                if not filename or filename.strip() == '':
                                    self.log(f"    ❌ ERREUR : filename vide, ignoré")
                                    continue
                                
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'keep_unsupported',
                                    'page': page_estimate,
                                    'context': context,
                                    'location': location,
                                    'file_type': file_ext
                                })
                        
                        # ========================================
                        # CAS 4 : AUTRE FICHIER (non-OLE)
                        # ========================================
                        else:
                            self.log(f"    📄 Fichier direct (non-OLE) détecté")
                            
                            file_ext = self.detect_file_type(file_data)
                            self.log(f"    🔍 Type détecté: {file_ext}")
                            
                            SUPPORTED_EXTENSIONS = {
                                '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls', 
                                '.pptx', '.ppt', '.msg', '.txt', '.zip', '.7z','.pptm','.htm', '.html' 
                            }
                            
                            if file_ext not in SUPPORTED_EXTENSIONS and file_ext != '.ole':
                                # NON SUPPORTÉ
                                self.log(f"    ⚠️ Type NON SUPPORTÉ: {file_ext}")
                                
                                # ✅ VÉRIFICATION : filename valide
                                if not filename or filename.strip() == '':
                                    self.log(f"    ❌ ERREUR : filename vide, ignoré")
                                    continue
                                
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'keep_unsupported',
                                    'page': page_estimate,
                                    'context': context,
                                    'location': location,
                                    'file_type': file_ext
                                })
                                
                                self.add_to_report('extracted', {
                                    'source_file': Path(docx_path).name,
                                    'extracted_file': f'[Non supporté: {file_ext}]',
                                    'position': f'Page {page_estimate}',
                                    'type': file_ext,
                                    'status': 'Conservé dans le document'
                                })
                            else:
                                # SUPPORTÉ
                                self.log(f"    ✅ Type SUPPORTÉ, extraction...")
                                
                                self.extracted_count += 1
                                output_name = self._generate_output_name(
                                    Path(docx_path).stem, None, file_ext
                                )
                                output_path = Path(output_dir) / output_name
                                
                                with open(output_path, 'wb') as f:
                                    f.write(file_data)
                                
                                self.log(f"    ✅ Extrait: {output_name}")
                                
                                # ✅ VÉRIFICATION : filename valide
                                if not filename or filename.strip() == '':
                                    self.log(f"    ❌ ERREUR : filename vide, ignoré")
                                    continue
                                
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'replace',
                                    'page': page_estimate,
                                    'context': context,
                                    'location': location,
                                    'extracted_name': output_name
                                })
                                
                                self.add_to_report('extracted', {
                                    'source_file': Path(docx_path).name,
                                    'extracted_file': output_name,
                                    'position': f'Page {page_estimate}',
                                    'type': file_ext,
                                    'status': 'Succès'
                                })
                    
                    # ========================================
                    # ÉTAPE 3 : CRÉER DOCX MODIFIÉ
                    # ========================================
                    if extracted_files:
                        self.log(f"\n  🔧 CRÉATION DU WORD MODIFIÉ...")
                        self.log(f"  📊 {len(extracted_files)} fichier(s) à traiter")
                        
                        # Vérification finale : tous les filenames sont valides
                        valid_files = []
                        for item in extracted_files:
                            if item.get('filename') and item['filename'].strip():
                                valid_files.append(item)
                            else:
                                self.log(f"    ⚠️ Fichier ignoré (filename invalide): {item}")
                        
                        self.log(f"  ✅ {len(valid_files)} fichier(s) valides pour modification")
                        
                        if valid_files:
                            self.create_modified_docx_exact_positions(
                                docx_path, output_dir, valid_files
                            )
                        else:
                            self.log(f"    ⚠️ Aucun fichier valide, copie de l'original")
                            original_copy = Path(output_dir) / Path(docx_path).name
                            shutil.copy2(docx_path, original_copy)
                            self.log(f"  ✅ Copie créée: {original_copy.name}")
                    else:
                        # Aucun fichier
                        self.log(f"\n  📄 Aucun fichier incorporé → Copie du fichier original")
                        original_copy = Path(output_dir) / Path(docx_path).name
                        shutil.copy2(docx_path, original_copy)
                        self.log(f"  ✅ Copie créée: {original_copy.name}")
                    
                    self.add_to_report('processed', {
                        'file_name': Path(docx_path).name,
                        'file_type': 'DOCX',
                        'files_found': len(extracted_files),
                        'status': 'Traité avec succès'
                    })
                
                else:
                    # Pas de dossier embeddings
                    self.log(f"\n  ℹ️ Aucun dossier embeddings/ → Aucun fichier incorporé")
                    original_copy = Path(output_dir) / Path(docx_path).name
                    shutil.copy2(docx_path, original_copy)
                    self.log(f"  ✅ Copie créée: {original_copy.name}")
                    
                    self.add_to_report('processed', {
                        'file_name': Path(docx_path).name,
                        'file_type': 'DOCX',
                        'files_found': 0,
                        'status': 'Aucun fichier incorporé'
                    })
        
        except Exception as e:
            self.log(f"  ❌ Erreur: {str(e)}")
            import traceback
            self.log(f"  📋 Détails: {traceback.format_exc()}")
            self.add_to_report('error', {
                'file_name': Path(docx_path).name,
                'error_type': 'Erreur d\'extraction',
                'error_message': str(e),
                'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            })
        
        return extracted_files
    def extract_ole_from_legacy(self, file_path, output_dir):
        """Extrait fichiers des anciens formats Office (.doc, .xls, .ppt)"""
        self.log(f"\n📄 Traitement ancien format: {Path(file_path).name}")
        self.reset_extraction_tracking(Path(file_path).name, output_dir)
        extracted_files = []
        
        try:
            ole = olefile.OleFileIO(file_path)
            all_streams = ole.listdir()
            position = 0
            
            for stream_entry in all_streams:
                stream_name = '/'.join(stream_entry)
                
                if any(x in stream_name.lower() for x in ['mbd', 'objinfo', 'ole', 'package']):
                    try:
                        stream_data = ole.openstream(stream_entry).read()
                        
                        if len(stream_data) > 100:
                            position += 1
                            
                            for offset in range(0, min(len(stream_data), 2000), 4):
                                chunk = stream_data[offset:]
                                
                                if chunk[:4] == b'\xd0\xcf\x11\xe0':
                                    extracted_name = self.extract_ole_content(
                                        chunk, output_dir, Path(file_path).stem, position
                                    )
                                    if extracted_name:
                                        extracted_files.append((position, extracted_name))
                                        self.add_to_report('extracted', {
                                            'source_file': Path(file_path).name,
                                            'extracted_file': extracted_name,
                                            'position': position,
                                            'type': Path(extracted_name).suffix,
                                            'status': 'Succès'
                                        })
                                        break
                                elif chunk[:4] == b'%PDF':
                                    self.extracted_count += 1
                                    
                                    # CORRECTION 3 : Utiliser _generate_output_name
                                    output_name = self._generate_output_name(
                                        Path(file_path).stem, None, '.pdf'
                                    )
                                    output_path = Path(output_dir) / output_name
                                    
                                    with open(output_path, 'wb') as f:
                                        f.write(chunk)
                                    
                                    extracted_files.append((position, output_name))
                                    self.log(f"  ✓ PDF extrait: {output_name}")
                                    self.add_to_report('extracted', {
                                        'source_file': Path(file_path).name,
                                        'extracted_file': output_name,
                                        'position': position,
                                        'type': '.pdf',
                                        'status': 'Succès'
                                    })
                                    break
                                elif chunk[:4] == b'PK\x03\x04':
                                    file_ext = self.detect_office_xml_type(chunk)
                                    self.extracted_count += 1
                                    
                                    # CORRECTION 3 : Utiliser _generate_output_name
                                    output_name = self._generate_output_name(
                                        Path(file_path).stem, None, file_ext
                                    )
                                    output_path = Path(output_dir) / output_name
                                    
                                    with open(output_path, 'wb') as f:
                                        f.write(chunk)
                                    
                                    extracted_files.append((position, output_name))
                                    self.log(f"  ✓ Fichier extrait: {output_name}")
                                    
                                    # CORRECTION 2 : Marquer ZIP pour remplacement
                                    if file_ext == '.zip':
                                        self.log(f"  📌 ZIP sera remplacé dans document modifié")
                                    
                                    self.add_to_report('extracted', {
                                        'source_file': Path(file_path).name,
                                        'extracted_file': output_name,
                                        'position': position,
                                        'type': file_ext,
                                        'status': 'Succès'
                                    })
                                    break
                    except Exception as stream_error:
                        continue
            
            ole.close()
            
            ext = Path(file_path).suffix
            
            
            if extracted_files:
                self.log(f"  ✓ {len(extracted_files)} fichier(s) extrait(s)")
            else:
                # Ne pas copier les fichiers legacy (.doc/.xls/.ppt) dans output_dir
                # La version convertie (.docx/.xlsx/.pptx) est gérée par process_file
                ext = Path(file_path).suffix.lower()
                if ext not in ['.doc', '.xls', '.ppt']:
                    self.log(f"\n  📄 Aucun fichier incorporé → Copie du fichier original")
                    original_copy = Path(output_dir) / Path(file_path).name
                    shutil.copy2(file_path, original_copy)
                    self.log(f"  ✅ Copie créée: {original_copy.name}")
                else:
                    self.log(f"\n  📄 Aucun fichier incorporé — format legacy ignoré (utiliser la version convertie)")

            self.add_to_report('processed', {
                'file_name': Path(file_path).name,
                'file_type': ext.upper().replace('.', ''),
                'files_found': len(extracted_files),
                'status': 'Traité avec succès' if not extracted_files else f'{len(extracted_files)} fichier(s) extrait(s)'
            })
        
        except Exception as e:
            self.log(f"  ❌ Erreur: {str(e)}")
            self.add_to_report('error', {
                'file_name': Path(file_path).name,
                'error_type': 'Erreur d\'extraction',
                'error_message': str(e),
                'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            })
            self.add_to_report('processed', {
                'file_name': Path(file_path).name,
                'file_type': Path(file_path).suffix.upper().replace('.', ''),
                'files_found': 0,
                'status': f'Erreur: {str(e)[:50]}'
            })
        
        return [f[1] for f in extracted_files]
    
    def extract_from_xlsx(self, xlsx_path, output_dir):
        """
        Extrait fichiers Excel avec POSITION EXACTE (feuille + cellule)
        """
        bin_cleaned = self.clean_bin_files_from_embeddings(xlsx_path)
        if bin_cleaned > 0:
            self.log(f"  ✅ Pré-traitement: {bin_cleaned} fichier(s) .bin vide(s) supprimé(s)")
        
        self.reset_extraction_tracking(Path(xlsx_path).name, output_dir)
        Path(output_dir).mkdir(parents=True, exist_ok=True)

        extracted_files = []

        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)
                
                with zipfile.ZipFile(xlsx_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_path)
                
                embeddings_dir = temp_path / "xl" / "embeddings"
                
                if embeddings_dir.exists():
                    # ========================================
                    # ÉTAPE 1 : MAPPER POSITIONS EXACTES
                    # ========================================
                    self.log(f"\n  🔍 MAPPING EXACT DES POSITIONS...")
                    
                    position_mapping = self.get_xlsx_ole_positions_exact(xlsx_path, temp_path)
                    
                    if not position_mapping:
                        self.log(f"    ⚠️ Aucun mapping trouvé")
                        original_copy = Path(output_dir) / Path(xlsx_path).name
                        shutil.copy2(xlsx_path, original_copy)
                        return []
                    
                    # ========================================
                    # ÉTAPE 2 : TRAITER CHAQUE FICHIER
                    # ========================================
                    self.log(f"\n  📦 EXTRACTION ET CLASSIFICATION...")
                    
                    for filename in sorted(position_mapping.keys(), key=lambda x: position_mapping[x]['position']):
                        location = position_mapping[filename]
                        position = location['position']
                        sheet_name = location.get('sheet_name', '?')
                        
                        embed_path = embeddings_dir / filename
                        
                        if not embed_path.exists():
                            continue
                        
                        self.log(f"\n  🔍 Position {position} (Feuille: {sheet_name})")
                        
                        with open(embed_path, 'rb') as f:
                            file_data = f.read()
                        
                        self.log(f"    📦 Fichier: {filename} ({len(file_data)} bytes)")
                        
                        # CAS 1 : VIDE
                        if len(file_data) == 0:
                            extracted_files.append({
                                'position': position,
                                'filename': filename,
                                'action': 'keep_empty',
                                'sheet_name': sheet_name,
                                'location': location
                            })
                            
                            self.add_to_report('extracted', {
                                'source_file': Path(xlsx_path).name,
                                'extracted_file': '[Vide - 0 Ko]',
                                'position': f'Feuille {sheet_name}',
                                'type': 'Vide',
                                'status': 'Ignoré'
                            })
                            continue
                        
                        # CAS 2 : IMAGE
                        if not self.is_file_not_image(filename, file_data):
                            self.log(f"    ⏭️ Image ignorée")
                            continue

                        # CAS 3 : ZIP / Office moderne (xlsx, docx, pptx…)
                        if file_data[:4] == b'PK\x03\x04':
                            self.log(f"    📦 Fichier ZIP/Office moderne détecté")
                            file_ext = self.detect_office_xml_type(file_data)
                            self.log(f"    🔍 Type: {file_ext}")

                            SUPPORTED_ZIP_XLSX = {
                                '.xlsx', '.xlsm', '.docx', '.pptx', '.pptm', '.zip', '.7z'
                            }

                            if file_ext in SUPPORTED_ZIP_XLSX:
                                self.extracted_count += 1
                                output_name = self._generate_output_name(
                                    Path(xlsx_path).stem, None, file_ext
                                )
                                output_path = Path(output_dir) / output_name
                                with open(output_path, 'wb') as f:
                                    f.write(file_data)
                                self.log(f"    ✅ Extrait (ZIP/Office): {output_name}")
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'replace',
                                    'sheet_name': sheet_name,
                                    'location': location,
                                    'extracted_name': output_name
                                })
                                self.add_to_report('extracted', {
                                    'source_file': Path(xlsx_path).name,
                                    'extracted_file': output_name,
                                    'position': f'Feuille {sheet_name}',
                                    'type': file_ext,
                                    'status': 'Succès'
                                })
                            else:
                                self.log(f"    ⚠️ Type ZIP non supporté: {file_ext}")
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'keep_unsupported',
                                    'sheet_name': sheet_name,
                                    'location': location,
                                    'file_type': file_ext
                                })
                                self.add_to_report('extracted', {
                                    'source_file': Path(xlsx_path).name,
                                    'extracted_file': f'[Non supporté: {file_ext}]',
                                    'position': f'Feuille {sheet_name}',
                                    'type': file_ext,
                                    'status': 'Conservé'
                                })
                            continue

                        # CAS 4 : OLE
                        if file_data[:4] == b'\xd0\xcf\x11\xe0':
                            file_ext = self.detect_file_type(file_data)
                            
                            SUPPORTED = {'.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls', '.pptx', '.ppt', '.msg', '.txt', '.zip', '.7z','.htm', '.html' }
                            
                            if file_ext not in SUPPORTED and file_ext != '.ole':
                                # NON SUPPORTÉ
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'keep_unsupported',
                                    'sheet_name': sheet_name,
                                    'location': location,
                                    'file_type': file_ext
                                })
                                
                                self.add_to_report('extracted', {
                                    'source_file': Path(xlsx_path).name,
                                    'extracted_file': f'[Non supporté: {file_ext}]',
                                    'position': f'Feuille {sheet_name}',
                                    'type': file_ext,
                                    'status': 'Conservé'
                                })
                                continue
                            
                            # SUPPORTÉ
                            extracted_name = self.extract_ole_content(
                                file_data, output_dir, Path(xlsx_path).stem, position
                            )
                            
                            if extracted_name:
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'replace',
                                    'sheet_name': sheet_name,
                                    'location': location,
                                    'extracted_name': extracted_name
                                })

                                self.add_to_report('extracted', {
                                    'source_file': Path(xlsx_path).name,
                                    'extracted_file': extracted_name,
                                    'position': f'Feuille {sheet_name}',
                                    'type': Path(extracted_name).suffix,
                                    'status': 'Succès'
                                })

                        # CAS 5 : Fichier direct (non-OLE, non-ZIP, ex: PDF brut)
                        else:
                            self.log(f"    📄 Fichier direct (non-OLE, non-ZIP) détecté")
                            file_ext = self.detect_file_type(file_data)
                            self.log(f"    🔍 Type: {file_ext}")

                            SUPPORTED_DIRECT_XLSX = {
                                '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls',
                                '.pptx', '.pptm', '.ppt', '.msg', '.txt', '.zip', '.7z',
                                '.htm', '.html'
                            }

                            if file_ext in SUPPORTED_DIRECT_XLSX:
                                self.extracted_count += 1
                                output_name = self._generate_output_name(
                                    Path(xlsx_path).stem, None, file_ext
                                )
                                output_path = Path(output_dir) / output_name
                                with open(output_path, 'wb') as f:
                                    f.write(file_data)
                                self.log(f"    ✅ Extrait (direct): {output_name}")
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'replace',
                                    'sheet_name': sheet_name,
                                    'location': location,
                                    'extracted_name': output_name
                                })
                                self.add_to_report('extracted', {
                                    'source_file': Path(xlsx_path).name,
                                    'extracted_file': output_name,
                                    'position': f'Feuille {sheet_name}',
                                    'type': file_ext,
                                    'status': 'Succès'
                                })
                            else:
                                self.log(f"    ⚠️ Type direct non supporté: {file_ext}")
                                extracted_files.append({
                                    'position': position,
                                    'filename': filename,
                                    'action': 'keep_unsupported',
                                    'sheet_name': sheet_name,
                                    'location': location,
                                    'file_type': file_ext
                                })
                                self.add_to_report('extracted', {
                                    'source_file': Path(xlsx_path).name,
                                    'extracted_file': f'[Non supporté: {file_ext}]',
                                    'position': f'Feuille {sheet_name}',
                                    'type': file_ext,
                                    'status': 'Conservé'
                                })

                    # ========================================
                    # ÉTAPE 3 : CRÉER XLSX MODIFIÉ
                    # ========================================
                    if extracted_files:
                        self.log(f"\n  🔧 CRÉATION DU EXCEL MODIFIÉ...")
                        self.create_modified_xlsx_exact_positions(
                            xlsx_path, output_dir, extracted_files
                        )
                    else:
                        original_copy = Path(output_dir) / Path(xlsx_path).name
                        shutil.copy2(xlsx_path, original_copy)
                    
                    self.add_to_report('processed', {
                        'file_name': Path(xlsx_path).name,
                        'file_type': Path(xlsx_path).suffix.upper().replace('.', ''),
                        'files_found': len(extracted_files),
                        'status': 'Traité avec succès'
                    })

                else:
                    # Pas de dossier xl/embeddings/ → aucun fichier incorporé
                    self.log(f"\n  ℹ️ Aucun dossier embeddings/ → Aucun fichier incorporé")
                    original_copy = Path(output_dir) / Path(xlsx_path).name
                    shutil.copy2(xlsx_path, original_copy)
                    self.log(f"  ✅ Copie créée: {original_copy.name}")
                    self.add_to_report('processed', {
                        'file_name': Path(xlsx_path).name,
                        'file_type': Path(xlsx_path).suffix.upper().replace('.', ''),
                        'files_found': 0,
                        'status': 'Aucun fichier incorporé'
                    })

        except Exception as e:
            self.log(f"  ❌ Erreur: {str(e)}")
            import traceback
            self.log(f"  📋 Détails: {traceback.format_exc()}")
            # Copie de secours si dossier de sortie vide
            try:
                original_copy = Path(output_dir) / Path(xlsx_path).name
                if not original_copy.exists():
                    shutil.copy2(xlsx_path, original_copy)
                    self.log(f"  ⚠️ Erreur → copie de l'original créée: {original_copy.name}")
            except Exception:
                pass

        return extracted_files


    def extract_from_pptx(self, pptx_path, output_dir):
        """
        Extrait fichiers PowerPoint avec POSITION EXACTE conservée.
        Supporte .pptx et .pptm (macros VBA).
        CORRECTIONS :
        - ÉTAPE 1 : lecture .rels pour mapper rId → embed_file
        - ÉTAPE 1 : stockage du shape_id (cNvPr id)
        - CAS 4   : support ZIP/xlsx direct + add_to_report si non supporté
        - CAS 4   : .pptm ajouté dans SUPPORTED_ZIP
        - CAS 5   : add_to_report appelé SYSTÉMATIQUEMENT (même si extraction échoue)
        - CAS 5   : SUPPORTED_OLE correctement défini (corrigé NameError)
        - CAS 4   : '.7z' avec le point (corrigé)
        """
        pptx_path = Path(pptx_path)
        file_ext_src = pptx_path.suffix.lower()
        file_type_label = 'PPTM' if file_ext_src == '.pptm' else 'PPTX'

        bin_cleaned = self.clean_bin_files_from_embeddings(pptx_path)
        if bin_cleaned > 0:
            self.log(f"  ✅ Pré-traitement: {bin_cleaned} fichier(s) .bin vide(s) supprimé(s)")

        self.reset_extraction_tracking(pptx_path.name, output_dir)
        Path(output_dir).mkdir(parents=True, exist_ok=True)

        extracted_files = []

        # ── Extensions supportées (partagées par tous les CAS) ──────────────
        SUPPORTED_ZIP = {
            '.xlsx', '.xlsm', '.docx', '.pptx', '.pptm',  # ← .pptm ajouté
            '.zip', '.7z'
        }
        SUPPORTED_OLE = {
            '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls',
            '.pptx', '.pptm', '.ppt', '.msg', '.txt', '.zip', '.7z' ,'.htm', '.html'  # ← .pptm ajouté
        }
        SUPPORTED_DIRECT = {
            '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls',
            '.pptx', '.pptm', '.ppt', '.msg', '.txt', '.zip', '.7z' ,'.htm', '.html'  # ← .pptm ajouté
        }

        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)

                with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_path)

                embeddings_dir = temp_path / "ppt" / "embeddings"

                if embeddings_dir.exists():
                    embed_list = list(embeddings_dir.iterdir())

                    def get_file_num(file_path):
                        match = re.search(r'oleObject(\d+)', file_path.name)
                        return int(match.group(1)) if match else 999

                    embed_list.sort(key=get_file_num)

                    self.log(f"  📂 {len(embed_list)} fichier(s) dans embeddings/")
                    for ef in embed_list:
                        self.log(f"      {ef.name} ({ef.stat().st_size} bytes)")

                    # ========================================
                    # ÉTAPE 1 : MAPPER POSITION → SLIDE + SHAPE
                    # via .rels + shape_id (cNvPr id)
                    # ========================================
                    self.log(f"\n  🔍 MAPPING EXACT DES POSITIONS...")

                    position_to_location = {}

                    slides_dir = temp_path / "ppt" / "slides"
                    slide_files = sorted(
                        [f for f in slides_dir.glob('slide*.xml')],
                        key=lambda x: int(re.search(r'slide(\d+)', x.name).group(1))
                    )

                    embed_name_to_file = {ef.name: ef for ef in embed_list}

                    global_position = 0

                    for slide_file in slide_files:
                        slide_num = int(re.search(r'slide(\d+)', slide_file.name).group(1))

                        # Lire le .rels : rId → nom fichier embed
                        rels_file = slides_dir / '_rels' / (slide_file.name + '.rels')
                        if not rels_file.exists():
                            continue

                        import xml.etree.ElementTree as _ET
                        _rels_root = _ET.parse(str(rels_file)).getroot()
                        _ns_rel = 'http://schemas.openxmlformats.org/package/2006/relationships'

                        rid_to_embed = {}
                        for _rel in _rels_root.findall(f'{{{_ns_rel}}}Relationship'):
                            _target = _rel.get('Target', '')
                            if 'embeddings/' in _target:
                                rid_to_embed[_rel.get('Id')] = _target.split('/')[-1]

                        if not rid_to_embed:
                            continue

                        self.log(f"    Slide {slide_num} : {len(rid_to_embed)} embedding(s) dans .rels")

                        # Parser le XML de la slide
                        from lxml import etree as _lxml_etree
                        _slide_root = _lxml_etree.parse(str(slide_file)).getroot()

                        _ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

                        _spTree = _slide_root.find(f'.//{{{_ns_p}}}spTree')
                        if _spTree is None:
                            _spTree = _slide_root.find('.//spTree')
                        if _spTree is None:
                            self.log(f"    ⚠️ spTree introuvable dans slide{slide_num}")
                            continue

                        shape_index = 0
                        for shape_elem in _spTree:
                            shape_index += 1
                            shape_xml = _lxml_etree.tostring(shape_elem, encoding='unicode')

                            r_ids = re.findall(r'r:id="([^"]+)"', shape_xml)
                            r_ids += re.findall(
                                r'\{http://schemas\.openxmlformats\.org/officeDocument/2006/relationships\}id="([^"]+)"',
                                shape_xml
                            )

                            for r_id in r_ids:
                                if r_id in rid_to_embed:
                                    embed_name = rid_to_embed[r_id]
                                    if embed_name in embed_name_to_file:
                                        global_position += 1
                                        embed_file = embed_name_to_file[embed_name]

                                        fn_match = re.search(r'oleObject(\d+)', embed_name)
                                        file_num = int(fn_match.group(1)) if fn_match else global_position

                                        # Récupérer le shape_id (cNvPr id)
                                        sp_element_id = None
                                        for _elem in shape_elem.iter():
                                            tag = _elem.tag
                                            if tag.endswith('}cNvPr') or tag == 'cNvPr':
                                                try:
                                                    sp_element_id = int(_elem.get('id', 0))
                                                except (ValueError, TypeError):
                                                    sp_element_id = None
                                                break

                                        position_to_location[global_position] = {
                                            'slide_number': slide_num,
                                            'slide_file': slide_file,
                                            'shape_index': shape_index,
                                            'sp_element_id': sp_element_id,
                                            'file_num': file_num,
                                            'embed_file': embed_file
                                        }
                                        self.log(f"    ✅ Position {global_position} → Slide {slide_num}, cNvPr id={sp_element_id} ({embed_name})")
                                        break

                    self.log(f"\n  📊 Mapping : {len(position_to_location)} position(s) trouvée(s)")

                    if not position_to_location:
                        self.log(f"    ⚠️ Mapping vide — extraction directe de {len(embed_list)} fichier(s) dans embeddings/")
                        # Fallback : extraire directement chaque fichier dans ppt/embeddings/
                        fallback_extracted = []
                        for pos_fb, ef in enumerate(embed_list, start=1):
                            with open(ef, 'rb') as f:
                                fb_data = f.read()
                            if len(fb_data) == 0:
                                continue
                            if not self.is_file_not_image(ef.name, fb_data):
                                continue
                            if fb_data[:4] == b'PK\x03\x04':
                                fb_ext = self.detect_office_xml_type(fb_data)
                            elif fb_data[:4] == b'\xd0\xcf\x11\xe0':
                                fb_ext = self.detect_file_type(fb_data)
                            elif fb_data[:4] == b'%PDF':
                                fb_ext = '.pdf'
                            else:
                                fb_ext = self.detect_file_type(fb_data)

                            SUPPORTED_FB = {'.pdf','.docx','.doc','.xlsx','.xlsm','.xls',
                                            '.pptx','.pptm','.ppt','.msg','.txt',
                                            '.zip','.7z','.htm','.html'}
                            if fb_ext not in SUPPORTED_FB:
                                self.log(f"    ⚠️ Fallback: type non supporté {fb_ext} ({ef.name})")
                                continue

                            self.extracted_count += 1
                            fb_name = self._generate_output_name(pptx_path.stem, None, fb_ext)
                            fb_path = Path(output_dir) / fb_name
                            with open(fb_path, 'wb') as f:
                                f.write(fb_data)
                            self.log(f"    ✅ Fallback extrait: {fb_name}")
                            fallback_extracted.append({
                                'position': pos_fb, 'filename': fb_name, 'action': 'replace',
                                'slide_number': 1, 'shape_index': pos_fb, 'sp_element_id': None
                            })
                            self.add_to_report('extracted', {
                                'source_file': pptx_path.name,
                                'extracted_file': fb_name,
                                'position': f'Embed {pos_fb}',
                                'type': fb_ext, 'status': 'Succès (fallback)'
                            })

                        original_copy = Path(output_dir) / pptx_path.name
                        if not original_copy.exists() or str(original_copy.resolve()) != str(pptx_path.resolve()):
                            shutil.copy2(pptx_path, original_copy)
                        self.add_to_report('processed', {
                            'file_name': pptx_path.name, 'file_type': file_type_label,
                            'files_found': len(fallback_extracted),
                            'status': 'Fallback (mapping vide)'
                        })
                        return fallback_extracted

                    # ========================================
                    # ÉTAPE 2 : TRAITER CHAQUE FICHIER
                    # ========================================
                    self.log(f"\n  📦 EXTRACTION ET CLASSIFICATION...")

                    for position in sorted(position_to_location.keys()):
                        location      = position_to_location[position]
                        embed_file    = location['embed_file']
                        slide_num     = location['slide_number']
                        shape_index   = location['shape_index']
                        sp_element_id = location.get('sp_element_id')

                        self.log(f"\n  🔍 Position {position} (Slide {slide_num}, shape_id={sp_element_id})")

                        with open(embed_file, 'rb') as f:
                            file_data = f.read()

                        self.log(f"    📦 Fichier: {embed_file.name} ({len(file_data)} bytes)")

                        # ── CAS 1 : FICHIER VIDE ─────────────────────────────────────
                        if len(file_data) == 0:
                            self.log(f"    ⚠️ Fichier VIDE détecté")
                            extracted_files.append({
                                'position': position,
                                'slide_number': slide_num,
                                'shape_index': shape_index,
                                'sp_element_id': sp_element_id,
                                'action': 'keep_empty',
                                'filename': None
                            })
                            self.add_to_report('extracted', {
                                'source_file': pptx_path.name,
                                'extracted_file': '[Fichier vide - 0 Ko]',
                                'position': f'Slide {slide_num}',
                                'type': 'Vide',
                                'status': 'Conservé dans le document — type non supporté SharePoint'
                            })
                            continue

                        # ── CAS 2 : IMAGE (ignorer) ───────────────────────────────────
                        if not self.is_file_not_image(embed_file.name, file_data):
                            self.log(f"    ⏭️ Image ignorée")
                            continue

                        # ── CAS 3 : FICHIER MÉDIA (ignorer) ──────────────────────────
                        if embed_file.suffix.lower() in ['.mp4', '.avi', '.mov', '.mp3', '.wav']:
                            self.log(f"    ⏭️ Fichier média ignoré")
                            continue

                        # ── CAS 4 : ZIP / Office moderne (xlsx, docx, pptx, pptm...) ─
                        if file_data[:4] == b'PK\x03\x04':
                            self.log(f"    📦 Fichier ZIP/Office moderne détecté")
                            file_ext = self.detect_office_xml_type(file_data)
                            self.log(f"    🔍 Type: {file_ext}")

                            if file_ext in SUPPORTED_ZIP:
                                self.extracted_count += 1
                                output_name = self._generate_output_name(
                                    pptx_path.stem, None, file_ext
                                )
                                output_path = Path(output_dir) / output_name
                                with open(output_path, 'wb') as f:
                                    f.write(file_data)
                                self.log(f"    ✅ Extrait (ZIP/Office): {output_name}")
                                extracted_files.append({
                                    'position': position,
                                    'slide_number': slide_num,
                                    'shape_index': shape_index,
                                    'sp_element_id': sp_element_id,
                                    'action': 'replace',
                                    'filename': output_name
                                })
                                self.add_to_report('extracted', {
                                    'source_file': pptx_path.name,
                                    'extracted_file': output_name,
                                    'position': f'Slide {slide_num}',
                                    'type': file_ext,
                                    'status': 'Succès'
                                })
                            else:
                                self.log(f"    ⚠️ Type ZIP non supporté: {file_ext}")
                                extracted_files.append({
                                    'position': position,
                                    'slide_number': slide_num,
                                    'shape_index': shape_index,
                                    'sp_element_id': sp_element_id,
                                    'action': 'keep_unsupported',
                                    'filename': None,
                                    'file_type': file_ext
                                })
                                self.add_to_report('extracted', {
                                    'source_file': pptx_path.name,
                                    'extracted_file': f'[Non supporté: {file_ext}]',
                                    'position': f'Slide {slide_num}',
                                    'type': file_ext,
                                    'status': 'Conservé dans le document — type non supporté SharePoint'
                                })
                            continue

                        # ── CAS 5 : OLE binaire ───────────────────────────────────────
                        if file_data[:4] == b'\xd0\xcf\x11\xe0':
                            self.log(f"    📦 Fichier OLE détecté")

                            file_ext = self.detect_file_type(file_data)
                            self.log(f"    🔍 Type OLE: {file_ext}")

                            # Type explicitement non supporté (CATIA, etc.)
                            if file_ext not in SUPPORTED_OLE and file_ext != '.ole':
                                self.log(f"    ⚠️ Type non supporté: {file_ext}")
                                extracted_files.append({
                                    'position': position,
                                    'slide_number': slide_num,
                                    'shape_index': shape_index,
                                    'sp_element_id': sp_element_id,
                                    'action': 'keep_unsupported',
                                    'filename': None,
                                    'file_type': file_ext
                                })
                                self.add_to_report('extracted', {
                                    'source_file': pptx_path.name,
                                    'extracted_file': f'[Non supporté: {file_ext}]',
                                    'position': f'Slide {slide_num}',
                                    'type': file_ext,
                                    'status': 'Conservé dans le document — type non supporté SharePoint'
                                })
                                continue

                            # Type .ole ou supporté → tenter extraction
                            extracted_name = self.extract_ole_content(
                                file_data, output_dir, pptx_path.stem, position
                            )

                            if extracted_name:
                                self.log(f"    ✅ Extrait (OLE): {extracted_name}")
                                extracted_files.append({
                                    'position': position,
                                    'slide_number': slide_num,
                                    'shape_index': shape_index,
                                    'sp_element_id': sp_element_id,
                                    'action': 'replace',
                                    'filename': extracted_name
                                })
                                self.add_to_report('extracted', {
                                    'source_file': pptx_path.name,
                                    'extracted_file': extracted_name,
                                    'position': f'Slide {slide_num}',
                                    'type': Path(extracted_name).suffix,
                                    'status': 'Succès'
                                })
                            else:
                                # Extraction échouée → conservé dans document
                                self.log(f"    ⚠️ Extraction OLE échouée (type: {file_ext}) — conservé dans document")
                                extracted_files.append({
                                    'position': position,
                                    'slide_number': slide_num,
                                    'shape_index': shape_index,
                                    'sp_element_id': sp_element_id,
                                    'action': 'keep_unsupported',
                                    'filename': None,
                                    'file_type': file_ext
                                })
                                self.add_to_report('extracted', {
                                    'source_file': pptx_path.name,
                                    'extracted_file': f'[Non extractible: {file_ext}]',
                                    'position': f'Slide {slide_num}',
                                    'type': file_ext,
                                    'status': 'Conservé dans le document — type non supporté SharePoint'
                                })
                            continue

                        # ── CAS 6 : Fichier direct (non-OLE, non-ZIP) ────────────────
                        self.log(f"    📄 Fichier direct (non-OLE) détecté")
                        file_ext = self.detect_file_type(file_data)
                        self.log(f"    🔍 Type: {file_ext}")

                        if file_ext not in SUPPORTED_DIRECT and file_ext != '.ole':
                            self.log(f"    ⚠️ Type non supporté: {file_ext}")
                            extracted_files.append({
                                'position': position,
                                'slide_number': slide_num,
                                'shape_index': shape_index,
                                'sp_element_id': sp_element_id,
                                'action': 'keep_unsupported',
                                'filename': None,
                                'file_type': file_ext
                            })
                            self.add_to_report('extracted', {
                                'source_file': pptx_path.name,
                                'extracted_file': f'[Non supporté: {file_ext}]',
                                'position': f'Slide {slide_num}',
                                'type': file_ext,
                                'status': 'Conservé dans le document — type non supporté SharePoint'
                            })
                        else:
                            self.extracted_count += 1
                            output_name = self._generate_output_name(
                                pptx_path.stem, None, file_ext
                            )
                            output_path = Path(output_dir) / output_name
                            with open(output_path, 'wb') as f:
                                f.write(file_data)
                            self.log(f"    ✅ Extrait: {output_name}")
                            extracted_files.append({
                                'position': position,
                                'slide_number': slide_num,
                                'shape_index': shape_index,
                                'sp_element_id': sp_element_id,
                                'action': 'replace',
                                'filename': output_name
                            })
                            self.add_to_report('extracted', {
                                'source_file': pptx_path.name,
                                'extracted_file': output_name,
                                'position': f'Slide {slide_num}',
                                'type': file_ext,
                                'status': 'Succès'
                            })

                    # ========================================
                    # ÉTAPE 3 : CRÉER PPTX/PPTM MODIFIÉ
                    # ========================================
                    if extracted_files:
                        self.log(f"\n  🔧 CRÉATION DU POWERPOINT MODIFIÉ...")
                        self.create_modified_pptx_exact_positions(
                            pptx_path, output_dir, extracted_files, temp_path
                        )
                    else:
                        self.log(f"\n  📄 Aucun fichier incorporé → Copie du fichier original")
                        original_copy = Path(output_dir) / pptx_path.name
                        shutil.copy2(pptx_path, original_copy)
                        self.log(f"  ✅ Copie créée: {original_copy.name}")

                    self.add_to_report('processed', {
                        'file_name': pptx_path.name,
                        'file_type': file_type_label,
                        'files_found': len(extracted_files),
                        'status': 'Traité avec succès'
                    })

                else:
                    self.log(f"\n  ℹ️ Aucun dossier embeddings/ → Aucun fichier incorporé")
                    original_copy = Path(output_dir) / pptx_path.name
                    shutil.copy2(pptx_path, original_copy)
                    self.log(f"  ✅ Copie créée: {original_copy.name}")
                    self.add_to_report('processed', {
                        'file_name': pptx_path.name,
                        'file_type': file_type_label,
                        'files_found': 0,
                        'status': 'Aucun fichier incorporé'
                    })

        except Exception as e:
            self.log(f"  ❌ Erreur: {str(e)}")
            import traceback
            self.log(f"  📋 Détails: {traceback.format_exc()}")
            self.add_to_report('error', {
                'file_name': pptx_path.name,
                'error_type': 'Erreur d\'extraction',
                'error_message': str(e),
                'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            })
            # Copie de secours si dossier de sortie vide
            try:
                original_copy = Path(output_dir) / pptx_path.name
                if not original_copy.exists():
                    shutil.copy2(pptx_path, original_copy)
                    self.log(f"  ⚠️ Erreur → copie de l'original créée: {original_copy.name}")
            except Exception:
                pass

        return extracted_files
    def validate_ole_mapping(self, docx_path, position_mapping):
        """
        Valide que le mapping des positions est cohérent
        """
        self.log(f"\n  🔍 VALIDATION DU MAPPING :")
        
        if not position_mapping:
            self.log(f"    ⚠️ Aucun mapping trouvé")
            return False
        
        # Vérifier que les positions sont séquentielles
        positions = sorted(position_mapping.values())
        expected = list(range(1, len(positions) + 1))
        
        if positions != expected:
            self.log(f"    ⚠️ AVERTISSEMENT : Positions non séquentielles")
            self.log(f"       Trouvées: {positions}")
            self.log(f"       Attendues: {expected}")
            return False
        
        self.log(f"    ✅ Mapping validé : {len(positions)} position(s) séquentielles")
        return True
    def create_modified_docx_exact_positions(self, original_path, output_dir, extracted_files):
        """
        Crée Word modifié en CONSERVANT LES POSITIONS EXACTES
        CORRECTION : Utilise les noms de fichiers pour mapper correctement
        """
        try:
            doc = Document(original_path)
            
            # ========================================
            # ÉTAPE 1 : CRÉER MAPPING FILENAME → PARAGRAPH/RUN
            # ========================================
            self.log(f"\n  🔍 MAPPING EXACT PAR NOM DE FICHIER...")
            
            # Charger le XML pour obtenir les noms de fichiers
            from lxml import etree
            import zipfile
            
            with zipfile.ZipFile(original_path, 'r') as zf:
                doc_xml = zf.read('word/document.xml')
                rels_xml = zf.read('word/_rels/document.xml.rels')
            
            doc_root = etree.fromstring(doc_xml)
            rels_root = etree.fromstring(rels_xml)
            
            ns = {
                'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'
            }
            
            # Mapping r:id → filename
            rid_to_filename = {}
            for rel in rels_root.xpath('rel:Relationship', namespaces={'rel': ns['rel']}):
                target = rel.get('Target')
                if target and 'embeddings/' in target:
                    r_id = rel.get('Id')
                    filename = target.split('/')[-1]
                    rid_to_filename[r_id] = filename
            
            # Créer mapping filename → (paragraph, run)
            filename_to_element = {}
            
            # Parcourir paragraphes
            for para_idx, para in enumerate(doc.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    # Chercher r:id dans le XML du run
                    run_xml = etree.tostring(run._element, encoding='unicode')
                    
                    for r_id, filename in rid_to_filename.items():
                        if r_id in run_xml:
                            filename_to_element[filename] = {
                                'paragraph': para,
                                'run': run,
                                'type': 'paragraph'
                            }
                            self.log(f"    ✓ Mappé: {filename} → Paragraphe #{para_idx}")
                            break
            
            # Parcourir tableaux
            for table_idx, table in enumerate(doc.tables):
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        for para_idx, para in enumerate(cell.paragraphs):
                            for run_idx, run in enumerate(para.runs):
                                run_xml = etree.tostring(run._element, encoding='unicode')
                                
                                for r_id, filename in rid_to_filename.items():
                                    if r_id in run_xml:
                                        filename_to_element[filename] = {
                                            'paragraph': para,
                                            'run': run,
                                            'type': 'table_cell',
                                            'table': table_idx,
                                            'row': row_idx,
                                            'cell': cell_idx
                                        }
                                        self.log(f"    ✓ Mappé: {filename} → Table {table_idx} L{row_idx} C{cell_idx}")
                                        break
            
            self.log(f"  ✅ {len(filename_to_element)} fichier(s) mappé(s)")
            
            # ========================================
            # ÉTAPE 2 : APPLIQUER MODIFICATIONS PAR FILENAME
            # ========================================
            self.log(f"\n  🔧 APPLICATION DES MODIFICATIONS...")
            
            modifications_applied = 0
            
            for item in extracted_files:
                filename = item.get('filename')
                action = item['action']
                page = item.get('page', '?')
                
                # VÉRIFICATION CRITIQUE : Le fichier doit être dans le mapping
                if filename not in filename_to_element:
                    self.log(f"    ⚠️ FICHIER NON TROUVÉ DANS DOCUMENT : {filename}")
                    self.log(f"       → Aucune modification appliquée pour cette position")
                    continue
                
                element_info = filename_to_element[filename]
                para = element_info['paragraph']
                run = element_info['run']
                
                # ========================================
                # CAS 1 : FICHIER VIDE
                # ========================================
                if action == 'keep_empty':
                    warning_para = para.insert_paragraph_before()
                    warning_run = warning_para.add_run()
                    warning_run.text = f"⚠️ ATTENTION (Page ~{page}): Fichier vide (0 Ko) ci-dessous - non archivé SharePoint"
                    warning_run.bold = True
                    warning_run.font.color.rgb = RGBColor(255, 153, 0)
                    warning_run.font.size = Pt(11)
                    
                    self.log(f"    ✅ Message VIDE : {filename} (Page ~{page})")
                    modifications_applied += 1
                
                # ========================================
                # CAS 2 : NON SUPPORTÉ
                # ========================================
                elif action == 'keep_unsupported':
                    warning_para = para.insert_paragraph_before()
                    warning_run = warning_para.add_run()
                    warning_run.text = f"⚠️ ATTENTION (Page ~{page}): Type de fichier ci-dessous non supporté SharePoint"
                    warning_run.bold = True
                    warning_run.font.color.rgb = RGBColor(255, 153, 0)
                    warning_run.font.size = Pt(11)
                    
                    self.log(f"    ✅ Message NON SUPPORTÉ : {filename} (Page ~{page})")
                    modifications_applied += 1
                
                # ========================================
                # CAS 3 : REMPLACER
                # ========================================
                elif action == 'replace':
                    extracted_name = item.get('extracted_name')
                    
                    if not extracted_name:
                        self.log(f"    ⚠️ Pas de nom extrait pour {filename}")
                        continue
                    
                    # Supprimer objets XML
                    for obj in run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}object'):
                        obj.getparent().remove(obj)
                    
                    for pict in run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pict'):
                        pict.getparent().remove(pict)
                    
                    # Remplacer par texte
                    run.text = ''
                    run.text = f"Voir  {Path(extracted_name).stem}"
                    run.bold = True
                    run.font.color.rgb = RGBColor(255, 0, 0)
                    run.font.size = Pt(12)
                    
                    self.log(f"    ✅ REMPLACÉ : {filename} → {extracted_name} (Page ~{page})")
                    modifications_applied += 1
            
            # Sauvegarder
            modified_path = Path(output_dir) / Path(original_path).name
            doc.save(str(modified_path))
            
            self.log(f"\n  ✅ Word modifié : {modified_path.name}")
            self.log(f"  ✅ {modifications_applied} modification(s) appliquée(s)")
            
            if modifications_applied != len(extracted_files):
                self.log(f"  ⚠️ ATTENTION : {len(extracted_files)} fichiers extraits mais seulement {modifications_applied} modifications appliquées")
            
        except Exception as e:
            self.log(f"  ❌ Erreur: {str(e)}")
            import traceback
            self.log(f"  📋 Détails: {traceback.format_exc()}")
            
            # Copie de secours
            try:
                modified_path = Path(output_dir) / Path(original_path).name
                shutil.copy2(original_path, modified_path)
                self.log(f"  ⚠️ Copie de l'original créée")
            except:
                pass
    def create_modified_xlsx_exact_positions(self, original_path, output_dir, extracted_files):
        """
        Crée Excel modifié :
        - Supprime les objets OLE extraits du worksheet XML et drawing XML
        - Insère "Voir [nom_fichier]" en rouge dans la cellule correspondante
        """
        try:
            from lxml import etree as _et

            original_path = Path(original_path)
            modified_path = Path(output_dir) / original_path.name

            ns_rel = 'http://schemas.openxmlformats.org/package/2006/relationships'
            ns_xdr = 'http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing'
            ns_r   = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

            with tempfile.TemporaryDirectory() as temp_dir:
                temp_path = Path(temp_dir)

                # ── 1. Extraire le XLSX ──────────────────────────────────
                with zipfile.ZipFile(original_path, 'r') as zf:
                    zf.extractall(temp_path)

                # ── 2. Grouper les items 'replace' par sheet_number ──────
                sheets_actions = {}
                for item in extracted_files:
                    if item['action'] != 'replace':
                        continue
                    sheet_num = item.get('location', {}).get('sheet_number')
                    if sheet_num is None:
                        continue
                    sheets_actions.setdefault(sheet_num, []).append(item)

                # ── 3. Pour chaque feuille : supprimer shapes OLE ────────
                cells_to_fill = []  # (sheet_name, col, row, text)

                for sheet_num, items in sheets_actions.items():
                    ws_path = (temp_path / 'xl' / 'worksheets'
                               / f'sheet{sheet_num}.xml')
                    ws_rels_path = (temp_path / 'xl' / 'worksheets' / '_rels'
                                    / f'sheet{sheet_num}.xml.rels')

                    if not ws_path.exists() or not ws_rels_path.exists():
                        self.log(f"    ⚠️ sheet{sheet_num}.xml ou rels introuvable")
                        continue

                    # ── 3a. Sheet rels → embed_filename → rel_id ─────────
                    # Les embeddings OLE sont référencés dans les sheet rels,
                    # PAS dans les drawing rels.
                    ws_rels_tree = _et.parse(str(ws_rels_path))
                    ws_rels_root = ws_rels_tree.getroot()

                    embed_to_relid = {}   # 'oleObject1.bin' → 'rId5'
                    relid_to_embed = {}   # 'rId5' → 'oleObject1.bin'
                    for rel in ws_rels_root.findall(f'{{{ns_rel}}}Relationship'):
                        tgt = rel.get('Target', '')
                        if 'embeddings/' in tgt:
                            basename = Path(tgt).name
                            rid = rel.get('Id')
                            embed_to_relid[basename] = rid
                            relid_to_embed[rid] = basename

                    items_by_embed = {item['filename']: item for item in items}
                    rel_ids_to_remove = {
                        embed_to_relid[fn]
                        for fn in items_by_embed
                        if fn in embed_to_relid
                    }

                    if not rel_ids_to_remove:
                        self.log(f"    ⚠️ Aucune relation embedding dans sheet rels "
                                 f"pour feuille {sheet_num} "
                                 f"(cherché: {list(items_by_embed.keys())}, "
                                 f"trouvé: {list(embed_to_relid.keys())})")
                        continue

                    # ── 3b. Worksheet XML → <oleObject r:id shapeId> ──────
                    ws_tree = _et.parse(str(ws_path))
                    ws_root = ws_tree.getroot()

                    relid_to_shapeid = {}     # 'rId5' → 1025
                    extra_relids_to_remove = set()  # image rels dans objectPr

                    for elem in ws_root.iter():
                        local = (elem.tag.split('}')[-1]
                                 if '}' in elem.tag else elem.tag)
                        if local != 'oleObjects':
                            continue
                        for ole_obj in list(elem):
                            ol = (ole_obj.tag.split('}')[-1]
                                  if '}' in ole_obj.tag else ole_obj.tag)
                            if ol != 'oleObject':
                                continue
                            r_id = ole_obj.get(f'{{{ns_r}}}id')
                            shape_id = ole_obj.get('shapeId')
                            if r_id in rel_ids_to_remove and shape_id:
                                relid_to_shapeid[r_id] = int(shape_id)
                                # Collecter les rel d'image (objectPr)
                                for child in ole_obj.iter():
                                    cl = (child.tag.split('}')[-1]
                                          if '}' in child.tag else child.tag)
                                    if cl == 'objectPr':
                                        pr_rid = child.get(f'{{{ns_r}}}id')
                                        if pr_rid:
                                            extra_relids_to_remove.add(pr_rid)
                                elem.remove(ole_obj)
                                self.log(f"    🗑️ <oleObject r:id={r_id} "
                                         f"shapeId={shape_id}> supprimé du XML")
                        # Supprimer <oleObjects> si vide
                        if len(elem) == 0:
                            parent = elem.getparent()
                            if parent is not None:
                                parent.remove(elem)

                    ws_tree.write(str(ws_path),
                                  encoding='UTF-8', xml_declaration=True)

                    shape_ids_to_remove = set(relid_to_shapeid.values())
                    self.log(f"    📐 shapeIds à supprimer: {shape_ids_to_remove}")

                    # ── 3c. Drawing XML → supprimer anchors par shapeId ───
                    drawing_path = None
                    for rel in ws_rels_root.findall(f'{{{ns_rel}}}Relationship'):
                        target = rel.get('Target', '')
                        if 'drawing' in target.lower():
                            if target.startswith('../'):
                                drawing_path = temp_path / 'xl' / target[3:]
                            else:
                                drawing_path = (temp_path / 'xl'
                                                / 'worksheets' / target)
                            if not drawing_path.exists():
                                drawing_path = (temp_path / 'xl' / 'drawings'
                                                / Path(target).name)
                            break

                    if drawing_path and drawing_path.exists() and shape_ids_to_remove:
                        drawing_tree = _et.parse(str(drawing_path))
                        drawing_root = drawing_tree.getroot()

                        for anchor_tag in [f'{{{ns_xdr}}}twoCellAnchor',
                                           f'{{{ns_xdr}}}oneCellAnchor',
                                           f'{{{ns_xdr}}}absoluteAnchor']:
                            for anchor in list(drawing_root.findall(anchor_tag)):
                                # Chercher un cNvPr dont id ∈ shape_ids_to_remove
                                shape_id_found = None
                                for cnv_pr in anchor.iter():
                                    tlocal = (cnv_pr.tag.split('}')[-1]
                                              if '}' in cnv_pr.tag else cnv_pr.tag)
                                    if tlocal == 'cNvPr':
                                        sp_id = cnv_pr.get('id')
                                        if sp_id and int(sp_id) in shape_ids_to_remove:
                                            shape_id_found = int(sp_id)
                                            break

                                if shape_id_found is None:
                                    continue

                                # Position cellule depuis <xdr:from>
                                col, row = 1, 1
                                from_elem = anchor.find(f'{{{ns_xdr}}}from')
                                if from_elem is not None:
                                    col_e = from_elem.find(f'{{{ns_xdr}}}col')
                                    row_e = from_elem.find(f'{{{ns_xdr}}}row')
                                    if col_e is not None:
                                        col = int(col_e.text) + 1
                                    if row_e is not None:
                                        row = int(row_e.text) + 1

                                # Retrouver l'item correspondant à ce shapeId
                                for r_id, sid in relid_to_shapeid.items():
                                    if sid != shape_id_found:
                                        continue
                                    embed_name = relid_to_embed.get(r_id, '')
                                    if embed_name not in items_by_embed:
                                        continue
                                    ext_name = items_by_embed[embed_name].get(
                                        'extracted_name', '')
                                    if ext_name:
                                        sname = (items_by_embed[embed_name]
                                                 .get('location', {})
                                                 .get('sheet_name',
                                                      f'Feuille{sheet_num}'))
                                        cells_to_fill.append(
                                            (sname, col, row,
                                             f"Voir {Path(ext_name).stem}"))
                                        self.log(f"    📍 {sname} col={col} "
                                                 f"row={row} → "
                                                 f"Voir {Path(ext_name).stem}")

                                drawing_root.remove(anchor)
                                self.log(f"    🗑️ Anchor supprimé du drawing "
                                         f"(shapeId={shape_id_found})")

                        drawing_tree.write(str(drawing_path),
                                           encoding='UTF-8', xml_declaration=True)

                    # ── 3d. Supprimer les rels embedding + image dans sheet ──
                    all_to_remove = rel_ids_to_remove | extra_relids_to_remove
                    for rel in list(ws_rels_root.findall(
                            f'{{{ns_rel}}}Relationship')):
                        if rel.get('Id') in all_to_remove:
                            ws_rels_root.remove(rel)

                    ws_rels_tree.write(str(ws_rels_path),
                                       encoding='UTF-8', xml_declaration=True)
                    self.log(f"    🗑️ {len(all_to_remove)} rel(s) supprimée(s) "
                             f"du sheet rels")

                # ── 4. Supprimer les fichiers embedding ──────────────────
                embeddings_dir = temp_path / 'xl' / 'embeddings'
                for item in extracted_files:
                    if item['action'] == 'replace' and item.get('filename'):
                        ep = embeddings_dir / item['filename']
                        if ep.exists():
                            ep.unlink()
                            self.log(f"    🗑️ Embedding supprimé: {item['filename']}")

                # ── 5. Repackager en XLSX ────────────────────────────────
                with zipfile.ZipFile(modified_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                    for fp in temp_path.rglob('*'):
                        if fp.is_file():
                            zout.write(fp, fp.relative_to(temp_path))

            # ── 6. Ajouter texte "Voir ..." via openpyxl ─────────────────
            if cells_to_fill:
                original_ext = original_path.suffix
                wb = openpyxl.load_workbook(str(modified_path),
                                            keep_vba=(original_ext == '.xlsm'))
                for sheet_name, col, row, text in cells_to_fill:
                    ws = (wb[sheet_name] if sheet_name in wb.sheetnames
                          else wb.active)
                    if ws is not None:
                        cell = ws.cell(row=row, column=col)
                        cell.value = text
                        cell.font = Font(bold=True, color='FF0000', size=12)
                        self.log(f"    ✅ Texte inséré: '{text}' → "
                                 f"{ws.title}!{cell.coordinate}")
                wb.save(str(modified_path))

            nb_replaced = sum(1 for i in extracted_files if i['action'] == 'replace')
            self.log(f"\n  ✅ Excel modifié : {modified_path.name}")
            self.log(f"  ✅ {nb_replaced} objet(s) OLE supprimé(s) et remplacé(s) par texte")

        except Exception as e:
            self.log(f"  ❌ Erreur: {str(e)}")
            import traceback
            self.log(f"  📋 Détails: {traceback.format_exc()}")
            try:
                modified_path = Path(output_dir) / Path(original_path).name
                shutil.copy2(original_path, modified_path)
                self.log(f"  ⚠️ Copie de l'original créée")
            except Exception:
                pass
    def create_modified_pptx_exact_positions(self, original_path, output_dir, extracted_files, temp_path):
            """
            Crée PowerPoint modifié en CONSERVANT LES POSITIONS EXACTES.
            
            RÈGLES :
            - Fichier EXTRAIT      → shape OLE supprimé + textbox rouge à la même position
            - Fichier NON SUPPORTÉ → shape OLE CONSERVÉ + textbox orange JUSTE À CÔTÉ (à droite)
            - Fichier VIDE         → shape OLE CONSERVÉ + textbox orange JUSTE AU-DESSUS
            
            CORRECTION : identification des shapes par shape.shape_id (cNvPr id)
            """
            try:
                from pptx import Presentation
                from pptx.util import Inches, Emu, Pt as PPTPt
                from pptx.dml.color import RGBColor as PPTRGBColor

                prs = Presentation(original_path)

                # Largeur de diapo (utile pour éviter de déborder)
                slide_width = prs.slide_width

                # Grouper les actions par slide
                slides_actions = {}
                for item in extracted_files:
                    slide_num = item['slide_number']
                    if slide_num not in slides_actions:
                        slides_actions[slide_num] = []
                    slides_actions[slide_num].append(item)

                # Collect VML spids of deleted OLE shapes for post-processing
                _vml_spids_to_remove = set()

                self.log(f"\n  📋 Plan de modification:")
                for slide_num in sorted(slides_actions.keys()):
                    self.log(f"    Slide {slide_num}: {len(slides_actions[slide_num])} modification(s)")

                for slide_num in sorted(slides_actions.keys()):
                    slide_idx = slide_num - 1
                    if slide_idx >= len(prs.slides):
                        self.log(f"    ⚠️ Slide {slide_num} introuvable")
                        continue

                    slide = prs.slides[slide_idx]
                    actions = slides_actions[slide_num]

                    self.log(f"\n  🔧 Slide {slide_num} ({len(slide.shapes)} shapes):")

                    # Trier par sp_element_id décroissant (évite décalages lors des suppressions)
                    actions.sort(key=lambda x: x.get('sp_element_id') or 0, reverse=True)

                    # Suivi des zones occupées par les textboxes déjà placées sur ce slide
                    # pour éviter le chevauchement (liste de (left, top, right, bottom))
                    _placed_boxes = []

                    def _find_free_top(left, top, w, h, gap=Inches(0.08)):
                        """Décale top vers le bas jusqu'à trouver une zone sans chevauchement."""
                        _top = top
                        for _ in range(50):  # max 50 tentatives
                            right  = left + w
                            bottom = _top + h
                            overlap = False
                            for (bl, bt, br, bb) in _placed_boxes:
                                if left < br and right > bl and _top < bb and bottom > bt:
                                    overlap = True
                                    _top = bb + gap
                                    break
                            if not overlap:
                                break
                        return _top

                    # Track names of added textboxes for reliable z-order identification
                    # (lxml proxy id() is not stable — we use shape name instead)
                    _added_tb_names = []
                    _tb_name_seq = [0]   # mutable counter accessible in nested scope

                    # ── Pre-grouping: detect 'replace' items at the same position ──
                    # (multiple OLE objects in the same table cell share similar left/top)
                    _MERGE_THR = int(914400 * 0.5)  # 0.5 inch in EMU
                    _sp_id_to_pos_key = {}
                    _pos_key_groups = {}
                    for _ai in actions:
                        if _ai.get('action') == 'replace':
                            _sid = _ai.get('sp_element_id')
                            _sh = None
                            if _sid is not None:
                                for _s in slide.shapes:
                                    if _s.shape_id == _sid:
                                        _sh = _s
                                        break
                            if _sh is not None:
                                _lb = round(_sh.left / _MERGE_THR)
                                _tb_val = round(_sh.top / _MERGE_THR)
                            else:
                                _lb = id(_ai)
                                _tb_val = 0
                            _key = (_lb, _tb_val)
                            _sp_id_to_pos_key[_sid] = _key
                            if _key not in _pos_key_groups:
                                _pos_key_groups[_key] = []
                            _pos_key_groups[_key].append(_ai)
                    _processed_merged_keys = set()

                    for action_item in actions:
                        action_type   = action_item['action']
                        sp_element_id = action_item.get('sp_element_id')
                        shape_index   = action_item.get('shape_index', 0)  # fallback

                        # ── Trouver le shape par shape_id ────────────────────────────
                        target_shape = None

                        if sp_element_id is not None:
                            for shape in slide.shapes:
                                if shape.shape_id == sp_element_id:
                                    target_shape = shape
                                    break

                        # Fallback par index si shape_id n'a pas fonctionné
                        if target_shape is None:
                            self.log(f"    ⚠️ shape_id={sp_element_id} non trouvé, fallback index {shape_index}")
                            current_idx = 0
                            for shape in slide.shapes:
                                current_idx += 1
                                if current_idx == shape_index:
                                    target_shape = shape
                                    break

                        if not target_shape:
                            self.log(f"    ❌ Shape introuvable (id={sp_element_id}, index={shape_index})")
                            continue

                        self.log(f"    ✅ Shape trouvé : id={target_shape.shape_id} name='{target_shape.name}'")

                        left   = target_shape.left
                        top    = target_shape.top
                        width  = target_shape.width
                        height = target_shape.height

                        # ────────────────────────────────────────────────────────────
                        # CAS 1 : FICHIER VIDE
                        # → shape conservé + texte orange AU-DESSUS
                        # ────────────────────────────────────────────────────────────
                        if action_type == 'keep_empty':
                            try:
                                w_w = width if width > 0 else Inches(1.8)
                                w_h = Inches(0.38)
                                warning_top = top - Inches(0.45)
                                if warning_top < 0:
                                    warning_top = top + height + Inches(0.1)
                                # Anti-chevauchement
                                warning_top = _find_free_top(left, warning_top, w_w, w_h)
                                _placed_boxes.append((left, warning_top, left + w_w, warning_top + w_h))

                                txBox = slide.shapes.add_textbox(left, warning_top, w_w, w_h)
                                txBox.fill.solid()
                                txBox.fill.fore_color.rgb = PPTRGBColor(50, 50, 50)
                                txBox.line.color.rgb = PPTRGBColor(200, 120, 0)
                                txBox.line.width = PPTPt(1.0)

                                tf = txBox.text_frame
                                tf.clear()
                                p = tf.paragraphs[0]
                                p.text = "⚠️ Fichier vide (0 Ko) — non archivé SharePoint"
                                p.font.bold = True
                                p.font.size = PPTPt(9)
                                p.font.color.rgb = PPTRGBColor(255, 200, 0)
                                tf.word_wrap = True

                                _tb_nm = f"_mlr_tb_{_tb_name_seq[0]}"
                                _tb_name_seq[0] += 1
                                txBox.name = _tb_nm
                                _added_tb_names.append(_tb_nm)

                                self.log(f"    ✅ Message VIDE ajouté au-dessus")
                            except Exception as e:
                                self.log(f"    ⚠️ Erreur message vide: {e}")

                        # ────────────────────────────────────────────────────────────
                        # CAS 2 : TYPE NON SUPPORTÉ (CATIA, etc.)
                        # → shape CONSERVÉ + texte orange JUSTE À CÔTÉ (à droite)
                        # ────────────────────────────────────────────────────────────
                        elif action_type == 'keep_unsupported':
                            try:
                                gap          = Inches(0.15)
                                label_left   = left + width + gap
                                label_width  = Inches(2.8)
                                label_height = max(height, Inches(0.4))

                                if label_left + label_width > slide_width:
                                    label_left = max(Emu(0), left - label_width - gap)

                                # Anti-chevauchement
                                label_top = _find_free_top(label_left, top, label_width, label_height)
                                _placed_boxes.append((label_left, label_top,
                                                      label_left + label_width, label_top + label_height))

                                txBox = slide.shapes.add_textbox(label_left, label_top, label_width, label_height)
                                txBox.fill.solid()
                                txBox.fill.fore_color.rgb = PPTRGBColor(50, 50, 50)
                                txBox.line.color.rgb = PPTRGBColor(200, 120, 0)
                                txBox.line.width = PPTPt(1.0)

                                tf = txBox.text_frame
                                tf.word_wrap = True
                                tf.clear()
                                p = tf.paragraphs[0]
                                p.text = "⚠️ Type de fichier non supporté par l'archivage SharePoint"
                                p.font.bold  = True
                                p.font.size  = PPTPt(9)
                                p.font.color.rgb = PPTRGBColor(255, 200, 0)

                                _tb_nm = f"_mlr_tb_{_tb_name_seq[0]}"
                                _tb_name_seq[0] += 1
                                txBox.name = _tb_nm
                                _added_tb_names.append(_tb_nm)

                                self.log(f"    ✅ Message NON SUPPORTÉ ajouté à droite du shape")
                            except Exception as e:
                                self.log(f"    ⚠️ Erreur message non supporté: {e}")

                        # ────────────────────────────────────────────────────────────
                        # CAS 3 : FICHIER EXTRAIT
                        # → shape OLE supprimé + textbox rouge à la même position
                        # ────────────────────────────────────────────────────────────
                        elif action_type == 'replace':
                            filename = action_item['filename']
                            _sid = action_item.get('sp_element_id')
                            _pos_key = _sp_id_to_pos_key.get(_sid)
                            _is_merged = _pos_key is not None and len(_pos_key_groups.get(_pos_key, [])) > 1
                            try:
                                sp = target_shape.element
                                # Extract VML spid BEFORE deletion (mc:Choice path)
                                _vml_spid = None
                                if sp is not None:
                                    # Method 1: iterate XML tree
                                    for _sub in sp.iter():
                                        _stag = _sub.tag
                                        if _stag.endswith('}oleObj') or _stag == 'oleObj':
                                            _vml_spid = _sub.get('spid')
                                            break
                                    # Method 2: fallback — scan raw XML bytes for spid="..."
                                    if not _vml_spid:
                                        try:
                                            import re as _re_spid
                                            from lxml import etree as _lxml_sp
                                            _raw = _lxml_sp.tostring(sp, encoding='unicode')
                                            _m = _re_spid.search(r'\bspid\s*=\s*["\']([^"\']+)["\']', _raw)
                                            if _m:
                                                _vml_spid = _m.group(1)
                                        except Exception:
                                            pass
                                    if _vml_spid:
                                        _vml_spids_to_remove.add(_vml_spid)
                                        self.log(f"    🔍 spid VML capturé: {_vml_spid}")
                                    else:
                                        self.log(f"    ⚠️ Aucun spid trouvé dans le shape OLE")
                                if sp is not None and sp.getparent() is not None:
                                    sp.getparent().remove(sp)
                                    self.log(f"    ✅ Shape OLE supprimé"
                                             + (f" (spid={_vml_spid})" if _vml_spid else ""))

                                # ── Supprimer TOUTES les formes OLE/icône à la même position ──
                                # Critères d'identification (du plus fiable au moins fiable) :
                                #   1. shape_type PICTURE / OLE / LINKED_PICTURE
                                #   2. tag XML == <p:pic> ou <p:graphicFrame>
                                #   3. <p:sp> avec un <blipFill> (image de fond)
                                # On préserve : nos textboxes (_mlr_), formes avec texte non vide.
                                _margin_ico = int(914400 * 0.3)   # 0.3 inch
                                _ole_r = left + max(width,  Inches(0.1))
                                _ole_b = top  + max(height, Inches(0.1))
                                try:
                                    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
                                    _DELETABLE_TYPES = {
                                        _MSO.PICTURE,
                                        _MSO.EMBEDDED_OLE_OBJECT,
                                        _MSO.LINKED_OLE_OBJECT,
                                        _MSO.LINKED_PICTURE,
                                    }
                                except Exception:
                                    _DELETABLE_TYPES = {13, 7, 10, 11}

                                for _ico_s in list(slide.shapes):
                                    try:
                                        if _ico_s.shape_id == (_sid or -1):
                                            continue
                                        _ico_nm = getattr(_ico_s, 'name', '')
                                        if _ico_nm.startswith('_mlr_'):
                                            continue
                                        # Filtre position (centre dans la zone OLE)
                                        _cx = _ico_s.left + _ico_s.width  / 2
                                        _cy = _ico_s.top  + _ico_s.height / 2
                                        if not ((left - _margin_ico) <= _cx <= (_ole_r + _margin_ico) and
                                                (top  - _margin_ico) <= _cy <= (_ole_b + _margin_ico)):
                                            continue
                                        # Préserver les formes avec du texte visible
                                        try:
                                            if (_ico_s.has_text_frame and
                                                    _ico_s.text_frame.text.strip()):
                                                continue
                                        except Exception:
                                            pass
                                        # Décider si on supprime
                                        _del = False
                                        # Critère 1 : shape_type
                                        try:
                                            if _ico_s.shape_type in _DELETABLE_TYPES:
                                                _del = True
                                        except Exception:
                                            pass
                                        # Critère 2 : tag XML
                                        if not _del:
                                            _xtag = _ico_s.element.tag
                                            if (_xtag.endswith('}pic')
                                                    or _xtag == 'pic'
                                                    or _xtag.endswith('}graphicFrame')
                                                    or _xtag == 'graphicFrame'):
                                                _del = True
                                        # Critère 3 : blipFill dans le XML
                                        if not _del:
                                            for _sub in _ico_s.element.iter():
                                                _st = _sub.tag
                                                if (_st.endswith('}blipFill')
                                                        or _st == 'blipFill'):
                                                    _del = True
                                                    break
                                        if _del:
                                            _ico_s.element.getparent().remove(
                                                _ico_s.element)
                                            self.log(
                                                f"    ✅ Élément OLE/icône supprimé: {_ico_nm}")
                                    except Exception:
                                        pass

                                # Si ce groupe de table a déjà eu son textbox fusionné, ne pas en créer un autre
                                if _is_merged and _pos_key in _processed_merged_keys:
                                    self.log(f"    ℹ️ Groupe table fusionné: textbox déjà créé")
                                    continue

                                slide_h = prs.slide_height
                                slide_w = prs.slide_width

                                # Noms à afficher : un seul fichier, ou tous les fichiers du groupe
                                if _is_merged:
                                    _group_filenames = [gi['filename'] for gi in _pos_key_groups[_pos_key]]
                                    _processed_merged_keys.add(_pos_key)
                                else:
                                    _group_filenames = [filename]

                                # Hauteur adaptée au nombre de fichiers dans le groupe
                                tb_width  = Inches(1.8)
                                tb_height = Inches(0.35) + Inches(0.28) * len(_group_filenames)

                                # Position de base : là où était le shape
                                clamped_left = left
                                clamped_top  = top

                                # ── Clamper LEFT ─────────────────────────────────────────
                                if clamped_left < 0:
                                    clamped_left = Inches(0.1)
                                if clamped_left + tb_width > slide_w:
                                    clamped_left = slide_w - tb_width - Inches(0.1)
                                if clamped_left < 0:
                                    clamped_left = Inches(0.1)

                                # ── Clamper TOP ──────────────────────────────────────────
                                if clamped_top < 0:
                                    clamped_top = Inches(0.1)
                                if clamped_top + tb_height > slide_h:
                                    clamped_top = slide_h - tb_height - Inches(0.1)
                                if clamped_top < 0:
                                    clamped_top = Inches(0.1)

                                # ── Anti-chevauchement : décaler vers le bas si nécessaire ─
                                clamped_top = _find_free_top(clamped_left, clamped_top, tb_width, tb_height)
                                # Clamper après décalage
                                if clamped_top + tb_height > slide_h:
                                    clamped_top = slide_h - tb_height - Inches(0.1)
                                if clamped_top < 0:
                                    clamped_top = Inches(0.1)
                                _placed_boxes.append((clamped_left, clamped_top,
                                                      clamped_left + tb_width, clamped_top + tb_height))

                                # ── Taille du texte adaptée selon l'espace disponible ───
                                font_size = PPTPt(10) if (width < Inches(1.5) or height < Inches(0.5)) else PPTPt(12)

                                txBox = slide.shapes.add_textbox(
                                    clamped_left, clamped_top, tb_width, tb_height
                                )

                                tf = txBox.text_frame
                                tf.word_wrap = True
                                tf.clear()

                                # Une ligne par fichier (retour à la ligne si groupe de table)
                                for _fn_idx, _fn in enumerate(_group_filenames):
                                    if _fn_idx == 0:
                                        p = tf.paragraphs[0]
                                    else:
                                        p = tf.add_paragraph()
                                    p.text = f"Voir {Path(_fn).stem}"
                                    p.font.bold      = True
                                    p.font.size      = font_size
                                    p.font.color.rgb = PPTRGBColor(255, 0, 0)

                                _tb_nm = f"_mlr_tb_{_tb_name_seq[0]}"
                                _tb_name_seq[0] += 1
                                txBox.name = _tb_nm
                                _added_tb_names.append(_tb_nm)

                                self.log(f"    ✅ Shape REMPLACÉ → {len(_group_filenames)} fichier(s) "
                                        f"(left={clamped_left//914400:.2f}\", top={clamped_top//914400:.2f}\")")
                            except Exception as e:
                                self.log(f"    ⚠️ Erreur remplacement: {e}")

                    # ── Passe finale z-order : reconstruire spTree avec nos textboxes EN DERNIER ──
                    # Utilise slide.shapes._spTree directement + identification par nom (pas par id() lxml)
                    if _added_tb_names:
                        try:
                            _spt = slide.shapes._spTree
                            _tb_name_set = set(_added_tb_names)

                            def _cNvPr_name(elem):
                                for _el in elem.iter():
                                    tag = _el.tag
                                    if tag.endswith('}cNvPr') or tag == 'cNvPr':
                                        return _el.get('name', '')
                                return ''

                            _all_ch = list(_spt)
                            _non_tb_ch = [e for e in _all_ch if _cNvPr_name(e) not in _tb_name_set]
                            _tb_ch     = [e for e in _all_ch if _cNvPr_name(e) in _tb_name_set]
                            if _tb_ch:
                                for _e in _all_ch:
                                    _spt.remove(_e)
                                for _e in _non_tb_ch + _tb_ch:
                                    _spt.append(_e)
                                self.log(f"    ✅ Z-order: {len(_tb_ch)} textbox(es) placés en avant-plan absolu")
                        except Exception as _ze:
                            self.log(f"    ⚠️ Passe z-order: {_ze}")

#                # ── Slide récapitulative (insérée en 1ère position) ──────────────────
#                if extracted_files:
#                    try:
#                        self.log(f"\n  📊 Ajout slide récapitulative...")
#                        slide_layout = (
#                            prs.slide_layouts[5]
#                            if len(prs.slide_layouts) > 5
#                            else prs.slide_layouts[0]
#                        )
#                        recap_slide = prs.slides.add_slide(slide_layout)
#
#                        xml_slides  = prs.slides._sldIdLst
#                        slides_list = list(xml_slides)
#                        xml_slides.remove(slides_list[-1])
#                        xml_slides.insert(0, slides_list[-1])
#
#                        if recap_slide.shapes.title:
#                            title_para = recap_slide.shapes.title.text_frame.paragraphs[0]
#                            title_para.text = "FICHIERS INCORPORÉS — RÉCAPITULATIF"
#                            title_para.font.bold  = True
#                            title_para.font.size  = PPTPt(26)
#                            title_para.font.color.rgb = PPTRGBColor(255, 0, 0)
#
#                        txBox = recap_slide.shapes.add_textbox(
#                            Inches(0.5), Inches(1.5), Inches(9), Inches(5)
#                        )
#                        tf = txBox.text_frame
#                        tf.clear()
#
#                        for slide_num in sorted(slides_actions.keys()):
#                            p = tf.add_paragraph()
#                            p.text = f"Slide {slide_num} :"
#                            p.font.bold = True
#                            p.font.size = PPTPt(13)
#
#                            for item in slides_actions[slide_num]:
#                                p2 = tf.add_paragraph()
#                                p2.font.size = PPTPt(11)
#                                if item['action'] == 'keep_empty':
#                                    p2.text = "  ⚠️ Fichier vide (0 Ko) — non archivé SharePoint"
#                                    p2.font.color.rgb = PPTRGBColor(255, 153, 0)
#                                elif item['action'] == 'keep_unsupported':
#                                    p2.text = "  ⚠️ Type de fichier non supporté par l'archivage SharePoint"
#                                    p2.font.color.rgb = PPTRGBColor(255, 153, 0)
#                                else:
#                                    p2.text = f"  ✅ Extrait → Voir {item['filename']}"
#                                    p2.font.color.rgb = PPTRGBColor(0, 176, 80)
#
#                        self.log(f"    ✅ Slide récapitulative ajoutée en position 1")
#                    except Exception as e:
#                        self.log(f"    ⚠️ Erreur slide récap: {e}")
#
                # ── Sauvegarde ───────────────────────────────────────────────────────
                modified_path = Path(output_dir) / Path(original_path).name
                prs.save(str(modified_path))

                # Post-process: remove VML OLE icon shapes from vmlDrawingN.vml files
                if _vml_spids_to_remove:
                    self._remove_vml_ole_shapes(modified_path, _vml_spids_to_remove)

                self.log(f"\n  ✅ PowerPoint modifié : {modified_path.name}")
                self.log(f"  ✅ {len(extracted_files)} modification(s) appliquée(s)")

            except Exception as e:
                self.log(f"  ❌ Erreur: {str(e)}")
                import traceback
                self.log(f"  📋 Détails: {traceback.format_exc()}")
                try:
                    modified_path = Path(output_dir) / Path(original_path).name
                    shutil.copy2(original_path, modified_path)
                    self.log(f"  ⚠️ Copie de l'original créée (fallback)")
                except Exception:
                    pass

    def _remove_vml_ole_shapes(self, pptx_path, vml_spids):
        """Post-process saved PPTX ZIP to remove VML OLE icon shapes from vmlDrawingN.vml.

        When an OLE object uses mc:Choice/VML, the icon lives in a separate .vml file
        inside the ZIP — python-pptx cannot touch it.  We rewrite the ZIP here.
        """
        import zipfile
        import re as _re

        if not vml_spids:
            return

        pptx_path = Path(pptx_path)
        tmp_path  = pptx_path.with_suffix('.tmp_vml_fix')
        total_removed = 0

        try:
            with zipfile.ZipFile(pptx_path, 'r') as zin:
                all_names = zin.namelist()
                vml_names = [n for n in all_names if n.lower().endswith('.vml')]
                self.log(f"    🔍 Fichiers VML dans ZIP: {vml_names if vml_names else 'aucun'}")
                self.log(f"    🔍 Spids à supprimer: {sorted(vml_spids)}")

                with zipfile.ZipFile(tmp_path, 'w') as zout:
                    for item in zin.infolist():
                        data = zin.read(item.filename)

                        if item.filename.lower().endswith('.vml'):
                            try:
                                text = data.decode('utf-8', errors='replace')
                                # Log a preview so we can see the content
                                preview = text.replace('\n', ' ').replace('\r', '')[:400]
                                self.log(f"    📄 {item.filename}: {preview}")

                                orig_text = text
                                for _spid in vml_spids:
                                    # Pattern 1: <v:shape id="SPID" ...>...</v:shape>
                                    _pat = _re.compile(
                                        r'<v:shape\b[^>]*\bid\s*=\s*["\']'
                                        + _re.escape(_spid)
                                        + r'["\'][^>]*>.*?</v:shape\s*>',
                                        _re.DOTALL | _re.IGNORECASE,
                                    )
                                    text, _n = _pat.subn('', text)
                                    if _n:
                                        total_removed += _n
                                        self.log(f"    ✅ VML shape supprimé (regex): {_spid}")
                                    else:
                                        # Pattern 2: self-closing <v:shape id="SPID" .../>
                                        _pat2 = _re.compile(
                                            r'<v:shape\b[^>]*\bid\s*=\s*["\']'
                                            + _re.escape(_spid)
                                            + r'["\'][^>]*/\s*>',
                                            _re.DOTALL | _re.IGNORECASE,
                                        )
                                        text, _n2 = _pat2.subn('', text)
                                        if _n2:
                                            total_removed += _n2
                                            self.log(f"    ✅ VML shape (self-closing) supprimé: {_spid}")
                                        else:
                                            self.log(f"    ⚠️ Spid '{_spid}' non trouvé dans {item.filename}")

                                if text != orig_text:
                                    data = text.encode('utf-8')

                            except Exception as _ve:
                                self.log(f"    ⚠️ Erreur traitement VML {item.filename}: {_ve}")

                        # Preserve original compression for each entry
                        info_out = zipfile.ZipInfo(item.filename)
                        info_out.compress_type = item.compress_type
                        info_out.date_time     = item.date_time
                        zout.writestr(info_out, data)

            if total_removed > 0:
                try:
                    pptx_path.unlink()
                except Exception:
                    pass
                try:
                    tmp_path.rename(pptx_path)
                    self.log(f"    ✅ PPTX ZIP recréé — {total_removed} icône(s) VML supprimée(s)")
                except Exception as _rn:
                    import shutil as _sh
                    _sh.move(str(tmp_path), str(pptx_path))
                    self.log(f"    ✅ PPTX déplacé (shutil) — {total_removed} icône(s) VML supprimée(s)")
            else:
                try:
                    tmp_path.unlink()
                except Exception:
                    pass
                self.log(f"    ℹ️ Aucune icône VML supprimée")

        except Exception as _e:
            self.log(f"    ⚠️ Erreur nettoyage VML: {_e}")
            import traceback as _tb
            self.log(f"    📋 {_tb.format_exc()}")
            try:
                tmp_path.unlink()
            except Exception:
                pass

    def _create_pptx_textbox(self, x, y, cx, cy, text, color, font_size):
        """Crée un élément textbox pour PowerPoint (XML lxml)"""
        from lxml import etree
        
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        
        # Créer la shape
        sp = etree.Element(f"{{{ns_p}}}sp")
        
        # nvSpPr (Non-Visual Shape Properties)
        nvSpPr = etree.SubElement(sp, f"{{{ns_p}}}nvSpPr")
        
        cNvPr = etree.SubElement(nvSpPr, f"{{{ns_p}}}cNvPr")
        cNvPr.set('id', str(9999))
        cNvPr.set('name', 'TextBox')
        
        cNvSpPr = etree.SubElement(nvSpPr, f"{{{ns_p}}}cNvSpPr")
        txBox = etree.SubElement(cNvSpPr, f"{{{ns_a}}}txBox")
        
        nvPr = etree.SubElement(nvSpPr, f"{{{ns_p}}}nvPr")
        
        # spPr (Shape Properties)
        spPr = etree.SubElement(sp, f"{{{ns_p}}}spPr")
        
        xfrm = etree.SubElement(spPr, f"{{{ns_a}}}xfrm")
        
        off = etree.SubElement(xfrm, f"{{{ns_a}}}off")
        off.set('x', str(x))
        off.set('y', str(y))
        
        ext = etree.SubElement(xfrm, f"{{{ns_a}}}ext")
        ext.set('cx', str(cx))
        ext.set('cy', str(cy))
        
        prstGeom = etree.SubElement(spPr, f"{{{ns_a}}}prstGeom")
        prstGeom.set('prst', 'rect')
        avLst = etree.SubElement(prstGeom, f"{{{ns_a}}}avLst")
        
        # txBody (Text Body)
        txBody = etree.SubElement(sp, f"{{{ns_p}}}txBody")
        
        bodyPr = etree.SubElement(txBody, f"{{{ns_a}}}bodyPr")
        bodyPr.set('wrap', 'square')
        bodyPr.set('rtlCol', '0')
        
        lstStyle = etree.SubElement(txBody, f"{{{ns_a}}}lstStyle")
        
        # Paragraphe
        p = etree.SubElement(txBody, f"{{{ns_a}}}p")
        
        # Run
        r = etree.SubElement(p, f"{{{ns_a}}}r")
        
        # Run Properties
        rPr = etree.SubElement(r, f"{{{ns_a}}}rPr")
        rPr.set('b', '1')  # Bold
        rPr.set('sz', str(font_size * 100))  # Font size en centièmes de point
        rPr.set('lang', 'fr-FR')
        
        # Couleur
        solidFill = etree.SubElement(rPr, f"{{{ns_a}}}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{{{ns_a}}}srgbClr")
        srgbClr.set('val', color)
        
        # Texte
        t = etree.SubElement(r, f"{{{ns_a}}}t")
        t.text = text
        
        return sp


    def _add_pptx_summary_slide_simple(self, pptx_temp_path, position_actions, ns):
        """
        Ajoute une slide récapitulative simple au début du PowerPoint
        """
        from lxml import etree
        
        try:
            # Lire le fichier presentation.xml pour obtenir la liste des slides
            presentation_file = pptx_temp_path / "ppt" / "presentation.xml"
            pres_tree = etree.parse(str(presentation_file))
            pres_root = pres_tree.getroot()
            
            # Trouver sldIdLst
            sldIdLst = pres_root.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
            
            if sldIdLst is not None:
                # Trouver le prochain ID disponible
                max_id = 256
                for sldId in sldIdLst:
                    current_id = int(sldId.get('id', '256'))
                    if current_id >= max_id:
                        max_id = current_id + 1
                
                # Créer une nouvelle slide simple avec titre et texte
                slides_dir = pptx_temp_path / "ppt" / "slides"
                
                # Trouver le prochain numéro de slide
                existing_slides = list(slides_dir.glob("slide*.xml"))
                slide_numbers = [int(re.search(r'slide(\d+)', s.name).group(1)) for s in existing_slides]
                next_slide_num = max(slide_numbers) + 1 if slide_numbers else 1
                
                new_slide_file = slides_dir / f"slide{next_slide_num}.xml"
                
                # Créer le XML de la nouvelle slide
                new_slide_xml = self._create_summary_slide_xml(position_actions)
                
                # Écrire la nouvelle slide
                with open(new_slide_file, 'wb') as f:
                    f.write(new_slide_xml)
                
                # Ajouter la référence dans presentation.xml
                new_sldId = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}sldId')
                new_sldId.set('id', str(max_id))
                new_sldId.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', f'rId{max_id}')
                
                # Insérer au début
                sldIdLst.insert(0, new_sldId)
                
                # Sauvegarder presentation.xml
                pres_tree.write(str(presentation_file), encoding='utf-8', xml_declaration=True)
                
                self.log(f"    ✅ Slide récapitulative ajoutée en première position")
        
        except Exception as e:
            self.log(f"    ⚠️ Impossible d'ajouter slide récapitulative: {e}")


    def _create_summary_slide_xml(self, position_actions):
        """
        Crée le XML pour une slide récapitulative simple
        """
        from lxml import etree
        
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        
        # Template XML minimal pour une slide
        slide_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
    <p:sld xmlns:p="{ns_p}" xmlns:a="{ns_a}">
        <p:cSld>
            <p:spTree>
                <p:nvGrpSpPr>
                    <p:cNvPr id="1" name=""/>
                    <p:cNvGrpSpPr/>
                    <p:nvPr/>
                </p:nvGrpSpPr>
                <p:grpSpPr>
                    <a:xfrm>
                        <a:off x="0" y="0"/>
                        <a:ext cx="0" cy="0"/>
                        <a:chOff x="0" y="0"/>
                        <a:chExt cx="0" cy="0"/>
                    </a:xfrm>
                </p:grpSpPr>
            </p:spTree>
        </p:cSld>
        <p:clrMapOvr>
            <a:masterClrMapping/>
        </p:clrMapOvr>
    </p:sld>'''
        
        root = etree.fromstring(slide_xml.encode('utf-8'))
        spTree = root.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree')
        
        # Ajouter titre
        title_shape = self._create_pptx_textbox(914400, 457200, 7315200, 914400, 
                                            "FICHIERS EXTRAITS - RÉCAPITULATIF", "FF0000", 28)
        spTree.append(title_shape)
        
        # Ajouter contenu
        y_pos = 1828800
        for position in sorted(position_actions.keys()):
            action, filename = position_actions[position]
            
            if action == 'keep_with_warning_empty':
                text = f"⚠️ Position {position}: Fichier vide (0 Ko) - non archivé"
                color = "FF9900"
            elif action == 'keep_with_warning_unsupported':
                text = f"⚠️ Position {position}: Type non supporté - reste dans document"
                color = "FF9900"
            else:
                text = f"✅ Position {position}: Extrait → Voir {filename}"
                color = "00B050"
            
            content_shape = self._create_pptx_textbox(914400, y_pos, 7315200, 457200, text, color, 14)
            spTree.append(content_shape)
            
            y_pos += 548640
        
        return etree.tostring(root, encoding='utf-8', xml_declaration=True)
    def extract_from_pdf(self, pdf_path, output_dir):
        """
        Extrait pièces jointes PDF avec numéro de page exact
        """
        self.reset_extraction_tracking(Path(pdf_path).name, output_dir)
        Path(output_dir).mkdir(parents=True, exist_ok=True)

        extracted_files = []

        try:
            pdf_document = fitz.open(pdf_path)
            
            self.log(f"\n  📖 Document: {len(pdf_document)} page(s)")
            
            # ========================================
            # PARCOURIR TOUTES LES PAGES
            # ========================================
            position = 0
            
            for page_num in range(len(pdf_document)):
                page = pdf_document[page_num]
                page_number = page_num + 1
                
                # Chercher annotations
                for annot in page.annots():
                    if annot.type[0] == fitz.PDF_ANNOT_FILE_ATTACHMENT:
                        try:
                            position += 1
                            
                            self.log(f"\n  🔍 Position {position} (Page {page_number})")
                            
                            # Extraire contenu
                            content = annot.get_file()
                            
                            if len(content) == 0:
                                self.log(f"    ⚠️ Pièce jointe VIDE")
                                
                                extracted_files.append({
                                    'position': position,
                                    'page': page_number,
                                    'action': 'keep_empty',
                                    'annot_rect': annot.rect
                                })
                                
                                self.add_to_report('extracted', {
                                    'source_file': Path(pdf_path).name,
                                    'extracted_file': '[Vide - 0 Ko]',
                                    'position': f'Page {page_number}',
                                    'type': 'Vide',
                                    'status': 'Ignoré'
                                })
                                continue
                            
                            # Détecter type
                            file_ext = self.detect_file_type(content)
                            
                            # Vérifier si image
                            if file_ext in ['.jpg', '.png', '.gif', '.bmp', '.tiff']:
                                self.log(f"    ⏭️ Image ignorée")
                                continue
                            
                            # Vérifier si supporté
                            SUPPORTED = {'.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls', '.pptx', '.ppt', '.msg', '.txt', '.zip', '.7z','.htm', '.html' }
                            
                            if file_ext not in SUPPORTED and file_ext != '.ole':
                                # NON SUPPORTÉ
                                self.log(f"    ⚠️ Type non supporté: {file_ext}")
                                
                                extracted_files.append({
                                    'position': position,
                                    'page': page_number,
                                    'action': 'keep_unsupported',
                                    'file_type': file_ext,
                                    'annot_rect': annot.rect
                                })
                                
                                self.add_to_report('extracted', {
                                    'source_file': Path(pdf_path).name,
                                    'extracted_file': f'[Non supporté: {file_ext}]',
                                    'position': f'Page {page_number}',
                                    'type': file_ext,
                                    'status': 'Conservé'
                                })
                                continue
                            
                            # SUPPORTÉ → EXTRAIRE
                            self.extracted_count += 1
                            output_name = self._generate_output_name(
                                Path(pdf_path).stem, None, file_ext
                            )
                            output_path = Path(output_dir) / output_name
                            
                            with open(output_path, 'wb') as f:
                                f.write(content)
                            
                            self.log(f"    ✅ Extrait: {output_name}")
                            
                            extracted_files.append({
                                'position': position,
                                'page': page_number,
                                'action': 'replace',
                                'extracted_name': output_name,
                                'annot_rect': annot.rect
                            })
                            
                            self.add_to_report('extracted', {
                                'source_file': Path(pdf_path).name,
                                'extracted_file': output_name,
                                'position': f'Page {page_number}',
                                'type': file_ext,
                                'status': 'Succès'
                            })
                            
                        except Exception as annot_error:
                            self.log(f"    ⚠️ Erreur pièce jointe: {annot_error}")
            
            pdf_document.close()
            
            # ========================================
            # CRÉER PDF MODIFIÉ
            # ========================================
            if extracted_files:
                self.log(f"\n  🔧 CRÉATION DU PDF MODIFIÉ...")
                self.create_modified_pdf_exact_positions(
                    pdf_path, output_dir, extracted_files
                )
            else:
                original_copy = Path(output_dir) / Path(pdf_path).name
                if Path(pdf_path).resolve() != original_copy.resolve():
                    try:
                        shutil.copy2(pdf_path, original_copy)
                    except PermissionError:
                        time.sleep(2)
                        shutil.copy2(pdf_path, original_copy)
                self.log(f"  ✅ Copie créée: {original_copy.name}")
            
            self.add_to_report('processed', {
                'file_name': Path(pdf_path).name,
                'file_type': 'PDF',
                'files_found': len(extracted_files),
                'status': 'Traité avec succès'
            })
        
        except Exception as e:
            self.log(f"  ❌ Erreur: {str(e)}")
            import traceback
            self.log(f"  📋 Détails: {traceback.format_exc()}")
            # Copie de secours si dossier de sortie vide
            try:
                original_copy = Path(output_dir) / Path(pdf_path).name
                if not original_copy.exists():
                    shutil.copy2(pdf_path, original_copy)
                    self.log(f"  ⚠️ Erreur → copie de l'original créée: {original_copy.name}")
            except Exception:
                pass

        return extracted_files


    def create_modified_pdf_exact_positions(self, original_path, output_dir, extracted_files):
        """
        Crée PDF modifié :
        - Insère uniquement le nom du fichier extrait en texte simple à côté de l'icône
        - Pas de page récapitulative
        - Pas de fond coloré / rectangle
        - Sauvegarde garantie
        """
        try:
            pdf_document = fitz.open(original_path)

            self.log(f"\n  📝 Ajout des noms de fichiers sur les pages...")

            for item in extracted_files:
                page_num   = item['page'] - 1
                action     = item['action']
                annot_rect = item.get('annot_rect')

                if page_num >= len(pdf_document):
                    continue

                page = pdf_document[page_num]

                # Position : juste à droite de l'icône pièce jointe
                if annot_rect:
                    x = annot_rect.x1 + 5
                    y = annot_rect.y0 + 10
                else:
                    x = page.rect.width - 220
                    y = 30

                # Texte et couleur selon le cas
                if action == 'keep_empty':
                    text  = "[Fichier vide - 0 Ko]"
                    color = (1.0, 0.5, 0.0)   # orange

                elif action == 'keep_unsupported':
                    ft    = item.get('file_type', '?')
                    text  = f"[Type non supporte : {ft}]"
                    color = (1.0, 0.5, 0.0)   # orange

                else:
                    extracted_name = item.get('extracted_name', '')
                    text  = Path(extracted_name).stem   # juste le nom, sans extension
                    color = (0.85, 0.0, 0.0)            # rouge

                try:
                    page.insert_text(
                        fitz.Point(x, y),
                        text,
                        fontsize=9,
                        color=color,
                        fontname="helv"
                    )
                    self.log(f"    ✅ Page {item['page']}: '{text}'")

                except Exception as err:
                    self.log(f"    ⚠️ Erreur texte page {item['page']}: {err}")

            # Sauvegarde — toujours exécutée
            modified_path = Path(output_dir) / Path(original_path).name
            pdf_document.save(str(modified_path))
            pdf_document.close()

            self.log(f"\n  ✅ PDF modifié sauvegardé : {modified_path.name}")

        except Exception as e:
            self.log(f"  ❌ Erreur : {e}")
            import traceback
            self.log(f"  📋 {traceback.format_exc()}")
            # Fallback : copier l'original
            try:
                fallback = Path(output_dir) / Path(original_path).name
                if not fallback.exists():
                    shutil.copy2(original_path, fallback)
                    self.log(f"  ⚠️ Copie originale créée : {fallback.name}")
            except Exception:
                pass
    def extract_from_msg_file(self, msg_path, output_dir):
        """
        Traite un fichier MSG :
        1. Crée le dossier de sortie si inexistant
        2. Extrait les pièces jointes (attachments) avec nomenclature FJ_X
        3. Convertit MSG → PDF en insérant "Voir document : nom" dans le PDF
        """
        self.log(f"\n📧 Traitement MSG: {Path(msg_path).name}")

        # ── Garantir que le dossier de sortie existe ──────────────────────
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        self.reset_extraction_tracking(Path(msg_path).name, output_dir)

        extracted_names = []

        try:
            # ================================================================
            # ÉTAPE 1 : EXTRAIRE LES PIÈCES JOINTES
            # ================================================================
            try:
                import win32com.client
                import pythoncom

                pythoncom.CoInitialize()

                try:
                    SUPPORTED_EXTENSIONS = {
                        '.pdf', '.docx', '.doc', '.xlsx', '.xlsm', '.xls',
                        '.pptx', '.ppt', '.msg', '.txt', '.zip', '.7z','.pptm','.htm', '.html' 
                    }

                    outlook = win32com.client.Dispatch("Outlook.Application")
                    msg_com = outlook.Session.OpenSharedItem(str(Path(msg_path).absolute()))

                    attachments = msg_com.Attachments
                    nb_att = attachments.Count
                    self.log(f"  📎 {nb_att} pièce(s) jointe(s) détectée(s)")

                    for i in range(1, nb_att + 1):
                        attachment = attachments.Item(i)
                        att_name   = attachment.FileName

                        if not att_name:
                            self.log(f"    ⚠️ Pièce jointe {i} sans nom, ignorée")
                            continue

                        att_ext = Path(att_name).suffix.lower()

                        if att_ext not in SUPPORTED_EXTENSIONS:
                            # Images from MSG are silently ignored — not added to Non Archivables
                            _IMG_EXTS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff',
                                         '.webp', '.svg', '.ico', '.emf', '.wmf'}
                            if att_ext in _IMG_EXTS:
                                self.log(f"    ⏭️ Image ignorée (MSG): {att_name}")
                                continue
                            self.log(f"    ⚠️ Extension non supportée : {att_name}")
                            self.add_to_report('extracted', {
                                'source_file': Path(msg_path).name,
                                'extracted_file': f'[Non supporté: {att_name}]',
                                'position': f'Pièce jointe {i}',
                                'type': att_ext,
                                'status': 'Conservé — type non supporté'
                            })
                            continue

                        try:
                            # ── Sauvegarde temporaire pour vérif MD5 avant nommage ─
                            import tempfile, hashlib, shutil
                            tmp = tempfile.mktemp(suffix=att_ext)
                            attachment.SaveAsFile(tmp)

                            with open(tmp, 'rb') as f:
                                att_bytes = f.read()

                            _att_hash = hashlib.md5(att_bytes).hexdigest()
                            if _att_hash in self._extracted_content_hashes:
                                self.log(f"    ⏭️ PJ déjà extraite ailleurs (MD5={_att_hash[:8]}...) : {att_name}")
                                os.remove(tmp)
                                continue

                            self._extracted_content_hashes.add(_att_hash)

                            # ── Nommage FJ_X APRÈS vérification doublon ───────
                            self.extracted_count += 1
                            output_name = self._generate_output_name(
                                Path(msg_path).stem, att_name, att_ext
                            )
                            output_name_clean = output_name.strip()
                            output_path = output_dir / output_name_clean

                            # ── Copie vers destination finale ─────────────────
                            shutil.copy2(tmp, str(output_path.absolute()))
                            os.remove(tmp)

                            self.log(f"    ✅ PJ extraite : {output_name_clean}")
                            extracted_names.append(output_name_clean)
                            self.add_to_report('extracted', {
                                'source_file': Path(msg_path).name,
                                'extracted_file': output_name_clean,
                                'position': f'Pièce jointe {i}',
                                'type': att_ext,
                                'status': 'Succès'
                            })
                        except Exception as save_err:
                            self.log(f"    ❌ Erreur sauvegarde PJ : {save_err}")

                    msg_com.Close(0)

                finally:
                    pythoncom.CoUninitialize()

            except ImportError:
                self.log(f"  ⚠️ Module pywin32 requis")
            except Exception as att_error:
                self.log(f"  ⚠️ Erreur extraction pièces jointes : {att_error}")

            # ================================================================
            # ÉTAPE 2 : CONVERSION MSG → PDF
            # ================================================================
            try:
                self.log(f"  🔄 Conversion MSG → PDF...")
                pdf_path = self.convert_msg_to_pdf(msg_path, output_dir, extracted_names)

                if pdf_path and Path(pdf_path).exists():
                    self.log(f"  ✓ PDF créé: {Path(pdf_path).name}")
                    self._processed_absolute_paths.add(str(Path(pdf_path).resolve()))
                else:
                    self.log(f"  ⚠ Conversion PDF échouée (Outlook requis)")
                # Toujours supprimer le MSG du dossier de sortie (avec ou sans PDF)
                msg_in_output = Path(output_dir) / Path(msg_path).name
                if msg_in_output.exists() and msg_in_output.resolve() != Path(msg_path).resolve():
                    try:
                        msg_in_output.unlink()
                        self.log(f"  🗑️ MSG supprimé du dossier de sortie: {msg_in_output.name}")
                    except Exception as del_err:
                        self.log(f"  ⚠️ Impossible de supprimer le MSG: {del_err}")

            except Exception as pdf_error:
                self.log(f"  ⚠ Erreur conversion PDF: {pdf_error}")

            self.add_to_report('processed', {
                'file_name': Path(msg_path).name,
                'file_type': 'MSG',
                'files_found': len(extracted_names),
                'status': f'{len(extracted_names)} pièce(s) jointe(s) extraite(s)' if extracted_names else 'Traité avec succès'
            })

        except Exception as e:
            self.log(f"  ❌ Erreur MSG: {str(e)}")
            self.add_to_report('error', {
                'file_name': Path(msg_path).name,
                'error_type': 'Erreur de traitement MSG',
                'error_message': str(e),
                'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S')
            })
            self.add_to_report('processed', {
                'file_name': Path(msg_path).name,
                'file_type': 'MSG',
                'files_found': 0,
                'status': f'Erreur: {str(e)[:50]}'
            })

        return extracted_names

    def convert_msg_to_pdf(self, msg_path, output_dir, extracted_names=None):
        """
        Convertit MSG en PDF par capture complète du contenu.
        Insère "Voir document : nom" dans le Word avant export PDF.
        Supprime le fichier .msg après création réussie du PDF.
        """
        try:
            import win32com.client
            import pythoncom

            pythoncom.CoInitialize()

            try:
                # Nettoyer le nom du PDF (supprimer espaces en fin de stem)
                stem_clean = Path(msg_path).stem.strip()
                pdf_output = Path(str(output_dir).strip()) / f"{stem_clean}.pdf"

                self.log(f"    📄 PDF cible: {pdf_output.name}")

                outlook = win32com.client.Dispatch("Outlook.Application")
                msg = outlook.Session.OpenSharedItem(str(Path(msg_path).absolute()))

                inspector = msg.GetInspector
                inspector.Display(False)

                time.sleep(2)

                word_doc = inspector.WordEditor

                if word_doc is None:
                    self.log(f"    ⚠ Impossible d'accéder au contenu")
                    inspector.Close(0)
                    msg.Close(0)
                    return None

                self.log(f"    📏 Mesure du contenu...")

                try:
                    max_width = 0

                    for table in word_doc.Tables:
                        table_width = 0
                        for col in table.Columns:
                            try:
                                table_width += col.Width
                            except:
                                pass
                        if table_width > max_width:
                            max_width = table_width

                    for para in word_doc.Paragraphs:
                        try:
                            para_width = para.Range.Information(10)
                            if para_width > max_width:
                                max_width = para_width
                        except:
                            pass

                    try:
                        num_pages = word_doc.ComputeStatistics(2)
                    except:
                        num_pages = 1

                    self.log(f"    📊 Contenu: largeur={max_width:.0f}pt, pages={num_pages}")

                    MAX_DIMENSION = 1584
                    needed_width = max_width + 100 if max_width > 0 else 842

                    if needed_width > 800 or num_pages > 2:
                        pdf_width  = min(needed_width, MAX_DIMENSION)
                        pdf_height = min(842, MAX_DIMENSION)
                        orientation = 1
                        if needed_width > MAX_DIMENSION:
                            pdf_width = MAX_DIMENSION
                            ratio = needed_width / MAX_DIMENSION
                            pdf_height = min(int(842 * ratio), MAX_DIMENSION)
                        self.log(f"    📐 Mode: PAYSAGE pour contenu large")
                    else:
                        pdf_width  = max(595, min(needed_width, MAX_DIMENSION))
                        pdf_height = max(842, min(842 * num_pages, MAX_DIMENSION))
                        orientation = 0
                        self.log(f"    📐 Mode: PORTRAIT standard")

                    pdf_width  = max(72, min(pdf_width,  MAX_DIMENSION))
                    pdf_height = max(72, min(pdf_height, MAX_DIMENSION))
                    self.log(f"    📐 PDF: {pdf_width:.0f}pt x {pdf_height:.0f}pt")

                except Exception as measure_error:
                    self.log(f"    ⚠ Mesure échouée, utilisation format large")
                    pdf_width  = 1190
                    pdf_height = 842
                    orientation = 1

                word = None
                new_doc = None
                try:
                    self.log(f"    🔄 Création PDF avec dimensions adaptées...")

                    # DispatchEx = instance Word SÉPARÉE, indépendante d'Outlook
                    word = win32com.client.DispatchEx("Word.Application")
                    word.Visible = False
                    time.sleep(1)

                    new_doc = word.Documents.Add()
                    time.sleep(0.5)

                    try:
                        new_doc.PageSetup.Orientation  = orientation
                        new_doc.PageSetup.PageWidth    = pdf_width
                        new_doc.PageSetup.PageHeight   = pdf_height
                        new_doc.PageSetup.TopMargin    = 36
                        new_doc.PageSetup.BottomMargin = 36
                        new_doc.PageSetup.LeftMargin   = 36
                        new_doc.PageSetup.RightMargin  = 36
                        self.log(f"    ✓ Page configurée")
                    except Exception as setup_error:
                        self.log(f"    ⚠ Configuration page: {setup_error}")
                        try:
                            new_doc.PageSetup.Orientation = 1
                            new_doc.PageSetup.PageWidth   = 1190
                            new_doc.PageSetup.PageHeight  = 842
                        except Exception:
                            pass  # Continuer sans mise en page

                    # ── Copie du contenu MSG ──────────────────────────────
                    try:
                        word_doc.Content.Select()
                        word_doc.Application.Selection.Copy()
                        new_doc.Content.Paste()
                        self.log(f"    ✓ Contenu copié (texte + images + tableaux)")
                        time.sleep(1)
                    except Exception as copy_error:
                        self.log(f"    ⚠ Copie échouée: {copy_error}")
                        try:
                            word_doc.Content.Copy()
                            new_doc.Range().Paste()
                        except:
                            pass

                    # ── Insertion "Voir document" au début du doc Word ────
                    if extracted_names:
                        try:
                            self.log(f"    🔄 Insertion des références pièces jointes...")
                            sep   = "-" * 55
                            lines = [sep, f"PIECES JOINTES EXTRAITES ({len(extracted_names)}) :"]
                            for name in extracted_names:
                                lines.append(f"Voir {Path(name).stem}")
                            lines.append(sep)

                            # Insérer en ordre inverse au début
                            for line in reversed(lines):
                                r = new_doc.Range(0, 0)
                                r.InsertBefore(line + "\n")

                            # Mettre en gras + couleur le bloc
                            bloc = new_doc.Range(0, 0)
                            for _ in range(len(lines)):
                                bloc.MoveEnd(5, 1)
                            bloc.Bold = True
                            bloc.Font.Color = 0xFF0000

                            self.log(f"    ✅ Références insérées dans le document Word")
                        except Exception as insert_err:
                            self.log(f"    ⚠ Erreur insertion références: {insert_err}")

                    # ── Export PDF ────────────────────────────────────────
                    self.log(f"    📄 Export PDF...")

                    new_doc.ExportAsFixedFormat(
                        OutputFileName   = str(pdf_output.absolute()),
                        ExportFormat     = 17,
                        OpenAfterExport  = False,
                        OptimizeFor      = 0,
                        CreateBookmarks  = 0,
                        DocStructureTags = True,
                        BitmapMissingFonts=True
                    )

                    inspector.Close(0)
                    msg.Close(0)

                    time.sleep(1)

                    if pdf_output.exists() and pdf_output.stat().st_size > 0:
                        self.log(f"    ✅ PDF créé avec succès")
                        self.log(f"    📦 Taille: {pdf_output.stat().st_size} bytes")
                        msg_p = Path(msg_path)
                        pdf_p = Path(pdf_output)
                        if msg_p.parent.resolve() == pdf_p.parent.resolve():
                            for _retry in range(5):
                                try:
                                    msg_p.unlink()
                                    self.log(f"    🗑️ Fichier MSG supprimé : {msg_p.name}")
                                    break
                                except PermissionError:
                                    if _retry < 4:
                                        time.sleep(1)
                                except Exception as del_err:
                                    self.log(f"    ⚠ Impossible de supprimer le MSG : {del_err}")
                                    break
                        return pdf_output
                    else:
                        self.log(f"    ❌ PDF non créé")
                        return None

                except Exception as create_error:
                    self.log(f"    ❌ Erreur création: {create_error}")
                    try:
                        inspector.Close(0)
                        msg.Close(0)
                    except:
                        pass
                    return None
                finally:
                    # Garantir fermeture de Word dans tous les cas
                    try:
                        if new_doc is not None:
                            new_doc.Close(False)
                    except Exception:
                        pass
                    try:
                        if word is not None:
                            word.Quit()
                    except Exception:
                        pass
                    time.sleep(0.5)

            finally:
                pythoncom.CoUninitialize()

        except ImportError:
            self.log(f"  ⚠ Module pywin32 requis")
            return None
        except Exception as e:
            self.log(f"  ❌ Erreur: {e}")
            return None


    def _convert_msg_to_pdf_with_attachments(self, msg_com_obj, msg_path, output_dir, extracted_names):
        """
        Convertit un MSG en PDF en construisant le contenu directement
        depuis les propriétés Outlook (sujet, expéditeur, corps, destinataires).
        Insère EN HAUT la liste des pièces jointes extraites.
        Supprime le fichier .msg après création réussie du PDF.
        """
        try:
            import win32com.client
            import pythoncom
            import time

            pdf_output = Path(output_dir) / f"{Path(msg_path).stem}.pdf"

            # ============================================================
            # ÉTAPE 1 : LIRE LES PROPRIÉTÉS DU MSG
            # ============================================================
            self.log(f"    📖 Lecture des propriétés du MSG...")

            try:
                subject    = str(msg_com_obj.Subject   or "")
                sender     = str(msg_com_obj.SenderName or "")
                sender_email = str(msg_com_obj.SenderEmailAddress or "")
                body_text  = str(msg_com_obj.Body      or "")
                sent_on    = ""
                try:
                    sent_on = str(msg_com_obj.SentOn)
                except Exception:
                    pass

                # Destinataires
                recipients_list = []
                try:
                    for i in range(1, msg_com_obj.Recipients.Count + 1):
                        r = msg_com_obj.Recipients.Item(i)
                        recipients_list.append(str(r.Name or r.Address or ""))
                except Exception:
                    pass
                recipients_str = " ; ".join(recipients_list) if recipients_list else ""

                self.log(f"    ✅ Sujet : {subject}")
                self.log(f"    ✅ Expéditeur : {sender}")
                self.log(f"    ✅ Corps : {len(body_text)} caractère(s)")

            except Exception as prop_err:
                self.log(f"    ⚠️ Erreur lecture propriétés : {prop_err}")
                subject = Path(msg_path).stem
                sender = ""
                sender_email = ""
                body_text = ""
                sent_on = ""
                recipients_str = ""

            # ============================================================
            # ÉTAPE 2 : CRÉER UN DOCUMENT WORD AVEC LE CONTENU
            # ============================================================
            self.log(f"    📄 Création du document Word...")

            word = win32com.client.DispatchEx("Word.Application")  # instance séparée
            word.Visible = False
            time.sleep(1)
            new_doc = word.Documents.Add()
            time.sleep(1)

            try:
                # Marges réduites
                new_doc.PageSetup.TopMargin    = 50
                new_doc.PageSetup.BottomMargin = 50
                new_doc.PageSetup.LeftMargin   = 60
                new_doc.PageSetup.RightMargin  = 60
            except Exception:
                pass

            # Helper pour ajouter un paragraphe formaté
            def add_paragraph(text, bold=False, font_size=11, color=None, separator=False):
                try:
                    para = new_doc.Content
                    para.Collapse(0)  # wdCollapseEnd

                    if separator:
                        para.InsertAfter("-" * 80)
                    else:
                        para.InsertAfter(text)

                    para.InsertParagraphAfter()

                    # Appliquer formatage sur le paragraphe ajouté
                    all_paras = new_doc.Paragraphs
                    last_para = all_paras.Item(all_paras.Count)
                    rng = last_para.Range

                    if bold:
                        rng.Bold = True
                    if font_size:
                        rng.Font.Size = font_size
                    if color:
                        rng.Font.Color = color  # RGB en format Word (BGR int)

                except Exception:
                    pass

            # ── SECTION 1 : PIÈCES JOINTES EXTRAITES (en tête) ──────────
            if extracted_names:
                add_paragraph("PIÈCES JOINTES EXTRAITES", bold=True, font_size=13,
                            color=0x0000CC)
                add_paragraph("", separator=True)

                for name in extracted_names:
                    name_sans_ext = Path(name).stem
                    add_paragraph(f"  Voir {name_sans_ext}",
                                bold=True, font_size=11, color=0x0000CC)

                add_paragraph("", separator=True)
                add_paragraph("")  # ligne vide

            # ── SECTION 2 : EN-TÊTE DU MESSAGE ──────────────────────────
            add_paragraph("INFORMATIONS DU MESSAGE", bold=True, font_size=12)
            add_paragraph("", separator=True)

            if subject:
                add_paragraph(f"Objet      : {subject}", bold=True, font_size=11)
            if sender:
                if sender_email and sender_email != sender:
                    add_paragraph(f"De         : {sender} <{sender_email}>",
                                bold=False, font_size=10)
                else:
                    add_paragraph(f"De         : {sender}", bold=False, font_size=10)
            if recipients_str:
                add_paragraph(f"À          : {recipients_str}", bold=False, font_size=10)
            if sent_on:
                add_paragraph(f"Date       : {sent_on}", bold=False, font_size=10)
            if extracted_names:
                add_paragraph(
                    f"Pièces jointes ({len(extracted_names)}) : "
                    + " | ".join(Path(n).stem for n in extracted_names),
                    bold=False, font_size=10
                )

            add_paragraph("", separator=True)
            add_paragraph("")  # ligne vide

            # ── SECTION 3 : CORPS DU MESSAGE ────────────────────────────
            add_paragraph("CORPS DU MESSAGE", bold=True, font_size=12)
            add_paragraph("", separator=True)
            add_paragraph("")

            if body_text.strip():
                # Découper en lignes pour éviter les très longs paragraphes
                lines = body_text.split('\n')
                for line in lines:
                    line = line.rstrip('\r')
                    try:
                        rng = new_doc.Content
                        rng.Collapse(0)
                        rng.InsertAfter(line if line.strip() else " ")
                        rng.InsertParagraphAfter()
                    except Exception:
                        pass
            else:
                add_paragraph("(Corps du message vide)", bold=False, font_size=10)

            # ── EXPORT PDF ───────────────────────────────────────────────
            self.log(f"    📤 Export PDF...")

            new_doc.ExportAsFixedFormat(
                OutputFileName=str(pdf_output.absolute()),
                ExportFormat=17,
                OpenAfterExport=False,
                OptimizeFor=0,
                CreateBookmarks=0,
                DocStructureTags=True,
                BitmapMissingFonts=True
            )

            new_doc.Close(False)
            word.Quit()

            time.sleep(1)

            if pdf_output.exists() and pdf_output.stat().st_size > 0:
                self.log(f"    ✅ PDF créé avec succès")
                self.log(f"    📦 Taille: {pdf_output.stat().st_size} bytes")
                # Suppression du .msg — seulement si le fichier est dans un dossier de sortie
                # (pas le fichier source original situé ailleurs)
                msg_p = Path(msg_path)
                pdf_p = Path(pdf_output)
                if msg_p.parent.resolve() == pdf_p.parent.resolve():
                    try:
                        msg_p.unlink()
                        self.log(f"    🗑️ Fichier MSG supprimé : {msg_p.name}")
                    except Exception as del_err:
                        self.log(f"    ⚠ Impossible de supprimer le MSG : {del_err}")
                return pdf_output
            else:
                self.log(f"    ⚠️ PDF vide ou absent")
                return None

        except Exception as e:
            self.log(f"    ❌ Erreur conversion PDF: {e}")
            import traceback
            self.log(f"    📋 {traceback.format_exc()}")
            try:
                word.Quit()
            except Exception:
                pass
            return None
    def generate_excel_report(self, output_path):
            """
            Génère un rapport Excel détaillé des résultats.
            
            FEUILLES :
            1. Résumé              — statistiques globales
            2. Fichiers Traités    — tous les fichiers sources traités
            3. Fichiers Extraits   — fichiers extraits avec succès
            4. Non Archivables     — fichiers conservés (type non supporté SharePoint ou vides)
            """
            try:
                wb = openpyxl.Workbook()

                # ================================================================
                # Couleurs / styles réutilisables
                # ================================================================
                COLOR_HEADER_DARK   = "2C3E50"
                COLOR_HEADER_BLUE   = "2980B9"
                COLOR_HEADER_GREEN  = "16A085"
                COLOR_HEADER_ORANGE = "E67E22"
                COLOR_HEADER_RED    = "C0392B"

                COLOR_ROW_GREEN  = "D5F4E6"
                COLOR_ROW_ORANGE = "FFF4E6"
                COLOR_ROW_RED    = "FADBD8"
                COLOR_ROW_GRAY   = "F2F3F4"

                def header_cell(cell, text, bg_color, font_size=11):
                    cell.value = text
                    cell.font      = Font(bold=True, color="FFFFFF", size=font_size)
                    cell.fill      = PatternFill(start_color=bg_color, end_color=bg_color, fill_type="solid")
                    cell.alignment = Alignment(horizontal='center', vertical='center')

                def color_row(ws, row_num, nb_cols, hex_color):
                    for c in range(1, nb_cols + 1):
                        ws.cell(row=row_num, column=c).fill = PatternFill(
                            start_color=hex_color, end_color=hex_color, fill_type="solid"
                        )

                # ================================================================
                # Séparer les fichiers extraits des non-archivables
                # ================================================================
                files_ok            = []   # extraits avec succès
                files_non_archivable = []  # type non supporté OU vides

                for entry in self.report_data['files_extracted']:
                    status = entry.get('status', '')
                    ext_file = entry.get('extracted_file', '')

                    is_non_archivable = (
                        'non supporté' in status.lower()
                        or 'conservé'  in status.lower()
                        or 'ignoré'    in status.lower()
                        or ext_file.startswith('[')         # ex: "[Non supporté: .ole]"
                    )

                    if is_non_archivable:
                        files_non_archivable.append(entry)
                    else:
                        files_ok.append(entry)

                total_processed       = len(self.report_data['files_processed'])
                total_extracted       = len(files_ok)
                total_non_archivable  = len(files_non_archivable)

                # ================================================================
                # FEUILLE 1 : RÉSUMÉ
                # ================================================================
                ws_sum = wb.active
                ws_sum.title = "Résumé"

                ws_sum.merge_cells('A1:D1')
                header_cell(ws_sum['A1'], "RAPPORT D'EXTRACTION DE FICHIERS INCORPORÉS",
                            COLOR_HEADER_DARK, font_size=15)
                ws_sum.row_dimensions[1].height = 32

                ws_sum.merge_cells('A2:D2')
                ws_sum['A2'] = f"Généré le : {datetime.now().strftime('%d/%m/%Y à %H:%M:%S')}"
                ws_sum['A2'].font      = Font(italic=True, size=10)
                ws_sum['A2'].alignment = Alignment(horizontal='center')

                # Bloc statistiques
                ws_sum.merge_cells('A4:B4')
                header_cell(ws_sum['A4'], "STATISTIQUES GLOBALES", COLOR_HEADER_BLUE, font_size=13)
                ws_sum.row_dimensions[4].height = 26

                stats = [
                    ("Volume du lot chargé :",          self.total_input_files, "8E44AD"),
                    ("Fichiers sources traités :",      total_processed,       "27AE60"),
                    ("Fichiers extraits avec succès :", total_extracted,       "2980B9"),
                    ("Fichiers non archivables :",      total_non_archivable,  "E67E22"),
                ]
                for i, (label, value, color) in enumerate(stats, start=6):
                    ws_sum[f'A{i}'] = label
                    ws_sum[f'A{i}'].font = Font(bold=True, size=11)
                    ws_sum[f'B{i}'] = value
                    ws_sum[f'B{i}'].font      = Font(bold=True, size=12, color=color)
                    ws_sum[f'B{i}'].alignment = Alignment(horizontal='center')
                    ws_sum.row_dimensions[i].height = 20

                ws_sum.column_dimensions['A'].width = 38
                ws_sum.column_dimensions['B'].width = 22

                # ================================================================
                # FEUILLE 2 : FICHIERS TRAITÉS
                # ================================================================
                ws_proc = wb.create_sheet("Fichiers Traités")
                headers = ['N°', 'Nom du fichier', 'Type', 'Fichiers trouvés', 'Statut']
                for col, h in enumerate(headers, 1):
                    header_cell(ws_proc.cell(1, col), h, COLOR_HEADER_DARK)
                ws_proc.row_dimensions[1].height = 24

                for idx, d in enumerate(self.report_data['files_processed'], 1):
                    row = idx + 1
                    ws_proc.cell(row, 1, idx)
                    ws_proc.cell(row, 2, d.get('file_name', ''))
                    ws_proc.cell(row, 3, d.get('file_type', ''))
                    ws_proc.cell(row, 4, d.get('files_found', 0))
                    ws_proc.cell(row, 5, d.get('status', ''))

                    ws_proc.cell(row, 1).alignment = Alignment(horizontal='center')
                    ws_proc.cell(row, 3).alignment = Alignment(horizontal='center')
                    ws_proc.cell(row, 4).alignment = Alignment(horizontal='center')

                    status = d.get('status', '')
                    if 'Erreur' in status:
                        color_row(ws_proc, row, 5, COLOR_ROW_RED)
                    elif 'extrait' in status or 'succès' in status.lower():
                        color_row(ws_proc, row, 5, COLOR_ROW_GREEN)

                ws_proc.column_dimensions['A'].width = 6
                ws_proc.column_dimensions['B'].width = 42
                ws_proc.column_dimensions['C'].width = 12
                ws_proc.column_dimensions['D'].width = 18
                ws_proc.column_dimensions['E'].width = 40

                # ================================================================
                # FEUILLE 3 : FICHIERS EXTRAITS (succès uniquement)
                # ================================================================
                ws_ext = wb.create_sheet("Fichiers Extraits")
                headers_ext = ['N°', 'Fichier source', 'Fichier extrait', 'Position', 'Type', 'Statut']
                for col, h in enumerate(headers_ext, 1):
                    header_cell(ws_ext.cell(1, col), h, COLOR_HEADER_GREEN)
                ws_ext.row_dimensions[1].height = 24

                if files_ok:
                    RETRAITEMENT_EXTS = {'.zip', '.7z', '.htm', '.html'}

                    for idx, d in enumerate(files_ok, 1):
                        row = idx + 1
                        ws_ext.cell(row, 1, idx)
                        ws_ext.cell(row, 2, d.get('source_file', ''))
                        ws_ext.cell(row, 3, d.get('extracted_file', ''))
                        ws_ext.cell(row, 4, d.get('position', ''))
                        file_type = d.get('type', '')
                        ws_ext.cell(row, 5, file_type)

                        # Statut enrichi pour les types nécessitant un retraitement
                        base_status = d.get('status', '')
                        if file_type.lower() in RETRAITEMENT_EXTS:
                            status = 'Succès + Retraitement'
                            ws_ext.cell(row, 6, status)
                            ws_ext.cell(row, 6).font = Font(bold=True, color="E67E22", size=10)
                        else:
                            ws_ext.cell(row, 6, base_status)
                            ws_ext.cell(row, 6).font = Font(bold=True, color="27AE60", size=10)
                        ws_ext.cell(row, 1).alignment = Alignment(horizontal='center')
                        ws_ext.cell(row, 4).alignment = Alignment(horizontal='center')
                        ws_ext.cell(row, 5).alignment = Alignment(horizontal='center')
                        ws_ext.cell(row, 6).alignment = Alignment(horizontal='center')

                        color_row(ws_ext, row, 6, COLOR_ROW_GREEN)
                else:
                    ws_ext.merge_cells('A2:F2')
                    ws_ext['A2'] = "Aucun fichier extrait"
                    ws_ext['A2'].font      = Font(italic=True, color="7F8C8D", size=11)
                    ws_ext['A2'].alignment = Alignment(horizontal='center')

                ws_ext.column_dimensions['A'].width = 6
                ws_ext.column_dimensions['B'].width = 38
                ws_ext.column_dimensions['C'].width = 42
                ws_ext.column_dimensions['D'].width = 14
                ws_ext.column_dimensions['E'].width = 10
                ws_ext.column_dimensions['F'].width = 16

                # ================================================================
                # FEUILLE 4 : FICHIERS NON ARCHIVABLES
                # (type non supporté SharePoint OU fichiers vides)
                # ================================================================
                ws_na = wb.create_sheet("Non Archivables")

                # Titre explicatif
                ws_na.merge_cells('A1:F1')
                header_cell(ws_na['A1'],
                            "FICHIERS NON ARCHIVABLES — Conservés dans le document source",
                            COLOR_HEADER_ORANGE, font_size=13)
                ws_na.row_dimensions[1].height = 28

                ws_na.merge_cells('A2:F2')
                ws_na['A2'] = (
                    "⚠️  Ces fichiers restent intacts dans leur document d'origine. "
                    "Ils ne peuvent pas être archivés sur SharePoint (type non supporté ou fichier vide)."
                )
                ws_na['A2'].font      = Font(italic=True, size=10, color="7E5109")
                ws_na['A2'].fill      = PatternFill(start_color="FDEBD0", end_color="FDEBD0", fill_type="solid")
                ws_na['A2'].alignment = Alignment(horizontal='left', wrap_text=True)
                ws_na.row_dimensions[2].height = 30

                # En-têtes colonnes
                headers_na = ['N°', 'Fichier source', 'Description', 'Position / Slide', 'Type détecté', 'Raison']
                for col, h in enumerate(headers_na, 1):
                    header_cell(ws_na.cell(4, col), h, COLOR_HEADER_ORANGE)
                ws_na.row_dimensions[4].height = 24

                if files_non_archivable:
                    for idx, d in enumerate(files_non_archivable, 1):
                        row = idx + 4
                        ws_na.cell(row, 1, idx)
                        ws_na.cell(row, 2, d.get('source_file', ''))

                        # Description lisible du fichier non archivable
                        ext_file = d.get('extracted_file', '')
                        if '[Non supporté:' in ext_file:
                            desc = ext_file.replace('[', '').replace(']', '')
                        elif '[Vide' in ext_file or '0 Ko' in ext_file:
                            desc = "Fichier vide (0 Ko)"
                        else:
                            desc = ext_file

                        ws_na.cell(row, 3, desc)
                        ws_na.cell(row, 4, d.get('position', ''))
                        ws_na.cell(row, 5, d.get('type', ''))

                        # Raison claire
                        status = d.get('status', '')
                        if '0 Ko' in status or 'vide' in status.lower():
                            raison = "Fichier vide — impossible à archiver"
                        else:
                            raison = "Type de fichier non supporté par l'archivage SharePoint"

                        ws_na.cell(row, 6, raison)

                        # Mise en forme
                        ws_na.cell(row, 1).alignment = Alignment(horizontal='center')
                        ws_na.cell(row, 4).alignment = Alignment(horizontal='center')
                        ws_na.cell(row, 5).alignment = Alignment(horizontal='center')
                        ws_na.cell(row, 3).font = Font(bold=True, color="7E5109")
                        ws_na.cell(row, 6).font = Font(italic=True, color="7E5109")

                        color_row(ws_na, row, 6, COLOR_ROW_ORANGE)
                else:
                    ws_na.merge_cells('A5:F5')
                    ws_na['A5'] = "✅ Aucun fichier non archivable détecté"
                    ws_na['A5'].font      = Font(bold=True, color="27AE60", size=12)
                    ws_na['A5'].alignment = Alignment(horizontal='center')
                    ws_na.row_dimensions[5].height = 28

                ws_na.column_dimensions['A'].width = 6
                ws_na.column_dimensions['B'].width = 38
                ws_na.column_dimensions['C'].width = 32
                ws_na.column_dimensions['D'].width = 18
                ws_na.column_dimensions['E'].width = 14
                ws_na.column_dimensions['F'].width = 50

                # ================================================================
                # Sauvegarder
                # ================================================================
                wb.save(output_path)
                return True

            except Exception as e:
                print(f"Erreur génération rapport Excel: {e}")
                traceback.print_exc()
                return False
    def process_file(self, file_path, output_dir):
            """Traite un fichier et extrait les objets incorporés"""
            file_path = Path(file_path)
            output_dir = Path(output_dir)

            self.log(f"📄 Traitement de: {file_path.name}")

            file_ext = file_path.suffix.lower()
            extracted_files = []

            try:
                # ========================================
                # CONVERSION ANCIEN FORMAT SI NÉCESSAIRE
                # ========================================
                if file_ext in ['.doc', '.ppt', '.xls']:
                    original_legacy_name = file_path.name
                    converted_path, was_converted = self.convert_legacy_to_modern_format(file_path)

                    if was_converted:
                        self.log(f"  ✅ Fichier converti, traitement du nouveau format")
                        file_path = converted_path
                        file_ext = file_path.suffix.lower()
                        # Marquer le fichier converti comme traité (évite double récursion)
                        self._processed_absolute_paths.add(str(file_path.resolve()))
                        # Supprimer toute copie legacy déjà présente dans output_dir
                        legacy_in_output = output_dir / original_legacy_name
                        if legacy_in_output.exists():
                            try:
                                legacy_in_output.unlink()
                                self.log(f"  🗑️ Fichier legacy supprimé de la sortie : {original_legacy_name}")
                            except Exception as del_err:
                                self.log(f"  ⚠️ Impossible de supprimer {original_legacy_name}: {del_err}")
                    else:
                        self.log(f"  ⚠️ Conversion échouée, traitement de l'ancien format")

                # ========================================
                # TRAITEMENT SELON LE FORMAT
                # ========================================
                if file_ext in ['.pptx', '.pptm']:
                    extracted_files = self.extract_from_pptx(file_path, output_dir)

                elif file_ext == '.docx':
                    extracted_files = self.extract_ole_from_docx(file_path, output_dir)

                elif file_ext in ['.xlsx', '.xlsm']:
                    extracted_files = self.extract_from_xlsx(file_path, output_dir)

                elif file_ext == '.pdf':
                    extracted_files = self.extract_from_pdf(file_path, output_dir)

                elif file_ext == '.msg':
                    extracted_files = self.extract_from_msg_file(file_path, output_dir)
                    
                # ========================================
                # TRAITEMENT DES ANCIENS FORMATS (fallback si conversion échouée)
                # ========================================
                elif file_ext in ['.doc', '.ppt', '.xls']:
                    self.log(f"  📦 Traitement ancien format (méthode legacy)")
                    extracted_files = self.extract_ole_from_legacy(file_path, output_dir)

                else:
                    self.log(f"  ⚠️ Format non supporté: {file_ext}")
                    return []

                # ── Fallback universel : dossier vide → copie de l'original ──
                try:
                    if output_dir.exists() and not any(output_dir.iterdir()):
                        shutil.copy2(file_path, output_dir / file_path.name)
                        self.log(f"  📋 Aucun fichier extrait — copie de l'original : {file_path.name}")
                except Exception:
                    pass

                return extracted_files

            except Exception as e:
                self.log(f"  ❌ Erreur traitement: {str(e)}")
                import traceback
                self.log(f"  📋 Détails: {traceback.format_exc()}")
                # Fallback : copier l'original si le dossier de sortie est vide
                try:
                    if output_dir.exists() and not any(output_dir.iterdir()):
                        shutil.copy2(file_path, output_dir / file_path.name)
                        self.log(f"  📋 Erreur → copie de l'original : {file_path.name}")
                except Exception:
                    pass
                return []
    def _safe_copy2(src, dst, retries=3, delay=2.0):
        import time as _t
        last_err = None
        for attempt in range(retries):
            try:
                shutil.copy2(src, dst)
                return
            except PermissionError as e:
                last_err = e
                if attempt < retries - 1:
                    _t.sleep(delay)
        raise last_err

class SimpleFileExtractorGUI:
    """Interface graphique avec activation Excel intégrée"""
    
    def __init__(self, root):
        self.root = root
        self.root.title("EXTRACTEUR DES FICHIERS INCORPORES")
        self.root.geometry("900x850")
        self.root.configure(bg='#ecf0f1')
        
        self.selected_files = []
        self.output_dir = tk.StringVar(value=str(Path.home() / "Documents" / "Fichiers_Extraits"))
        
        self.reference_file = None
        self.reference_config = None
        self.stop_event = threading.Event()
        self.extraction_thread = None
        self.setup_ui()
    def _validate_index_input(self, new_value):
        """Accepte uniquement les entiers positifs ou la chaîne vide"""
        if new_value == "":
            return True
        try:
            v = int(new_value)
            return v >= 1
        except ValueError:
            return False
    
    def setup_ui(self):
        """Configure l'interface"""
        style = ttk.Style()
        style.theme_use('clam')
        
        # En-tête
        header_frame = tk.Frame(self.root, bg='#2c3e50', height=80)
        header_frame.pack(fill=tk.X)
        header_frame.pack_propagate(False)

        title_label = tk.Label(header_frame, text="📎 EXTRACTEUR DES FICHIERS INCORPORES",
                               font=('Arial', 16, 'bold'), bg='#2c3e50', fg='white')
        title_label.pack(side=tk.LEFT, padx=20, pady=20)

        version_label = tk.Label(header_frame, text="SES MAROC  V1.0.6",
                                 font=('Arial', 10, 'bold'), bg='#2c3e50', fg='#aab7c4')
        version_label.pack(side=tk.RIGHT, padx=20, pady=20)
        
        # Frame principal
        main_frame = ttk.Frame(self.root, padding="20")
        main_frame.pack(fill=tk.BOTH, expand=True)
        
        # Section 1: Fichiers
        file_section = ttk.LabelFrame(main_frame, text="📁 Fichiers à traiter", padding="15")
        file_section.pack(fill=tk.BOTH, expand=True, pady=(0, 10))
        
        button_frame = ttk.Frame(file_section)
        button_frame.pack(fill=tk.X, pady=(0, 10))
        
        tk.Button(button_frame, text="➕ Ajouter fichiers", command=self.add_files,
                bg='#3498db', fg='white', font=('Arial', 10, 'bold'),
                padx=15, pady=8).pack(side=tk.LEFT, padx=5)
        
        tk.Button(button_frame, text="🗑️ Vider", command=self.clear_files,
                bg='#e74c3c', fg='white', font=('Arial', 10, 'bold'),
                padx=15, pady=8).pack(side=tk.LEFT, padx=5)
        
        list_frame = ttk.Frame(file_section)
        list_frame.pack(fill=tk.BOTH, expand=True)
        
        scrollbar = ttk.Scrollbar(list_frame)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        self.file_listbox = tk.Listbox(list_frame, yscrollcommand=scrollbar.set,
                                    font=('Consolas', 9), bg='#ecf0f1',
                                    selectmode=tk.EXTENDED, height=8)
        self.file_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar.config(command=self.file_listbox.yview)
        
        # Section 2: Configuration
        config_section = ttk.LabelFrame(main_frame, text="⚙️ Configuration", padding="15")
        config_section.pack(fill=tk.X, pady=(0, 10))
        
        # Dossier sortie
        output_frame = ttk.Frame(config_section)
        output_frame.pack(fill=tk.X, pady=5)
        
        tk.Label(output_frame, text="📂 Dossier de sortie :", 
                font=('Arial', 10), bg='#ecf0f1').pack(side=tk.LEFT, padx=(0, 10))
        
        ttk.Entry(output_frame, textvariable=self.output_dir, width=40).pack(side=tk.LEFT, padx=5)
        
        tk.Button(output_frame, text="📁", command=self.select_output,
                bg='#95a5a6', fg='white', font=('Arial', 10, 'bold'),
                padx=10, pady=5).pack(side=tk.LEFT, padx=5)
        
        tk.Button(output_frame, text="Ouvrir", command=self.open_output_folder,
                bg='#95a5a6', fg='white', font=('Arial', 10),
                padx=10, pady=5).pack(side=tk.LEFT, padx=5)
        
        # Fichier de référence
        ref_frame = ttk.Frame(config_section)
        ref_frame.pack(fill=tk.X, pady=5)
        
        tk.Label(ref_frame, text="📊 Fichier de référence (optionnel) :", 
                font=('Arial', 10), bg='#ecf0f1').pack(side=tk.LEFT, padx=(0, 10))
        
        self.ref_label = tk.Label(ref_frame, text="Aucun", 
                                font=('Arial', 9), bg='#ecf0f1', fg='#7f8c8d')
        self.ref_label.pack(side=tk.LEFT, padx=5)
        
        tk.Button(ref_frame, text="Charger référence", command=self.load_reference,
                bg='#9b59b6', fg='white', font=('Arial', 10),
                padx=10, pady=5).pack(side=tk.LEFT, padx=5)

        # ← NOUVEAU BOUTON SUPPRIMER RÉFÉRENCE
        tk.Button(ref_frame, text="🗑️ Supprimer référence", command=self.clear_reference,
                bg='#e74c3c', fg='white', font=('Arial', 10),
                padx=10, pady=5).pack(side=tk.LEFT, padx=5)

        # Index de départ manuel
        index_frame = ttk.Frame(config_section)
        index_frame.pack(fill=tk.X, pady=5)

        tk.Label(index_frame, text="🔢 Index de départ (FJ_X) :",
                font=('Arial', 10), bg='#ecf0f1').pack(side=tk.LEFT, padx=(0, 10))

        self.start_index_var = tk.StringVar(value="")
        self.start_index_entry = ttk.Entry(
            index_frame,
            textvariable=self.start_index_var,
            width=10,
            font=('Arial', 10)
        )
        self.start_index_entry.pack(side=tk.LEFT, padx=5)

        tk.Label(index_frame,
                text="(optionnel — remplace le fichier de référence si renseigné)",
                font=('Arial', 9, 'italic'), bg='#ecf0f1', fg='#7f8c8d').pack(side=tk.LEFT, padx=5)

        # Validation : uniquement des chiffres
        vcmd = (self.root.register(self._validate_index_input), '%P')
        self.start_index_entry.config(validate='key', validatecommand=vcmd)
        
        # Section 3: Actions
        action_frame = ttk.Frame(main_frame)
        action_frame.pack(fill=tk.X, pady=10)
        
        btn_row = ttk.Frame(action_frame)
        btn_row.pack()

        self.extract_btn = tk.Button(btn_row, text="🚀 ACTIVER EXCEL + EXTRAIRE", 
                                    command=self.extract_files,
                                    bg='#27ae60', fg='white', 
                                    font=('Arial', 13, 'bold'),
                                    padx=30, pady=15)
        self.extract_btn.pack(side=tk.LEFT, padx=10)

        self.stop_btn = tk.Button(btn_row, text="⛔ ARRÊTER",
                                command=self.stop_extraction,
                                bg='#c0392b', fg='white',
                                font=('Arial', 13, 'bold'),
                                padx=20, pady=15,
                                state='disabled')
        self.stop_btn.pack(side=tk.LEFT, padx=10)
        
        info_label = tk.Label(action_frame, 
                            text="✅ Activation Excel | ✅ ZIP remplacés | ✅ Nomenclature FJ",
                            font=('Arial', 9, 'italic'), bg='#ecf0f1', fg='#7f8c8d')
        info_label.pack(pady=5)

        # Section progression
        progress_section = ttk.LabelFrame(main_frame, text="📊 Progression", padding="10")
        progress_section.pack(fill=tk.X, pady=(0, 10))

        self.progress_label = tk.Label(progress_section, text="En attente...",
                                        font=('Arial', 10, 'bold'), bg='#ecf0f1', fg='#2c3e50')
        self.progress_label.pack(anchor='w')

        self.progress_bar = ttk.Progressbar(progress_section, mode='determinate', length=100)
        self.progress_bar.pack(fill=tk.X, pady=5)

        self.step_label = tk.Label(progress_section, text="",
                                    font=('Arial', 9, 'italic'), bg='#ecf0f1', fg='#7f8c8d')
        self.step_label.pack(anchor='w')

        # Section 4: Log
        log_section = ttk.LabelFrame(main_frame, text="📋 Journal", padding="10")
        log_section.pack(fill=tk.BOTH, expand=True)
        
        self.log_text = scrolledtext.ScrolledText(log_section, height=10, width=80,
                                                bg='#2c3e50', fg='#ecf0f1',
                                                font=('Consolas', 9), wrap=tk.WORD)
        self.log_text.pack(fill=tk.BOTH, expand=True)
        
        self.log("✓ Application prête")
        self.log("✅ CORRECTION 1 : Activation Excel fonctionnelle")
        self.log("✅ CORRECTION 2 : Fichiers ZIP remplacés dans documents")
        self.log("✅ CORRECTION 3 : Nomenclature FJ corrigée")
        
    def log(self, message):
        """Affiche message dans le log"""
        self.log_text.insert(tk.END, f"{message}\n")
        self.log_text.see(tk.END)
        self.root.update_idletasks()
    def clear_reference(self):
        """Supprime le fichier de référence chargé"""
        self.reference_file = None
        self.reference_config = None
        self.ref_label.config(text="Aucun", fg='#7f8c8d')
        self.log("✓ Fichier de référence supprimé")
    def stop_extraction(self):
        """Demande l'arrêt de l'extraction"""
        self.stop_event.set()
        self.stop_btn.config(state='disabled')
        self.log("⛔ Arrêt demandé — en attente de fin de l'opération courante...")

    def _update_progress(self, current, total, label_text, step_text=""):
        """Met à jour la barre de progression (thread-safe)"""
        def _update():
            if total > 0:
                pct = int(current / total * 100)
                self.progress_bar['value'] = pct
                self.progress_label.config(text=f"{label_text} — {current}/{total} ({pct}%)")
            else:
                self.progress_label.config(text=label_text)
            if step_text:
                self.step_label.config(text=step_text)
        self.root.after(0, _update)

    def _append_log(self, message):
        """Ajoute un message au log (thread-safe)"""
        def _update():
            self.log_text.insert(tk.END, f"{message}\n")
            self.log_text.see(tk.END)
        self.root.after(0, _update)


    def load_reference(self):
        """Charge fichier de référence avec sélection des colonnes"""
        file_path = filedialog.askopenfilename(
            title="Sélectionner le fichier Excel de référence",
            filetypes=[("Excel", "*.xlsx;*.xlsm;*.xls"), ("Tous", "*.*")]
        )
        
        if file_path:
            try:
                dialog = SheetColumnSelectorDialog(self.root, file_path)
                config = dialog.get_result()
                
                if config:
                    self.reference_file = file_path
                    self.reference_config = config
                    
                    filename = Path(file_path).name
                    self.ref_label.config(text=f"✓ {filename}", fg='#27ae60')
                    
                    self.log("✓ Référence chargée avec succès:")
                    self.log(f"  Fichier: {filename}")
                    self.log(f"  Feuille: {config['sheet_name']}")
                    self.log(f"  Colonne noms: {config['filename_col_letter']}")
                    self.log(f"  Colonne enfants: {config['children_col_letter']}")
                else:
                    self.log("⚠ Configuration de référence annulée")
            
            except Exception as e:
                messagebox.showerror("Erreur", f"Erreur lors du chargement:\n{e}")
                self.log(f"❌ Erreur: {e}")
    
    def add_files(self):
        """Ajoute fichiers"""
        files = filedialog.askopenfilenames(
            title="Sélectionner documents",
            filetypes=[
                ("Tous supportés", "*.docx;*.doc;*.xlsx;*.xlsm;*.xls;*.pptx;*.pptm;*.ppt;*.pdf;*.msg"),
                ("Word", "*.docx;*.doc"),
                ("Excel", "*.xlsx;*.xlsm;*.xls"),
                ("PowerPoint", "*.pptx;*.pptm;*.ppt"),
                ("PDF", "*.pdf"),
                ("Outlook MSG", "*.msg"),
                ("Tous", "*.*")
            ]
        )
        
        for file in files:
            if file not in self.selected_files:
                self.selected_files.append(file)
                self.file_listbox.insert(tk.END, f"📄 {Path(file).name}")
        
        if files:
            self.log(f"✓ {len(files)} fichier(s) ajouté(s)")
    
    def clear_files(self):
        """Vide liste"""
        self.selected_files = []
        self.file_listbox.delete(0, tk.END)
        self.log("✓ Liste vidée")
    
    def select_output(self):
        """Sélectionne dossier sortie"""
        folder = filedialog.askdirectory(title="Dossier de sortie")
        if folder:
            self.output_dir.set(folder)
            self.log(f"✓ Dossier: {folder}")
    
    def extract_files(self):
            """Lance activation Excel puis extraction dans un thread séparé"""
            if not self.selected_files:
                messagebox.showwarning("Attention", "Aucun fichier sélectionné!")
                return

            # Reset stop event et lancement thread
            self.stop_event.clear()
            self.extract_btn.config(state='disabled', text="⏳ En cours...")
            self.stop_btn.config(state='normal')
            self.progress_bar['value'] = 0
            self.progress_label.config(text="Démarrage...")
            self.step_label.config(text="")

            self.extraction_thread = threading.Thread(
                target=self._run_extraction, daemon=True
            )
            self.extraction_thread.start()

    def _run_extraction(self):
            """Exécution de l'extraction dans un thread secondaire"""
            import hashlib

            def log(msg):
                self._append_log(msg)

            def progress(current, total, label, step=""):
                self._update_progress(current, total, label, step)

            session_dir = None
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            files_with_excel = []
            total_extracted = 0
            extractor = None
            stopped_early = False
            files_processed_count = 0
            directories_to_recurse = []

            try:
                # Créer dossier de sortie
                output_base = Path(self.output_dir.get())
                session_dir = output_base / f"Extraction_{timestamp}"
                session_dir.mkdir(parents=True, exist_ok=True)

                log("")
                log("=" * 70)
                log(f"DÉBUT - {datetime.now().strftime('%H:%M:%S')}")
                log(f"Dossier: {session_dir}")
                log("=" * 70)

                total_files = len(self.selected_files)

               

                # =================================================================
                # ÉTAPE 2 : EXTRACTION INITIALE
                # =================================================================
                if not stopped_early:
                    log("")
                    log("📦 ÉTAPE 2 : EXTRACTION INITIALE")
                    log("=" * 70)
                    progress(0, total_files, "Étape 2/4 — Extraction initiale", "Démarrage...")

                    indexation_manager = None
                    if self.reference_file and self.reference_config:
                        indexation_manager = IndexationManager(
                            self.reference_file, self.reference_config
                        )

                    extractor = EmbeddedFileExtractor(indexation_manager)
                    extractor.total_input_files = len(self.selected_files)
                    extractor._processed_absolute_paths = set()

                    # ── Index de départ manuel ──────────────────────────────────
                    manual_index_str = self.start_index_var.get().strip()
                    if manual_index_str:
                        try:
                            manual_start = int(manual_index_str)
                            extractor.manual_start_index = manual_start
                            log(f"  🔢 Index de départ manuel : FJ_{manual_start}")
                        except ValueError:
                            extractor.manual_start_index = None
                    else:
                        extractor.manual_start_index = None

                    for i, file_path in enumerate(self.selected_files):
                        if self.stop_event.is_set():
                            log("⛔ Arrêt demandé — extraction initiale interrompue.")
                            stopped_early = True
                            break

                        try:
                            progress(i + 1, total_files,
                                    "Étape 2/4 — Extraction initiale",
                                    f"Traitement : {Path(file_path).name}")

                            # Tronquer à 40 chars pour éviter MAX_PATH (260) sur Windows
                            stem_safe = re.sub(r'[<>:"/\\|?*]', '_', Path(file_path).stem.strip())[:80]
                            file_dir = session_dir / stem_safe
                            if file_dir.exists():
                                file_dir = session_dir / f"{stem_safe}_{i+1}"
                            file_dir.mkdir(parents=True, exist_ok=True)

                            extracted = extractor.process_file(file_path, file_dir)

                            # Marquer le fichier source comme traité (évite retraitement en récursion)
                            try:
                                extractor._processed_absolute_paths.add(str(Path(file_path).resolve()))
                            except Exception:
                                extractor._processed_absolute_paths.add(str(file_path))

                            for msg in extractor.log_messages:
                                log(msg)
                            extractor.log_messages = []

                            count = len(extracted) if extracted else 0
                            total_extracted += count
                            files_processed_count += 1

                            if count > 0:
                                already_processed = set()
                                output_copy = file_dir / Path(file_path).name
                                if output_copy.exists():
                                    try:
                                        stat = output_copy.stat()
                                        with open(output_copy, 'rb') as f:
                                            head = f.read(2048)
                                        sig = (f"{output_copy.name}|{stat.st_size}|"
                                            f"{hashlib.md5(head).hexdigest()}")
                                        already_processed.add(sig)
                                        # Marquer aussi le chemin absolu de la copie de sortie
                                        extractor._processed_absolute_paths.add(str(output_copy.resolve()))
                                        log(f"  🔒 Copie source marquée : {output_copy.name}")
                                    except Exception:
                                        already_processed.add(str(output_copy.absolute()))

                                directories_to_recurse.append((file_dir, already_processed))

                        except Exception as file_err:
                            log(f"\n❌ Erreur inattendue sur {Path(file_path).name}: {file_err}")
                            for msg in extractor.log_messages:
                                log(msg)
                            extractor.log_messages = []
                            files_processed_count += 1

                # =================================================================
                # ÉTAPE 3 : EXTRACTION RÉCURSIVE
                # =================================================================
                if not stopped_early and directories_to_recurse:
                    log("")
                    log("🔄 ÉTAPE 3 : EXTRACTION RÉCURSIVE (3 niveaux)")
                    log("=" * 70)

                    total_dirs = len(directories_to_recurse)

                    for i, (directory, already_processed) in enumerate(directories_to_recurse):
                        if self.stop_event.is_set():
                            log("⛔ Arrêt demandé — récursion interrompue.")
                            stopped_early = True
                            break

                        progress(i + 1, max(total_dirs, 1),
                                "Étape 3/4 — Récursion",
                                f"Dossier : {directory.name}")
                        log(f"\n🔄 Dossier : {directory.name}")

                        total = extractor.process_extracted_files_recursively(
                            directory,
                            processed_signatures=already_processed,
                            depth=0,
                            max_depth=3
                        )

                        for msg in extractor.log_messages:
                            log(msg)
                        extractor.log_messages = []

                        log(f"  ✅ Récursion terminée : {total} fichier(s) traité(s)")

            except Exception as e:
                log(f"\n❌ ERREUR INATTENDUE : {e}")
                log(traceback.format_exc())

            finally:
                # =================================================================
                # ÉTAPE 4 : RAPPORT EXCEL (toujours généré, même si arrêt)
                # =================================================================
                if session_dir is not None:
                    if stopped_early:
                        log("")
                        log("⛔ TRAITEMENT INTERROMPU — génération du rapport partiel...")
                        progress(0, 1, "⛔ Arrêté — Rapport Excel partiel...", "Génération en cours...")
                    else:
                        log("")
                        log("📊 ÉTAPE 4 : GÉNÉRATION DU RAPPORT EXCEL...")
                        progress(0, 1, "Étape 4/4 — Rapport Excel", "Génération en cours...")

                    report_path = session_dir / f"Rapport_Extraction_{timestamp}.xlsx"

                    try:
                        if extractor is not None and extractor.generate_excel_report(report_path):
                            log(f"✅ Rapport créé : {report_path.name}")
                        else:
                            log("⚠ Erreur lors de la création du rapport")
                    except Exception as e:
                        log(f"❌ Erreur rapport : {str(e)}")

                    progress(1, 1,
                            "⛔ Arrêté avec rapport partiel" if stopped_early else "✅ Terminé !",
                            "Partiel — fichiers déjà traités sauvegardés." if stopped_early else "Extraction complète.")

                    # Bilan final
                    log("")
                    log("=" * 70)
                    if stopped_early:
                        log("⛔ ARRÊTÉ PAR L'UTILISATEUR (rapport partiel généré)")
                    else:
                        log("✓ TERMINÉ")
                    log(f"✓ Fichiers sources traités    : {files_processed_count}/{len(self.selected_files)}")
                    if files_with_excel:
                        log(f"⚡ Fichiers avec Excel activés : {len(files_with_excel)}")
                    log(f"✓ Extraction initiale          : {total_extracted} fichier(s)")
                    if not stopped_early:
                        log(f"✓ Récursion 3 niveaux effectuée")
                    log("=" * 70)

                    # Popup finale (identique arrêt ou fin normale)
                    _stopped = stopped_early
                    _files_with_excel = files_with_excel
                    _total_extracted = total_extracted
                    _files_processed = files_processed_count
                    _total_files = len(self.selected_files)
                    _session_dir = session_dir

                    def _show_done():
                        if _stopped:
                            message  = "⛔ Extraction interrompue par l'utilisateur.\n\n"
                            message += f"📁 {_files_processed}/{_total_files} fichier(s) traité(s)\n"
                        else:
                            message  = "✅ Extraction terminée !\n\n"
                            message += f"📁 {_files_processed} fichier(s) source(s) traité(s)\n"

                        if _files_with_excel:
                            message += f"⚡ {len(_files_with_excel)} fichier(s) avec Excel activés\n"
                        message += f"📎 {_total_extracted} fichier(s) extrait(s)\n"
                        if not _stopped:
                            message += f"🔄 Récursion automatique sur 3 niveaux effectuée\n"
                        message += f"📊 Rapport Excel généré\n\n"
                        message += "Ouvrir le dossier ?"

                        title = "Interrompu" if _stopped else "Terminé"
                        if messagebox.askyesno(title, message):
                            self.open_folder(_session_dir)

                    self.root.after(0, _show_done)

                self._finish_extraction()

    def _finish_extraction(self):
        """Remet l'interface en état initial (thread-safe)"""
        def _reset():
            self.extract_btn.config(state='normal', text="🚀 ACTIVER EXCEL + EXTRAIRE")
            self.stop_btn.config(state='disabled')
        self.root.after(0, _reset)
    
    def open_output_folder(self):
        """Ouvre dossier sortie"""
        folder = Path(self.output_dir.get())
        if folder.exists():
            self.open_folder(folder)
        else:
            messagebox.showinfo("Info", "Le dossier n'existe pas")
    
    def open_folder(self, folder_path):
        """Ouvre dossier"""
        folder_path = Path(folder_path)
        if sys.platform == 'win32':
            os.startfile(folder_path)
        elif sys.platform == 'darwin':
            subprocess.run(['open', str(folder_path)])
        else:
            subprocess.run(['xdg-open', str(folder_path)])


def test_nomenclature_fj():
    """Teste la nomenclature FJ"""
    print("\n" + "="*70)
    print("TEST NOMENCLATURE FJ")
    print("="*70)
    
    # Créer extracteur de test
    extractor = EmbeddedFileExtractor()
    
    # Test 1 : Nom avec FJ
    extractor.current_source_file = "DER-348818_3_FJ_2.docx"
    extractor.extracted_count = 1
    
    result1 = extractor._generate_output_name("DER-348818_3_FJ_2", "Budget.xlsx", ".xlsx")
    print(f"\n✅ Test 1 - Nom avec FJ:")
    print(f"  Base: DER-348818_3_FJ_2")
    print(f"  Fichier original: Budget.xlsx")
    print(f"  Résultat: {result1}")
    print(f"  Attendu: DER-348818_3_FJ_3_Budget.xlsx")
    assert result1 == "DER-348818_3_FJ_3_Budget.xlsx", f"ERREUR: {result1}"
    
    # Test 2 : Deuxième extraction avec FJ
    extractor.extracted_count = 2
    result2 = extractor._generate_output_name("DER-348818_3_FJ_2", "Stats.xlsx", ".xlsx")
    print(f"\n✅ Test 2 - Deuxième fichier FJ:")
    print(f"  Résultat: {result2}")
    print(f"  Attendu: DER-348818_3_FJ_4_Stats.xlsx")
    assert result2 == "DER-348818_3_FJ_4_Stats.xlsx", f"ERREUR: {result2}"
    
    # Test 3 : Nom normal
    extractor.current_source_file = "Document_1.docx"
    extractor.extracted_count = 1
    
    result3 = extractor._generate_output_name("Document_1", "Budget.xlsx", ".xlsx")
    print(f"\n✅ Test 3 - Nom normal:")
    print(f"  Base: Document_1")
    print(f"  Fichier original: Budget.xlsx")
    print(f"  Résultat: {result3}")
    print(f"  Attendu: Document_1_1_Budget.xlsx")
    assert result3 == "Document_1_1_Budget.xlsx", f"ERREUR: {result3}"
    
    print("\n" + "="*70)
    print("✅ TOUS LES TESTS RÉUSSIS!")
    print("="*70 + "\n")


def main():
    """Fonction principale"""
    root = tk.Tk()
    app = SimpleFileExtractorGUI(root)
    
    root.update_idletasks()
    width = root.winfo_width()
    height = root.winfo_height()
    x = (root.winfo_screenwidth() // 2) - (width // 2)
    y = (root.winfo_screenheight() // 2) - (height // 2)
    root.geometry(f'{width}x{height}+{x}+{y}')
    
    root.mainloop()


if __name__ == "__main__":
    # Test nomenclature FJ
    if len(sys.argv) > 1 and sys.argv[1] == "--test":
        test_nomenclature_fj()
        sys.exit(0)
    
    # Lancement normal
    print("=" * 70)
    print("Extracteur de Fichiers Incorporés - VERSION CORRIGÉE")
    print("=" * 70)
    print("\n✅ CORRECTION 1 : Activation Excel fonctionnelle")
    print("   → Scripts VBS optimisés pour activation rapide")
    print("   → Détection intelligente des objets Excel uniquement")
    print("")
    print("✅ CORRECTION 2 : Fichiers ZIP remplacés dans documents")
    print("   → Les ZIP extraits sont marqués pour remplacement")
    print("   → Apparaissent dans les documents modifiés")
    print("")
    print("✅ CORRECTION 3 : Nomenclature FJ corrigée")
    print("   → DER-348818_3_FJ_2 → FJ_3, FJ_4, FJ_5...")
    print("   → Incrémentation après FJ au lieu d'ajouter un compteur")
    print("")
    print("✅ Ordre réel des fichiers respecté (PPTX, DOCX, XLSX)")
    print("✅ Fichiers modifiés avec nom identique à l'original")
    print("✅ Support DOCX, DOC, XLSX, XLSM, XLS, PPTX, PPT, PDF, MSG")
    print("✅ Indexation personnalisée via fichier de référence")
    print("✅ Rapport Excel détaillé automatique")
    print("\nLancement...")
    print("=" * 70)
    
    main()